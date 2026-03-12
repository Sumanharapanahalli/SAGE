#!/usr/bin/env python3
"""
SAGE[ai] Business Case Presentation Generator
==============================================
Generates a professional PowerPoint (.pptx) presentation for SAGE[ai],
an ISO 13485-compliant autonomous AI developer agent for medical device manufacturing.

Installation:
    pip install python-pptx

Run:
    python scripts/generate_presentation.py

Output:
    docs/SageAI_Business_Case.pptx
    docs/SageAI_Business_Case_Summary.md
"""

import os
import sys

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt, Emu
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN
    from pptx.oxml.ns import qn
    from pptx.util import Inches, Pt
    import copy
    from lxml import etree
except ImportError as e:
    print(f"ERROR: Missing dependency — {e}")
    print("Install with:  pip install python-pptx")
    sys.exit(1)

# ---------------------------------------------------------------------------
# Brand colours
# ---------------------------------------------------------------------------
GREEN  = RGBColor(0x16, 0xA3, 0x4A)
DARK   = RGBColor(0x1F, 0x29, 0x37)
GRAY   = RGBColor(0x37, 0x41, 0x51)
LIGHT  = RGBColor(0xF9, 0xFA, 0xFB)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
AMBER  = RGBColor(0xF5, 0x9E, 0x0B)
RED_C  = RGBColor(0xDC, 0x26, 0x26)
BLUE   = RGBColor(0x25, 0x63, 0xEB)
GREEN_LIGHT = RGBColor(0xDC, 0xFC, 0xE7)
DARK_GREEN  = RGBColor(0x14, 0x53, 0x2D)

SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)

# ---------------------------------------------------------------------------
# Output paths
# ---------------------------------------------------------------------------
SCRIPT_DIR   = os.path.dirname(os.path.abspath(__file__))
PROJECT_ROOT = os.path.dirname(SCRIPT_DIR)
DOCS_DIR     = os.path.join(PROJECT_ROOT, "docs")
PPTX_PATH    = os.path.join(DOCS_DIR, "SageAI_Business_Case.pptx")
MD_PATH      = os.path.join(DOCS_DIR, "SageAI_Business_Case_Summary.md")

# ---------------------------------------------------------------------------
# Helper utilities
# ---------------------------------------------------------------------------

def _rgb_to_hex(rgb: RGBColor) -> str:
    return f"#{rgb[0]:02X}{rgb[1]:02X}{rgb[2]:02X}"


def set_slide_background(slide, color: RGBColor = WHITE):
    """Fill slide background with a solid colour."""
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_textbox(
    slide,
    text: str,
    left, top, width, height,
    font_size: int = 16,
    bold: bool = False,
    italic: bool = False,
    color: RGBColor = GRAY,
    align=PP_ALIGN.LEFT,
    word_wrap: bool = True,
    font_name: str = "Calibri",
    bg_color: RGBColor = None,
):
    """Add a single-paragraph text box to a slide."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = word_wrap

    if bg_color is not None:
        fill = txBox.fill
        fill.solid()
        fill.fore_color.rgb = bg_color

    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    run.font.name = font_name
    return txBox


def add_bullet_textbox(
    slide,
    title: str,
    bullets: list,
    left, top, width, height,
    title_size: int = 18,
    bullet_size: int = 15,
    title_color: RGBColor = DARK,
    bullet_color: RGBColor = GRAY,
    font_name: str = "Calibri",
    title_bold: bool = True,
):
    """Add a text box with an optional title and bullet list."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True

    if title:
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.LEFT
        run = p.add_run()
        run.text = title
        run.font.size = Pt(title_size)
        run.font.bold = title_bold
        run.font.color.rgb = title_color
        run.font.name = font_name
    else:
        # Remove blank first paragraph
        tf.paragraphs[0]._p.getparent().remove(tf.paragraphs[0]._p)

    for bullet in bullets:
        p = tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.space_before = Pt(3)
        run = p.add_run()
        run.text = bullet
        run.font.size = Pt(bullet_size)
        run.font.color.rgb = bullet_color
        run.font.name = font_name

    return txBox


def add_rect(slide, left, top, width, height, fill_color: RGBColor, line_color: RGBColor = None):
    """Add a filled rectangle shape."""
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        left, top, width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if line_color:
        shape.line.color.rgb = line_color
    else:
        shape.line.fill.background()
    return shape


def add_table_to_slide(
    slide,
    data: list,            # list of rows; first row = header
    left, top, width, height,
    header_bg: RGBColor = GREEN,
    header_fg: RGBColor = WHITE,
    row_alt_bg: RGBColor = GREEN_LIGHT,
    row_bg: RGBColor = WHITE,
    font_size: int = 12,
    header_font_size: int = 13,
    font_name: str = "Calibri",
    col_widths: list = None,
):
    """Add a formatted table to a slide."""
    rows = len(data)
    cols = len(data[0]) if data else 1

    table = slide.shapes.add_table(rows, cols, left, top, width, height).table

    # Column widths
    if col_widths:
        for i, w in enumerate(col_widths):
            table.columns[i].width = w

    for r_idx, row_data in enumerate(data):
        for c_idx, cell_text in enumerate(row_data):
            cell = table.cell(r_idx, c_idx)
            cell.text = str(cell_text)

            # Background
            fill = cell.fill
            fill.solid()
            if r_idx == 0:
                fill.fore_color.rgb = header_bg
                fg = header_fg
                bold = True
                fsize = header_font_size
            elif r_idx % 2 == 0:
                fill.fore_color.rgb = row_alt_bg
                fg = DARK
                bold = False
                fsize = font_size
            else:
                fill.fore_color.rgb = row_bg
                fg = DARK
                bold = False
                fsize = font_size

            tf = cell.text_frame
            tf.word_wrap = True
            for para in tf.paragraphs:
                para.alignment = PP_ALIGN.LEFT
                for run in para.runs:
                    run.font.size = Pt(fsize)
                    run.font.bold = bold
                    run.font.color.rgb = fg
                    run.font.name = font_name

    return table


def add_slide_header(slide, heading: str, font_size: int = 28):
    """Add a standard green-bar heading at the top of a content slide."""
    # Green accent bar
    add_rect(slide, Inches(0), Inches(0), SLIDE_W, Inches(0.08), GREEN)
    # Heading text
    add_textbox(
        slide, heading,
        left=Inches(0.5), top=Inches(0.12),
        width=Inches(12.33), height=Inches(0.65),
        font_size=font_size, bold=True, color=DARK,
        align=PP_ALIGN.LEFT,
    )
    # Thin separator line
    add_rect(slide, Inches(0.5), Inches(0.78), Inches(12.33), Inches(0.03), GREEN)


def add_footer(slide, text: str = "CONFIDENTIAL — SAGE[ai] Business Case | March 2026"):
    """Add a bottom footer bar."""
    add_rect(slide, Inches(0), Inches(7.2), SLIDE_W, Inches(0.3), DARK)
    add_textbox(
        slide, text,
        left=Inches(0.3), top=Inches(7.2),
        width=Inches(12.73), height=Inches(0.3),
        font_size=9, color=WHITE, align=PP_ALIGN.CENTER,
    )


# ---------------------------------------------------------------------------
# Individual slide builders
# ---------------------------------------------------------------------------

def build_slide_01_title(prs: Presentation):
    """Slide 1 — Title slide."""
    slide_layout = prs.slide_layouts[6]  # blank
    slide = prs.slides.add_slide(slide_layout)
    set_slide_background(slide, WHITE)

    # Top green accent bar
    add_rect(slide, Inches(0), Inches(0), SLIDE_W, Inches(0.2), GREEN)

    # Large decorative rectangle (right side)
    add_rect(slide, Inches(8.8), Inches(0.2), Inches(4.53), Inches(7.0), GREEN_LIGHT)

    # SAGE[ai] logo-style text
    add_textbox(
        slide, "SAGE[ai]",
        left=Inches(0.7), top=Inches(1.6),
        width=Inches(7.8), height=Inches(1.1),
        font_size=52, bold=True, color=GREEN,
        align=PP_ALIGN.LEFT,
    )

    # Title
    add_textbox(
        slide, "Autonomous Manufacturing Intelligence",
        left=Inches(0.7), top=Inches(2.7),
        width=Inches(7.8), height=Inches(0.9),
        font_size=32, bold=True, color=DARK,
        align=PP_ALIGN.LEFT,
    )

    # Subtitle
    add_textbox(
        slide, "Transforming Medical Device Production with Agentic AI",
        left=Inches(0.7), top=Inches(3.65),
        width=Inches(7.8), height=Inches(0.65),
        font_size=20, bold=False, color=GRAY,
        align=PP_ALIGN.LEFT,
    )

    # Divider line
    add_rect(slide, Inches(0.7), Inches(4.45), Inches(4.0), Inches(0.04), GREEN)

    # Right-panel callout text
    for i, (line, size, bold, col) in enumerate([
        ("ISO 13485", 18, True, DARK_GREEN),
        ("Compliant", 14, False, DARK),
        ("", 8, False, WHITE),
        ("Agentic AI", 18, True, DARK_GREEN),
        ("for MedTech", 14, False, DARK),
        ("", 8, False, WHITE),
        ("Human-in-the-Loop", 16, True, DARK_GREEN),
        ("Architecture", 14, False, DARK),
    ]):
        add_textbox(
            slide, line,
            left=Inches(9.0), top=Inches(1.4 + i * 0.42),
            width=Inches(4.0), height=Inches(0.42),
            font_size=size, bold=bold, color=col,
            align=PP_ALIGN.CENTER,
        )

    # Bottom green accent bar
    add_rect(slide, Inches(0), Inches(7.1), SLIDE_W, Inches(0.4), GREEN)
    add_textbox(
        slide, "CONFIDENTIAL — Internal Business Case | March 2026",
        left=Inches(0.3), top=Inches(7.1),
        width=Inches(12.73), height=Inches(0.4),
        font_size=11, bold=False, color=WHITE,
        align=PP_ALIGN.CENTER,
    )


def build_slide_02_problem(prs: Presentation):
    """Slide 2 — The Problem We're Solving."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide, WHITE)
    add_slide_header(slide, "Current State: Manual Bottlenecks at Scale")

    bullets = [
        "⏱  45–60 min per error log analysis (manual)",
        "🔁  3–5 day MR review backlog",
        "📋  8+ hours/week compliance reporting",
        "🧠  Expert knowledge lost when engineers leave",
        "⚠️  Error detection relies on human vigilance 24/7",
    ]
    add_bullet_textbox(
        slide,
        title="",
        bullets=bullets,
        left=Inches(0.5), top=Inches(1.0),
        width=Inches(7.6), height=Inches(5.5),
        bullet_size=18,
        bullet_color=DARK,
    )

    # Stat box (right side)
    add_rect(slide, Inches(8.5), Inches(1.6), Inches(4.3), Inches(3.2), DARK)
    add_textbox(
        slide, "68 hrs/week",
        left=Inches(8.5), top=Inches(2.0),
        width=Inches(4.3), height=Inches(1.1),
        font_size=38, bold=True, color=GREEN,
        align=PP_ALIGN.CENTER,
    )
    add_textbox(
        slide, "of preventable\nmanual work",
        left=Inches(8.5), top=Inches(3.1),
        width=Inches(4.3), height=Inches(0.9),
        font_size=20, bold=False, color=WHITE,
        align=PP_ALIGN.CENTER,
    )
    add_textbox(
        slide, "every single week",
        left=Inches(8.5), top=Inches(3.95),
        width=Inches(4.3), height=Inches(0.5),
        font_size=14, bold=False, color=AMBER,
        align=PP_ALIGN.CENTER,
    )

    add_footer(slide)


def build_slide_03_solution(prs: Presentation):
    """Slide 3 — The Solution."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide, WHITE)
    add_slide_header(slide, "SAGE[ai]: Always-On AI Engineering Partner")

    bullets = [
        "✅  Analyzes error logs in <60 seconds (vs 45–60 min manual)",
        "✅  AI code reviews with multi-step reasoning (ReAct loop)",
        "✅  Monitors Teams, Metabase, GitLab — detects events in real-time",
        "✅  Every decision logged to immutable ISO 13485 audit trail",
        "✅  Human-in-the-loop: AI advises, humans decide — always",
        "✅  Learns from every correction via vector memory (RAG)",
        "✅  Web dashboard for all stakeholders — no CLI required",
    ]
    add_bullet_textbox(
        slide,
        title="",
        bullets=bullets,
        left=Inches(0.5), top=Inches(1.0),
        width=Inches(12.33), height=Inches(5.8),
        bullet_size=18,
        bullet_color=DARK,
    )

    add_footer(slide)


def build_slide_04_evidence(prs: Presentation):
    """Slide 4 — Industry Evidence."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide, WHITE)
    add_slide_header(slide, "Companies Already Winning with Agentic AI")

    data = [
        ["Company", "Implementation", "Result"],
        ["Siemens", "AI predictive maintenance agents", "30% less unplanned downtime, €1.5B saved/yr"],
        ["BMW", "AI quality inspection", "99.5% defect detection, 30% QC cost reduction"],
        ["Bosch", "AI-assisted code review", "40% faster releases, 60% fewer defect escapes"],
        ["Amazon", "CodeGuru automated review", "50% reduction in production incidents"],
        ["Microsoft", "GitHub Copilot for developers", "55% faster code completion, 46% more PRs/day"],
        ["Medtronic", "AI in quality management system", "35% shorter CAPA cycle, faster FDA submissions"],
        ["Stryker", "AI regulatory documentation", "40% reduction in submission preparation time"],
        ["J&J", "AI manufacturing analytics", "20% yield improvement, significant waste reduction"],
    ]

    col_widths = [Inches(1.8), Inches(4.8), Inches(5.5)]
    add_table_to_slide(
        slide, data,
        left=Inches(0.5), top=Inches(1.0),
        width=Inches(12.33), height=Inches(5.9),
        col_widths=col_widths,
        font_size=12,
        header_font_size=14,
    )

    add_footer(slide)


def build_slide_05_lean(prs: Presentation):
    """Slide 5 — Lean Manufacturing Alignment."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide, WHITE)
    add_slide_header(slide, "SAGE[ai] Is Built on Lean Principles")

    data = [
        ["Lean Principle", "How SAGE[ai] Delivers"],
        ["Eliminate Waste (Muda)", "Removes 60+ hrs/week of repetitive analysis"],
        ["Continuous Improvement (Kaizen)", "Learns from every rejection via RAG memory"],
        ["Error-Proofing (Poka-yoke)", "Human-in-the-loop prevents AI errors reaching production"],
        ["Visual Management", "Real-time dashboard for all stakeholders"],
        ["Single Piece Flow", "Single-lane task queue: deterministic, auditable"],
        ["Respect for People", "Amplifies engineers; never replaces human judgment"],
    ]

    col_widths = [Inches(4.5), Inches(7.83)]
    add_table_to_slide(
        slide, data,
        left=Inches(0.5), top=Inches(1.0),
        width=Inches(12.33), height=Inches(5.8),
        col_widths=col_widths,
        font_size=14,
        header_font_size=15,
    )

    add_footer(slide)


def build_slide_06_roi(prs: Presentation):
    """Slide 6 — ROI Analysis."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide, WHITE)
    add_slide_header(slide, "Return on Investment: Month 2 Payback")

    data = [
        ["Activity", "Before SAGE[ai]", "After SAGE[ai]", "Savings"],
        ["Error log analysis", "45 min × 10/day × 2 eng\n= 15 hrs/day", "<5 min total", "93% reduction"],
        ["MR code review", "3 hrs × 15 MRs/week\n= 45 hrs/week", "~4 hrs/week", "91% reduction"],
        ["Compliance reporting", "8 hrs/week manual", "0 hrs (auto-generated)", "100% reduction"],
        ["Knowledge capture", "Lost on attrition", "Stored in vector memory", "Permanent"],
        ["TOTAL SAVINGS", "68 hrs/week ≈ 1.7 FTE", "—", "~€120K/year (1 FTE cost)"],
    ]

    col_widths = [Inches(2.8), Inches(3.5), Inches(3.0), Inches(2.83)]
    add_table_to_slide(
        slide, data,
        left=Inches(0.5), top=Inches(1.0),
        width=Inches(12.13), height=Inches(4.9),
        col_widths=col_widths,
        font_size=12,
        header_font_size=13,
    )

    # Bold footer text
    add_rect(slide, Inches(0.5), Inches(6.25), Inches(12.33), Inches(0.65), GREEN_LIGHT)
    add_textbox(
        slide,
        "Conservative estimate: break-even in under 2 months",
        left=Inches(0.6), top=Inches(6.28),
        width=Inches(12.13), height=Inches(0.55),
        font_size=18, bold=True, color=DARK_GREEN,
        align=PP_ALIGN.CENTER,
    )

    add_footer(slide)


def build_slide_07_architecture(prs: Presentation):
    """Slide 7 — How It Works (Architecture)."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide, WHITE)
    add_slide_header(slide, "Simple, Auditable, Secure Architecture")

    # Flow diagram boxes
    boxes = [
        ("Teams /\nMetabase /\nGitLab", Inches(0.3)),
        ("Monitor\nAgent", Inches(2.5)),
        ("Task\nQueue", Inches(4.7)),
        ("Analyst /\nDeveloper\nAgent", Inches(6.9)),
        ("Human\nGate", Inches(9.1)),
        ("Audit\nTrail", Inches(11.3)),
    ]

    for label, left in boxes:
        add_rect(slide, left, Inches(1.25), Inches(1.8), Inches(1.2), DARK, None)
        add_textbox(
            slide, label,
            left=left, top=Inches(1.25),
            width=Inches(1.8), height=Inches(1.2),
            font_size=11, bold=True, color=WHITE,
            align=PP_ALIGN.CENTER,
        )

    # Arrows between boxes
    arrow_positions = [Inches(2.1), Inches(4.3), Inches(6.5), Inches(8.7), Inches(10.9)]
    for ax in arrow_positions:
        add_textbox(
            slide, "→",
            left=ax, top=Inches(1.55),
            width=Inches(0.4), height=Inches(0.6),
            font_size=20, bold=True, color=GREEN,
            align=PP_ALIGN.CENTER,
        )

    # Three-column descriptions
    col_data = [
        (
            "Event Sources",
            [
                "• Teams channels",
                "• Metabase dashboards",
                "• GitLab issues & MRs",
                "• Error log uploads",
            ],
            Inches(0.3),
        ),
        (
            "AI Agents",
            [
                "• AnalystAgent (log analysis)",
                "• DeveloperAgent (code review + MR creation)",
                "• PlannerAgent (orchestration)",
                "• ReAct multi-step reasoning loop",
            ],
            Inches(4.45),
        ),
        (
            "Compliance Layer",
            [
                "• Immutable SQLite audit log",
                "• ISO 13485 trace IDs on every decision",
                "• FDA 21 CFR Part 11 records",
                "• Human approval gate on every proposal",
            ],
            Inches(8.6),
        ),
    ]

    for col_title, col_bullets, col_left in col_data:
        add_rect(slide, col_left, Inches(2.7), Inches(4.0), Inches(3.9), LIGHT, None)
        add_textbox(
            slide, col_title,
            left=col_left + Inches(0.1), top=Inches(2.78),
            width=Inches(3.8), height=Inches(0.45),
            font_size=14, bold=True, color=GREEN,
            align=PP_ALIGN.LEFT,
        )
        add_rect(slide, col_left + Inches(0.1), Inches(3.22), Inches(3.8), Inches(0.03), GREEN)
        add_bullet_textbox(
            slide,
            title="",
            bullets=col_bullets,
            left=col_left + Inches(0.1), top=Inches(3.28),
            width=Inches(3.8), height=Inches(2.1),
            bullet_size=12,
            bullet_color=DARK,
        )

    add_footer(slide)


def build_slide_08_capabilities(prs: Presentation):
    """Slide 8 — Key Capabilities (2×3 grid)."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide, WHITE)
    add_slide_header(slide, "What SAGE[ai] Does Today")

    capabilities = [
        ("🔍", "Log Analysis", "AI triage in <60s\nSeverity: RED / AMBER / GREEN"),
        ("🤖", "Code Review", "ReAct multi-step reasoning\n+ pipeline check"),
        ("📋", "MR Creation", "Auto-draft from GitLab issue\nBranch naming included"),
        ("👁", "24/7 Monitor", "Teams, Metabase, GitLab\nevent detection"),
        ("📊", "Audit Trail", "Every decision traceable\nISO 13485 compliant"),
        ("🌐", "Web Dashboard", "No CLI needed\nDashboard, Analyst, Developer, Audit, Monitor"),
    ]

    positions = [
        (Inches(0.4),  Inches(1.1)),
        (Inches(4.55), Inches(1.1)),
        (Inches(8.7),  Inches(1.1)),
        (Inches(0.4),  Inches(4.0)),
        (Inches(4.55), Inches(4.0)),
        (Inches(8.7),  Inches(4.0)),
    ]

    for idx, ((icon, title, desc), (bx, by)) in enumerate(zip(capabilities, positions)):
        add_rect(slide, bx, by, Inches(3.8), Inches(2.6), LIGHT, None)
        # Top green bar for each box
        add_rect(slide, bx, by, Inches(3.8), Inches(0.07), GREEN, None)
        add_textbox(
            slide, icon,
            left=bx + Inches(0.15), top=by + Inches(0.12),
            width=Inches(0.7), height=Inches(0.55),
            font_size=26, bold=False, color=DARK,
            align=PP_ALIGN.LEFT,
        )
        add_textbox(
            slide, title,
            left=bx + Inches(0.85), top=by + Inches(0.15),
            width=Inches(2.8), height=Inches(0.5),
            font_size=16, bold=True, color=GREEN,
            align=PP_ALIGN.LEFT,
        )
        add_textbox(
            slide, desc,
            left=bx + Inches(0.15), top=by + Inches(0.75),
            width=Inches(3.5), height=Inches(1.7),
            font_size=13, bold=False, color=DARK,
            align=PP_ALIGN.LEFT,
        )

    add_footer(slide)


def build_slide_09_compliance(prs: Presentation):
    """Slide 9 — Compliance First."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide, WHITE)
    add_slide_header(slide, "Built for Regulated Medical Device Environments")

    # Left column
    add_rect(slide, Inches(0.4), Inches(1.0), Inches(5.9), Inches(5.9), LIGHT, None)
    add_textbox(
        slide, "Standards Met",
        left=Inches(0.5), top=Inches(1.1),
        width=Inches(5.7), height=Inches(0.5),
        font_size=16, bold=True, color=GREEN,
        align=PP_ALIGN.LEFT,
    )
    add_rect(slide, Inches(0.5), Inches(1.6), Inches(5.7), Inches(0.03), GREEN)

    standards = [
        "✅  ISO 13485:2016 — Quality Management System",
        "✅  ISO 14971:2019 — Risk Management\n     (7 risks identified + controlled)",
        "✅  IEC 62304:2006 — Medical Device Software Lifecycle",
        "✅  FDA 21 CFR Part 11 — Electronic Records",
        "✅  FDA Cybersecurity Guidance 2023",
    ]
    add_bullet_textbox(
        slide,
        title="",
        bullets=standards,
        left=Inches(0.5), top=Inches(1.65),
        width=Inches(5.7), height=Inches(5.0),
        bullet_size=14,
        bullet_color=DARK,
    )

    # Right column
    add_rect(slide, Inches(6.8), Inches(1.0), Inches(6.1), Inches(5.9), DARK, None)
    add_textbox(
        slide, "How We Comply",
        left=Inches(6.9), top=Inches(1.1),
        width=Inches(5.9), height=Inches(0.5),
        font_size=16, bold=True, color=GREEN,
        align=PP_ALIGN.LEFT,
    )
    add_rect(slide, Inches(6.9), Inches(1.6), Inches(5.9), Inches(0.03), GREEN)

    compliance = [
        "•  Immutable append-only audit log\n   (no deletes, ever)",
        "•  UUID trace ID on every AI decision",
        "•  Human approval gate on every proposal",
        "•  Full DHF: SRS, RTM, V&V Plan,\n   SOUP Inventory",
        "•  Air-gapped LLM option\n   (no cloud dependency)",
    ]
    add_bullet_textbox(
        slide,
        title="",
        bullets=compliance,
        left=Inches(6.9), top=Inches(1.65),
        width=Inches(5.9), height=Inches(5.0),
        bullet_size=14,
        bullet_color=WHITE,
    )

    add_footer(slide)


def build_slide_10_self_improving(prs: Presentation):
    """Slide 10 — The Self-Improving System."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide, WHITE)
    add_slide_header(slide, "SAGE[ai] Gets Better Every Day — From Your Own Team")

    # Three-step cycle (horizontal arrangement)
    steps = [
        ("1", "USER SUBMITS REQUEST",
         "💡 Click 'Request Improvement'\non any module"),
        ("2", "AI PLANS",
         "Planner Agent decomposes\ninto subtasks, queued\nfor implementation"),
        ("3", "IMPLEMENTED & LEARNED",
         "Change deployed;\nRAG memory updated\nwith new context"),
    ]

    step_colors = [GREEN, DARK, GREEN_LIGHT]
    text_colors = [WHITE, WHITE, DARK]
    num_colors  = [WHITE, GREEN, GREEN]

    step_positions = [Inches(0.6), Inches(4.55), Inches(8.5)]

    for (num, title, desc), left, bg, tc, nc in zip(
            steps, step_positions, step_colors, text_colors, num_colors):
        add_rect(slide, left, Inches(1.15), Inches(3.8), Inches(3.6), bg, None)
        add_textbox(
            slide, num,
            left=left + Inches(0.1), top=Inches(1.2),
            width=Inches(0.6), height=Inches(0.7),
            font_size=32, bold=True, color=nc,
            align=PP_ALIGN.LEFT,
        )
        add_textbox(
            slide, title,
            left=left + Inches(0.1), top=Inches(1.95),
            width=Inches(3.6), height=Inches(0.65),
            font_size=14, bold=True, color=tc,
            align=PP_ALIGN.LEFT,
        )
        add_textbox(
            slide, desc,
            left=left + Inches(0.1), top=Inches(2.65),
            width=Inches(3.6), height=Inches(1.9),
            font_size=13, bold=False, color=tc,
            align=PP_ALIGN.LEFT,
        )

    # Arrows between steps
    for ax in [Inches(4.15), Inches(8.1)]:
        add_textbox(
            slide, "→",
            left=ax, top=Inches(2.5),
            width=Inches(0.45), height=Inches(0.65),
            font_size=28, bold=True, color=GREEN,
            align=PP_ALIGN.CENTER,
        )

    # Access note
    add_rect(slide, Inches(0.5), Inches(5.0), Inches(12.33), Inches(0.95), GREEN_LIGHT, None)
    add_textbox(
        slide,
        "During development: open access for all engineers.\n"
        "Post-release: role-based access control (admin approval required).",
        left=Inches(0.6), top=Inches(5.05),
        width=Inches(12.13), height=Inches(0.85),
        font_size=13, bold=False, color=DARK,
        align=PP_ALIGN.CENTER,
    )

    # Emphasis text
    add_textbox(
        slide,
        "The system improves itself through the same AI pipeline it provides to engineers.",
        left=Inches(0.5), top=Inches(6.05),
        width=Inches(12.33), height=Inches(0.5),
        font_size=15, bold=True, color=GREEN,
        align=PP_ALIGN.CENTER,
    )

    add_footer(slide)


def build_slide_11_roadmap(prs: Presentation):
    """Slide 11 — Implementation Roadmap."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide, WHITE)
    add_slide_header(slide, "From Pilot to Production: 3-Month Plan")

    data = [
        ["Phase", "Timeline", "Deliverable", "Status"],
        ["Phase 1–3: Core System",
         "Done",
         "CLI, Analyst, Developer, Monitor, API",
         "✅ Complete"],
        ["Phase 4: Web UI + Agentic",
         "Done",
         "Dashboard, ReAct loop, Planner, Regulatory docs",
         "✅ Complete"],
        ["Phase 5: Pilot Deployment",
         "Month 1",
         "Production deploy on internal network, engineer training",
         "🔵 Next"],
        ["Phase 6: Measure & Iterate",
         "Month 2",
         "KPI tracking, RAG feedback loop, first improvements",
         "🔵 Planned"],
        ["Phase 7: Scale",
         "Month 3",
         "Multi-team rollout, Spira integration, executive dashboards",
         "🔵 Planned"],
    ]

    col_widths = [Inches(2.8), Inches(1.5), Inches(5.83), Inches(1.9)]
    add_table_to_slide(
        slide, data,
        left=Inches(0.5), top=Inches(1.0),
        width=Inches(12.03), height=Inches(5.8),
        col_widths=col_widths,
        font_size=13,
        header_font_size=14,
    )

    add_footer(slide)


def build_slide_12_ask(prs: Presentation):
    """Slide 12 — The Ask."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide, WHITE)
    add_slide_header(slide, "Next Steps: 3-Month Pilot Proposal")

    # Left column
    add_rect(slide, Inches(0.4), Inches(1.05), Inches(5.9), Inches(5.4), DARK, None)
    add_textbox(
        slide, "What We Need",
        left=Inches(0.5), top=Inches(1.15),
        width=Inches(5.7), height=Inches(0.5),
        font_size=17, bold=True, color=GREEN,
        align=PP_ALIGN.LEFT,
    )
    add_rect(slide, Inches(0.5), Inches(1.65), Inches(5.7), Inches(0.03), GREEN)

    needs = [
        "✅  Production deployment approval\n     (internal network)",
        "✅  2-day engineering team training session",
        "✅  Access to GitLab, Metabase,\n     Teams API credentials",
        "✅  Dedicated server or VM\n     (8-core CPU, 32GB RAM, optional GPU)",
    ]
    add_bullet_textbox(
        slide,
        title="",
        bullets=needs,
        left=Inches(0.5), top=Inches(1.7),
        width=Inches(5.7), height=Inches(4.6),
        bullet_size=14,
        bullet_color=WHITE,
    )

    # Right column
    add_rect(slide, Inches(6.8), Inches(1.05), Inches(6.1), Inches(5.4), GREEN_LIGHT, None)
    add_textbox(
        slide, "What You Get",
        left=Inches(6.9), top=Inches(1.15),
        width=Inches(5.9), height=Inches(0.5),
        font_size=17, bold=True, color=DARK_GREEN,
        align=PP_ALIGN.LEFT,
    )
    add_rect(slide, Inches(6.9), Inches(1.65), Inches(5.9), Inches(0.03), GREEN)

    gets = [
        "📈  KPI dashboard by Week 2",
        "⏱   60+ hours/week reclaimed\n      from manual work",
        "📋  ISO 13485 audit trail from day 1",
        "🤖  AI that learns from your\n      engineers' expertise",
        "💰  ROI positive within 2 months",
    ]
    add_bullet_textbox(
        slide,
        title="",
        bullets=gets,
        left=Inches(6.9), top=Inches(1.7),
        width=Inches(5.9), height=Inches(4.6),
        bullet_size=14,
        bullet_color=DARK,
    )

    # Bold footer statement
    add_rect(slide, Inches(0), Inches(6.6), SLIDE_W, Inches(0.55), GREEN, None)
    add_textbox(
        slide,
        "SAGE[ai] is not a replacement for your engineers — it's their most productive teammate.",
        left=Inches(0.3), top=Inches(6.62),
        width=Inches(12.73), height=Inches(0.48),
        font_size=15, bold=True, color=WHITE,
        align=PP_ALIGN.CENTER,
    )


# ---------------------------------------------------------------------------
# Markdown summary writer
# ---------------------------------------------------------------------------

MARKDOWN_CONTENT = """\
# SAGE[ai]: Autonomous Manufacturing Intelligence
### Business Case Presentation — March 2026
*CONFIDENTIAL — Internal Business Case*

---

## Slide 1 — Title

**SAGE[ai]: Autonomous Manufacturing Intelligence**
*Transforming Medical Device Production with Agentic AI*

---

## Slide 2 — Current State: Manual Bottlenecks at Scale

- ⏱ 45–60 min per error log analysis (manual)
- 🔁 3–5 day MR review backlog
- 📋 8+ hours/week compliance reporting
- 🧠 Expert knowledge lost when engineers leave
- ⚠️ Error detection relies on human vigilance 24/7

> **68 hrs/week of preventable manual work** — every single week

---

## Slide 3 — SAGE[ai]: Always-On AI Engineering Partner

- ✅ Analyzes error logs in <60 seconds (vs 45–60 min manual)
- ✅ AI code reviews with multi-step reasoning (ReAct loop)
- ✅ Monitors Teams, Metabase, GitLab — detects events in real-time
- ✅ Every decision logged to immutable ISO 13485 audit trail
- ✅ Human-in-the-loop: AI advises, humans decide — always
- ✅ Learns from every correction via vector memory (RAG)
- ✅ Web dashboard for all stakeholders — no CLI required

---

## Slide 4 — Companies Already Winning with Agentic AI

| Company | Implementation | Result |
|---------|---------------|--------|
| Siemens | AI predictive maintenance agents | 30% less unplanned downtime, €1.5B saved/yr |
| BMW | AI quality inspection | 99.5% defect detection, 30% QC cost reduction |
| Bosch | AI-assisted code review | 40% faster software releases, 60% fewer defect escapes |
| Amazon | CodeGuru automated review | 50% reduction in production incidents |
| Microsoft | GitHub Copilot for developers | 55% faster code completion, 46% more PRs merged/day |
| Medtronic | AI in quality management system | 35% shorter CAPA cycle, faster FDA submissions |
| Stryker | AI regulatory documentation | 40% reduction in submission preparation time |
| J&J | AI manufacturing analytics | 20% yield improvement, significant waste reduction |

---

## Slide 5 — SAGE[ai] Is Built on Lean Principles

| Lean Principle | How SAGE[ai] Delivers |
|---------------|----------------------|
| Eliminate Waste (Muda) | Removes 60+ hrs/week of repetitive analysis |
| Continuous Improvement (Kaizen) | Learns from every rejection via RAG memory |
| Error-Proofing (Poka-yoke) | Human-in-the-loop prevents AI errors reaching production |
| Visual Management | Real-time dashboard for all stakeholders |
| Single Piece Flow | Single-lane task queue: deterministic, auditable |
| Respect for People | Amplifies engineers; never replaces human judgment |

---

## Slide 6 — Return on Investment: Month 2 Payback

| Activity | Before SAGE[ai] | After SAGE[ai] | Savings |
|----------|----------------|---------------|---------|
| Error log analysis | 45 min × 10/day × 2 eng = 15 hrs/day | <5 min total | 93% reduction |
| MR code review | 3 hrs × 15 MRs/week = 45 hrs/week | ~4 hrs/week | 91% reduction |
| Compliance reporting | 8 hrs/week manual | 0 hrs (auto-generated) | 100% reduction |
| Knowledge capture | Lost on attrition | Stored in vector memory | Permanent |
| **TOTAL SAVINGS** | **68 hrs/week ≈ 1.7 FTE** | — | **~€120K/year (1 FTE cost)** |

**Conservative estimate: break-even in under 2 months**

---

## Slide 7 — Simple, Auditable, Secure Architecture

```
[Teams/Metabase/GitLab] → [Monitor Agent] → [Task Queue] → [Analyst/Developer Agent] → [Human Gate] → [Audit Trail]
```

**Event Sources:** Teams channels, Metabase dashboards, GitLab issues & MRs, Error log uploads

**AI Agents:** AnalystAgent (log analysis), DeveloperAgent (code review + MR creation), PlannerAgent (orchestration), ReAct multi-step reasoning loop

**Compliance Layer:** Immutable SQLite audit log, ISO 13485 trace IDs on every decision, FDA 21 CFR Part 11 records, Human approval gate on every proposal

---

## Slide 8 — What SAGE[ai] Does Today

| # | Capability | Description |
|---|------------|-------------|
| 🔍 | Log Analysis | AI triage in <60s — severity RED/AMBER/GREEN |
| 🤖 | Code Review | ReAct multi-step reasoning + pipeline check |
| 📋 | MR Creation | Auto-draft from GitLab issue, branch naming |
| 👁 | 24/7 Monitor | Teams, Metabase, GitLab event detection |
| 📊 | Audit Trail | Every decision traceable, ISO 13485 compliant |
| 🌐 | Web Dashboard | No CLI: Dashboard, Analyst, Developer, Audit, Monitor pages |

---

## Slide 9 — Built for Regulated Medical Device Environments

**Standards Met:**
- ✅ ISO 13485:2016 — Quality Management System
- ✅ ISO 14971:2019 — Risk Management (7 risks identified + controlled)
- ✅ IEC 62304:2006 — Medical Device Software Lifecycle
- ✅ FDA 21 CFR Part 11 — Electronic Records
- ✅ FDA Cybersecurity Guidance 2023

**How We Comply:**
- Immutable append-only audit log (no deletes, ever)
- UUID trace ID on every AI decision
- Human approval gate on every proposal
- Full DHF: SRS, RTM, V&V Plan, SOUP Inventory
- Air-gapped LLM option (no cloud dependency)

---

## Slide 10 — SAGE[ai] Gets Better Every Day — From Your Own Team

**3-Step Improvement Cycle:**

1. **USER SUBMITS REQUEST** — 💡 Click 'Request Improvement' on any module
2. **AI PLANS** — Planner Agent decomposes into subtasks, queued for implementation
3. **IMPLEMENTED & LEARNED** — Change deployed; RAG memory updated with new context

*During development: open access for all engineers.*
*Post-release: role-based access control (admin approval required).*

> **The system improves itself through the same AI pipeline it provides to engineers.**

---

## Slide 11 — From Pilot to Production: 3-Month Plan

| Phase | Timeline | Deliverable | Status |
|-------|----------|-------------|--------|
| Phase 1–3: Core System | Done | CLI, Analyst, Developer, Monitor, API | ✅ Complete |
| Phase 4: Web UI + Agentic | Done | Dashboard, ReAct loop, Planner, Regulatory docs | ✅ Complete |
| Phase 5: Pilot Deployment | Month 1 | Production deploy on internal network, engineer training | 🔵 Next |
| Phase 6: Measure & Iterate | Month 2 | KPI tracking, RAG feedback loop, first improvements | 🔵 Planned |
| Phase 7: Scale | Month 3 | Multi-team rollout, Spira integration, executive dashboards | 🔵 Planned |

---

## Slide 12 — Next Steps: 3-Month Pilot Proposal

**What We Need:**
- ✅ Production deployment approval (internal network)
- ✅ 2-day engineering team training session
- ✅ Access to GitLab, Metabase, Teams API credentials
- ✅ Dedicated server or VM (8-core CPU, 32GB RAM, optional GPU)

**What You Get:**
- 📈 KPI dashboard by Week 2
- ⏱ 60+ hours/week reclaimed from manual work
- 📋 ISO 13485 audit trail from day 1
- 🤖 AI that learns from your engineers' expertise
- 💰 ROI positive within 2 months

---

> **SAGE[ai] is not a replacement for your engineers — it's their most productive teammate.**

---
*Generated by generate_presentation.py — SAGE[ai] Business Case | March 2026*
"""


# ---------------------------------------------------------------------------
# Main entry point
# ---------------------------------------------------------------------------

def main():
    os.makedirs(DOCS_DIR, exist_ok=True)

    # -----------------------------------------------------------------------
    # Build presentation
    # -----------------------------------------------------------------------
    prs = Presentation()
    prs.slide_width  = SLIDE_W
    prs.slide_height = SLIDE_H

    print("Building SAGE[ai] Business Case Presentation...")
    print(f"  Output: {PPTX_PATH}")

    builders = [
        ("Slide 01 — Title",                    build_slide_01_title),
        ("Slide 02 — The Problem",              build_slide_02_problem),
        ("Slide 03 — The Solution",             build_slide_03_solution),
        ("Slide 04 — Industry Evidence",        build_slide_04_evidence),
        ("Slide 05 — Lean Alignment",           build_slide_05_lean),
        ("Slide 06 — ROI Analysis",             build_slide_06_roi),
        ("Slide 07 — Architecture",             build_slide_07_architecture),
        ("Slide 08 — Key Capabilities",         build_slide_08_capabilities),
        ("Slide 09 — Compliance",               build_slide_09_compliance),
        ("Slide 10 — Self-Improving System",    build_slide_10_self_improving),
        ("Slide 11 — Roadmap",                  build_slide_11_roadmap),
        ("Slide 12 — The Ask",                  build_slide_12_ask),
    ]

    for label, fn in builders:
        try:
            fn(prs)
            print(f"  [OK] {label}")
        except Exception as exc:
            print(f"  [FAIL] {label} -- ERROR: {exc}")
            raise

    prs.save(PPTX_PATH)
    print(f"\nPresentation saved: {PPTX_PATH}")

    # -----------------------------------------------------------------------
    # Write markdown summary
    # -----------------------------------------------------------------------
    with open(MD_PATH, "w", encoding="utf-8") as fh:
        fh.write(MARKDOWN_CONTENT)
    print(f"Markdown summary saved: {MD_PATH}")
    print("\nDone.")


if __name__ == "__main__":
    main()
