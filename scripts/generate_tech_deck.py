"""
SAGE[ai] — Technical Pitch Deck Generator (Fresh, March 2026)

White + Emerald Green theme — NO blue, NO black.
Generates: docs/SageAI_Tech_Pitch.pptx

Run:  python docs/generate_tech_deck.py
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
import os

# ── SAGE Emerald Theme — NO blue, NO black ──────────────────────────────────
SAGE_GREEN   = RGBColor(0x10, 0xB9, 0x81)   # emerald-500 — primary accent
SAGE_GREEN_D = RGBColor(0x04, 0x78, 0x57)   # emerald-700 — dark accent
SAGE_GREEN_L = RGBColor(0xEC, 0xFD, 0xF5)   # emerald-50  — light bg
SAGE_MINT    = RGBColor(0xD1, 0xFA, 0xE5)   # green-100   — row stripe
SAGE_DARK    = RGBColor(0x1F, 0x2A, 0x37)   # gray-800 (dark charcoal, NOT black)
SAGE_GRAY    = RGBColor(0x4B, 0x55, 0x63)   # gray-600
SAGE_LGRAY   = RGBColor(0x9C, 0xA3, 0xAF)   # gray-400
SAGE_AMBER   = RGBColor(0xF5, 0x9E, 0x0B)   # amber-500
SAGE_RED     = RGBColor(0xEF, 0x44, 0x44)   # red-500
WHITE        = RGBColor(0xFF, 0xFF, 0xFF)
SLIDE_BG     = RGBColor(0xFA, 0xFA, 0xFA)   # near-white
SAGE_WARN_BG = RGBColor(0xFE, 0xF3, 0xC7)   # amber-100

SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)

# ── Helpers ──────────────────────────────────────────────────────────────────

def new_prs():
    prs = Presentation()
    prs.slide_width  = SLIDE_W
    prs.slide_height = SLIDE_H
    return prs

def blank_layout(prs):
    return prs.slide_layouts[6]

def add_rect(slide, l, t, w, h, fill_color=None, line_color=None, line_width=Pt(0)):
    shape = slide.shapes.add_shape(1, l, t, w, h)
    shape.line.width = line_width
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    else:
        shape.fill.background()
    if line_color:
        shape.line.color.rgb = line_color
    else:
        shape.line.fill.background()
    return shape

def add_text_box(slide, text, l, t, w, h,
                 font_size=18, bold=False, color=SAGE_DARK,
                 align=PP_ALIGN.LEFT, wrap=True, italic=False):
    txBox = slide.shapes.add_textbox(l, t, w, h)
    tf = txBox.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return txBox

def add_multiline_box(slide, lines, l, t, w, h,
                      font_size=14, color=SAGE_DARK, line_spacing=1.3,
                      bold_first=False, align=PP_ALIGN.LEFT):
    """Add a text box with multiple lines (list of strings)."""
    txBox = slide.shapes.add_textbox(l, t, w, h)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, line in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.alignment = align
        p.space_after = Pt(2)
        run = p.add_run()
        run.text = line
        run.font.size = Pt(font_size)
        run.font.color.rgb = color
        run.font.bold = bold_first and i == 0
    return txBox

def slide_bg(slide, color=SLIDE_BG):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def header_band(slide, title, subtitle=None):
    """White header band with emerald accent line."""
    add_rect(slide, 0, 0, SLIDE_W, Inches(1.5), fill_color=WHITE)
    add_rect(slide, 0, Inches(1.48), SLIDE_W, Inches(0.04), fill_color=SAGE_GREEN)
    add_text_box(slide, title,
                 Inches(0.5), Inches(0.15), Inches(12.3), Inches(0.8),
                 font_size=30, bold=True, color=SAGE_DARK)
    if subtitle:
        add_text_box(slide, subtitle,
                     Inches(0.5), Inches(0.88), Inches(12.3), Inches(0.5),
                     font_size=15, color=SAGE_GRAY)

def footer_band(slide, text="SAGE[ai] — Open Source (MIT) · github.com/Sumanharapanahalli/sage · March 2026"):
    add_rect(slide, 0, Inches(7.12), SLIDE_W, Inches(0.38), fill_color=WHITE)
    add_rect(slide, 0, Inches(7.1), SLIDE_W, Inches(0.03), fill_color=SAGE_GREEN_L)
    add_text_box(slide, text,
                 Inches(0.3), Inches(7.15), Inches(10), Inches(0.3),
                 font_size=9, color=SAGE_LGRAY)

def bullet_slide(prs, title, subtitle, bullets):
    slide = prs.slides.add_slide(blank_layout(prs))
    slide_bg(slide)
    header_band(slide, title, subtitle)
    top = Inches(1.7)
    for b in bullets:
        if b == "":
            top += Inches(0.15)
            continue
        # Emerald bullet marker
        add_text_box(slide, "●", Inches(0.5), top, Inches(0.3), Inches(0.4),
                     font_size=10, color=SAGE_GREEN, bold=True)
        add_text_box(slide, b, Inches(0.85), top, Inches(11.8), Inches(0.42),
                     font_size=15, color=SAGE_DARK)
        top += Inches(0.48)
    footer_band(slide)
    return slide

def table_slide(prs, title, subtitle, headers, rows, col_widths=None, note=None):
    slide = prs.slides.add_slide(blank_layout(prs))
    slide_bg(slide)
    header_band(slide, title, subtitle)

    n_cols = len(headers)
    n_rows = len(rows)
    if col_widths is None:
        col_widths = [Inches(12.3 / n_cols)] * n_cols

    tbl_top  = Inches(1.7)
    tbl_left = Inches(0.5)
    row_h    = Inches(0.45)
    tbl_h    = row_h * (n_rows + 1)

    table = slide.shapes.add_table(
        n_rows + 1, n_cols, tbl_left, tbl_top,
        sum(col_widths), tbl_h
    ).table

    for ci, w in enumerate(col_widths):
        table.columns[ci].width = w

    # Header row — emerald background, white text
    for ci, h in enumerate(headers):
        cell = table.cell(0, ci)
        cell.fill.solid()
        cell.fill.fore_color.rgb = SAGE_GREEN
        p = cell.text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = h
        run.font.bold = True
        run.font.size = Pt(12)
        run.font.color.rgb = WHITE

    # Data rows — alternating white / mint
    for ri, row in enumerate(rows):
        bg = WHITE if ri % 2 == 0 else SAGE_GREEN_L
        for ci, val in enumerate(row):
            cell = table.cell(ri + 1, ci)
            cell.fill.solid()
            cell.fill.fore_color.rgb = bg
            p = cell.text_frame.paragraphs[0]
            p.alignment = PP_ALIGN.LEFT
            run = p.add_run()
            run.text = str(val)
            run.font.size = Pt(11)
            run.font.color.rgb = SAGE_DARK

    if note:
        top_pos = tbl_top + tbl_h + Inches(0.1)
        add_text_box(slide, note,
                     Inches(0.5), top_pos, Inches(12.3), Inches(0.4),
                     font_size=10, italic=True, color=SAGE_LGRAY)

    footer_band(slide)
    return slide

def stat_card(slide, x, y, w, h, number, label, accent=SAGE_GREEN):
    """Rounded stat card — emerald accent."""
    add_rect(slide, x, y, w, h, fill_color=WHITE, line_color=SAGE_GREEN_L, line_width=Pt(1.5))
    add_text_box(slide, number,
                 x + Inches(0.15), y + Inches(0.15), w - Inches(0.3), Inches(0.6),
                 font_size=32, bold=True, color=accent, align=PP_ALIGN.CENTER)
    add_text_box(slide, label,
                 x + Inches(0.1), y + Inches(0.75), w - Inches(0.2), Inches(0.35),
                 font_size=11, color=SAGE_GRAY, align=PP_ALIGN.CENTER)

def diagram_box(slide, x, y, w, h, text, fill=SAGE_GREEN_L, text_color=SAGE_GREEN_D, font_size=11):
    """Diagram building block — emerald themed."""
    add_rect(slide, x, y, w, h, fill_color=fill, line_color=SAGE_GREEN, line_width=Pt(1))
    add_text_box(slide, text,
                 x + Inches(0.05), y + Inches(0.05), w - Inches(0.1), h - Inches(0.1),
                 font_size=font_size, color=text_color, align=PP_ALIGN.CENTER, bold=True)

def arrow_down(slide, x, y, length=Inches(0.4)):
    """Simple down arrow using text."""
    add_text_box(slide, "↓", x, y, Inches(0.3), length,
                 font_size=18, color=SAGE_GREEN, bold=True, align=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════════════════════════════
#  SLIDE DECK
# ═══════════════════════════════════════════════════════════════════════════════

def build_tech_deck():
    prs = new_prs()

    # ─────────────────────────────────────────────────────────────────────────
    # SLIDE 1 — Cover
    # ─────────────────────────────────────────────────────────────────────────
    slide = prs.slides.add_slide(blank_layout(prs))
    slide_bg(slide, WHITE)
    # Top emerald band
    add_rect(slide, 0, 0, SLIDE_W, Inches(0.06), fill_color=SAGE_GREEN)
    # Bottom emerald band
    add_rect(slide, 0, Inches(7.44), SLIDE_W, Inches(0.06), fill_color=SAGE_GREEN)

    add_text_box(slide, "SAGE[ai]",
                 Inches(1), Inches(1.6), Inches(11), Inches(1.2),
                 font_size=72, bold=True, color=SAGE_GREEN, align=PP_ALIGN.CENTER)
    add_text_box(slide, "Smart Agentic-Guided Empowerment",
                 Inches(1), Inches(2.9), Inches(11), Inches(0.6),
                 font_size=24, color=SAGE_DARK, align=PP_ALIGN.CENTER)
    add_text_box(slide, "The Lean AI Agent Framework for Regulated Industries",
                 Inches(1), Inches(3.55), Inches(11), Inches(0.5),
                 font_size=18, color=SAGE_GRAY, align=PP_ALIGN.CENTER)

    # Stat cards row
    cards = [
        ("136", "API Endpoints"),
        ("30", "UI Pages"),
        ("840+", "Tests Passing"),
        ("20+", "Agent Roles"),
        ("14", "Industry Domains"),
        ("7", "LLM Providers"),
    ]
    card_w = Inches(1.7)
    gap = Inches(0.22)
    total_w = len(cards) * card_w + (len(cards) - 1) * gap
    start_x = (SLIDE_W - total_w) / 2
    for i, (num, label) in enumerate(cards):
        x = start_x + i * (card_w + gap)
        stat_card(slide, x, Inches(4.5), card_w, Inches(1.15), num, label)

    add_text_box(slide, "Open Source (MIT)  ·  Self-Hosted  ·  Zero API Keys Required  ·  Air-Gappable",
                 Inches(1), Inches(6.0), Inches(11), Inches(0.4),
                 font_size=14, color=SAGE_GRAY, align=PP_ALIGN.CENTER)
    add_text_box(slide, "github.com/Sumanharapanahalli/sage",
                 Inches(1), Inches(6.45), Inches(11), Inches(0.35),
                 font_size=13, italic=True, color=SAGE_GREEN, align=PP_ALIGN.CENTER)

    # ─────────────────────────────────────────────────────────────────────────
    # SLIDE 2 — The Vision
    # ─────────────────────────────────────────────────────────────────────────
    slide = prs.slides.add_slide(blank_layout(prs))
    slide_bg(slide, WHITE)
    header_band(slide, "The Vision", "1 human + AI agent team = billion-dollar operation")

    add_text_box(slide, "\"Every company function becomes an agent role.\nThe founder reviews only decisions that require human judgment.\nEverything else runs autonomously with a complete audit trail.\"",
                 Inches(1), Inches(2.0), Inches(11), Inches(1.5),
                 font_size=20, italic=True, color=SAGE_GREEN_D, align=PP_ALIGN.CENTER)

    # Three pillars
    pillars = [
        ("Lean Development", "Eliminate waste, shorten feedback\nloops, amplify human judgment"),
        ("Agentic AI", "Autonomous agents that propose,\nlearn, and compound intelligence"),
        ("Human-in-the-Loop", "Every AI write action requires\nhuman approval — always"),
    ]
    pw = Inches(3.4)
    pgap = Inches(0.4)
    px_start = (SLIDE_W - 3*pw - 2*pgap) / 2
    for i, (title, desc) in enumerate(pillars):
        x = px_start + i * (pw + pgap)
        add_rect(slide, x, Inches(3.9), pw, Inches(2.2),
                 fill_color=SAGE_GREEN_L, line_color=SAGE_GREEN, line_width=Pt(1))
        add_text_box(slide, title,
                     x + Inches(0.2), Inches(4.1), pw - Inches(0.4), Inches(0.5),
                     font_size=17, bold=True, color=SAGE_GREEN_D, align=PP_ALIGN.CENTER)
        add_text_box(slide, desc,
                     x + Inches(0.2), Inches(4.65), pw - Inches(0.4), Inches(1.2),
                     font_size=13, color=SAGE_GRAY, align=PP_ALIGN.CENTER)

    footer_band(slide)

    # ─────────────────────────────────────────────────────────────────────────
    # SLIDE 3 — The Five Laws
    # ─────────────────────────────────────────────────────────────────────────
    bullet_slide(prs,
        "The Five Laws of SAGE",
        "The engineering principles that every feature is built on",
        [
            "1. Agents propose. Humans decide. Always.",
            "    The HITL approval gate is not bureaucracy — it IS the product",
            "",
            "2. Eliminate waste at every layer.",
            "    If it can be automated correctly, it must be. If not, it stays human.",
            "",
            "3. Compounding intelligence over cold-start.",
            "    Every correction feeds the vector store — agents improve without retraining",
            "",
            "4. Vertical slices, not horizontal layers.",
            "    Every task produces a working, reviewable end-to-end slice of value",
            "",
            "5. Atomic verification is non-negotiable.",
            "    Every agent action has a defined verification step. Build in the check.",
        ]
    )

    # ─────────────────────────────────────────────────────────────────────────
    # SLIDE 4 — Architecture Overview (Diagram)
    # ─────────────────────────────────────────────────────────────────────────
    slide = prs.slides.add_slide(blank_layout(prs))
    slide_bg(slide)
    header_band(slide, "Architecture Overview", "Data flows one way: UI → API → Agents → LLM → Audit Log")

    # Left column: Solution YAML
    diagram_box(slide, Inches(0.4), Inches(1.8), Inches(2.2), Inches(0.5),
                "solutions/<name>/", fill=SAGE_MINT)
    diagram_box(slide, Inches(0.4), Inches(2.4), Inches(2.2), Inches(0.4),
                "project.yaml", fill=WHITE, text_color=SAGE_GRAY, font_size=10)
    diagram_box(slide, Inches(0.4), Inches(2.85), Inches(2.2), Inches(0.4),
                "prompts.yaml", fill=WHITE, text_color=SAGE_GRAY, font_size=10)
    diagram_box(slide, Inches(0.4), Inches(3.3), Inches(2.2), Inches(0.4),
                "tasks.yaml", fill=WHITE, text_color=SAGE_GRAY, font_size=10)

    # Center column: Framework
    diagram_box(slide, Inches(3.2), Inches(1.8), Inches(3.2), Inches(0.5),
                "React Dashboard (30 pages)", fill=SAGE_GREEN_L)
    arrow_down(slide, Inches(4.65), Inches(2.35))
    diagram_box(slide, Inches(3.2), Inches(2.8), Inches(3.2), Inches(0.5),
                "FastAPI — 136 endpoints", fill=SAGE_GREEN_L)
    arrow_down(slide, Inches(4.65), Inches(3.35))
    diagram_box(slide, Inches(3.2), Inches(3.8), Inches(3.2), Inches(0.5),
                "Agent Engine (20 roles)", fill=SAGE_GREEN_L)
    arrow_down(slide, Inches(4.65), Inches(4.35))
    diagram_box(slide, Inches(3.2), Inches(4.8), Inches(3.2), Inches(0.5),
                "LLM Gateway (7 providers)", fill=SAGE_GREEN_L)

    # Right column: Integrations & Memory
    diagram_box(slide, Inches(7.0), Inches(1.8), Inches(2.6), Inches(0.5),
                "Integrations (13)", fill=SAGE_GREEN_L)
    items_right = ["LangGraph", "AutoGen", "MCP Tools", "Slack", "Temporal",
                   "n8n", "OpenSWE", "LangChain"]
    top = Inches(2.4)
    for item in items_right:
        add_text_box(slide, f"· {item}", Inches(7.1), top, Inches(2.4), Inches(0.28),
                     font_size=10, color=SAGE_GRAY)
        top += Inches(0.25)

    # Far right: Memory
    diagram_box(slide, Inches(10.2), Inches(1.8), Inches(2.6), Inches(0.5),
                "Memory & Audit", fill=SAGE_GREEN_L)
    mem_items = ["ChromaDB Vector Store", "SQLite Audit Log", "Compounding Feedback",
                 "Per-Solution .sage/"]
    top = Inches(2.4)
    for item in mem_items:
        add_text_box(slide, f"· {item}", Inches(10.3), top, Inches(2.4), Inches(0.28),
                     font_size=10, color=SAGE_GRAY)
        top += Inches(0.25)

    # The Lean Loop at bottom
    add_rect(slide, Inches(0.4), Inches(5.5), Inches(12.5), Inches(1.3),
             fill_color=SAGE_GREEN_L, line_color=SAGE_GREEN, line_width=Pt(1))
    add_text_box(slide, "The SAGE Lean Loop",
                 Inches(0.6), Inches(5.55), Inches(3), Inches(0.35),
                 font_size=12, bold=True, color=SAGE_GREEN_D)
    phases = ["① SURFACE", "② CONTEXTUALIZE", "③ PROPOSE", "④ DECIDE (HITL)", "⑤ COMPOUND"]
    phase_w = Inches(2.3)
    for i, phase in enumerate(phases):
        x = Inches(0.6) + i * phase_w
        add_text_box(slide, phase,
                     x, Inches(5.95), phase_w, Inches(0.35),
                     font_size=12, bold=True, color=SAGE_GREEN_D, align=PP_ALIGN.CENTER)
    add_text_box(slide, "→  Phase 5 feeds back into Phase 2. Every human decision compounds intelligence for the next task.",
                 Inches(0.6), Inches(6.35), Inches(12), Inches(0.35),
                 font_size=11, italic=True, color=SAGE_GRAY)

    footer_band(slide)

    # ─────────────────────────────────────────────────────────────────────────
    # SLIDE 5 — By the Numbers
    # ─────────────────────────────────────────────────────────────────────────
    slide = prs.slides.add_slide(blank_layout(prs))
    slide_bg(slide, WHITE)
    header_band(slide, "SAGE by the Numbers", "Everything built. Everything tested. Everything live.")

    stats = [
        ("136", "API Endpoints", "Full REST API — config, agents,\nknowledge, build, eval, audit"),
        ("30", "UI Pages", "React 18 + TypeScript dashboard\nwith dark mode + theming"),
        ("20", "Agent Roles", "5 functional teams — hire\nany role via YAML or LLM"),
        ("14", "Industry Domains", "Auto-detected from description\nwith compliance injection"),
        ("32", "Build Task Types", "SW, HW, firmware, QA, legal,\nbusiness, design, ops, ML"),
        ("7", "LLM Providers", "Gemini, Claude, Ollama, local\nGGUF, generic CLI + pool"),
        ("840+", "Tests Passing", "Unit + integration + system\n+ browser E2E (Playwright)"),
        ("16+", "Solution Templates", "Medical, automotive, avionics,\nfintech, IoT, SaaS, games..."),
        ("13", "Integration Modules", "LangGraph, AutoGen, MCP,\nSlack, Temporal, OpenSWE..."),
        ("20+", "Compliance Standards", "IEC 62304, ISO 13485, ISO 26262,\nDO-178C, PCI DSS, HIPAA..."),
        ("5", "Agentic Patterns", "Sequential, parallel, critic,\nhierarchical, swarm"),
        ("12+", "Integration Phases", "Langfuse → LangChain → MCP →\nn8n → Slack → Build Orch."),
    ]
    card_w = Inches(2.8)
    card_h = Inches(1.5)
    cols = 4
    gap_x = Inches(0.3)
    gap_y = Inches(0.2)
    total_grid_w = cols * card_w + (cols - 1) * gap_x
    start_x = (SLIDE_W - total_grid_w) / 2
    start_y = Inches(1.7)

    for idx, (num, label, desc) in enumerate(stats):
        r = idx // cols
        c = idx % cols
        x = start_x + c * (card_w + gap_x)
        y = start_y + r * (card_h + gap_y)
        add_rect(slide, x, y, card_w, card_h,
                 fill_color=SAGE_GREEN_L, line_color=SAGE_MINT, line_width=Pt(1))
        add_text_box(slide, num,
                     x + Inches(0.1), y + Inches(0.08), card_w - Inches(0.2), Inches(0.45),
                     font_size=26, bold=True, color=SAGE_GREEN)
        add_text_box(slide, label,
                     x + Inches(0.1), y + Inches(0.5), card_w - Inches(0.2), Inches(0.3),
                     font_size=12, bold=True, color=SAGE_DARK)
        add_text_box(slide, desc,
                     x + Inches(0.1), y + Inches(0.82), card_w - Inches(0.2), Inches(0.6),
                     font_size=9, color=SAGE_GRAY)

    footer_band(slide)

    # ─────────────────────────────────────────────────────────────────────────
    # SLIDE 6 — HITL Approval System
    # ─────────────────────────────────────────────────────────────────────────
    slide = prs.slides.add_slide(blank_layout(prs))
    slide_bg(slide)
    header_band(slide, "Human-in-the-Loop Approval System",
                "Every AI write action is risk-ranked and requires human approval before execution")

    # Two-tier table
    table_data_tiers = [
        ["Framework Control", "config/switch, llm/switch, config/modules", "Executes immediately", "No proposal"],
        ["Agent Proposals", "yaml_edit, code_diff, knowledge_delete, agent_hire", "Requires HITL approval", "Creates Proposal"],
    ]
    table_slide_inline = slide  # reuse

    # Risk tiers
    risk_tiers = [
        ("INFORMATIONAL", "Read-only query", "1 hour", SAGE_GREEN_L, SAGE_GREEN_D),
        ("EPHEMERAL", "LLM provider switch", "8 hours", SAGE_GREEN_L, SAGE_GREEN_D),
        ("STATEFUL", "Knowledge base edit", "7 days", SAGE_WARN_BG, SAGE_AMBER),
        ("EXTERNAL", "GitLab MR creation", "14 days", SAGE_WARN_BG, SAGE_AMBER),
        ("DESTRUCTIVE", "Drop data / delete agent", "Never expires", RGBColor(0xFE, 0xE2, 0xE2), SAGE_RED),
    ]

    top = Inches(1.8)
    add_text_box(slide, "5 Risk Tiers — proposals auto-expire based on risk level",
                 Inches(0.5), top, Inches(12), Inches(0.35),
                 font_size=14, bold=True, color=SAGE_GREEN_D)
    top += Inches(0.5)

    for tier, example, expiry, bg_color, txt_color in risk_tiers:
        add_rect(slide, Inches(0.5), top, Inches(12.3), Inches(0.55),
                 fill_color=bg_color, line_color=SAGE_GREEN_L, line_width=Pt(0.5))
        add_text_box(slide, tier,
                     Inches(0.7), top + Inches(0.08), Inches(2.2), Inches(0.4),
                     font_size=13, bold=True, color=txt_color)
        add_text_box(slide, example,
                     Inches(3.0), top + Inches(0.08), Inches(4), Inches(0.4),
                     font_size=12, color=SAGE_DARK)
        add_text_box(slide, f"Expiry: {expiry}",
                     Inches(8.5), top + Inches(0.08), Inches(4), Inches(0.4),
                     font_size=12, color=SAGE_GRAY)
        top += Inches(0.6)

    # Flow diagram
    top += Inches(0.3)
    add_text_box(slide, "Approval Flow:   Agent proposes  →  Proposal created (risk-ranked)  →  /approvals inbox  →  Human reviews  →  Approve / Reject with feedback  →  Audit log + vector memory updated",
                 Inches(0.5), top, Inches(12.3), Inches(0.5),
                 font_size=12, color=SAGE_DARK)

    top += Inches(0.55)
    add_text_box(slide, "Slack integration:  Block Kit approval cards → /webhook/slack callback → two-way approval without leaving Slack",
                 Inches(0.5), top, Inches(12.3), Inches(0.4),
                 font_size=11, italic=True, color=SAGE_GRAY)

    footer_band(slide)

    # ─────────────────────────────────────────────────────────────────────────
    # SLIDE 7 — Build Orchestrator (0→1→N)
    # ─────────────────────────────────────────────────────────────────────────
    slide = prs.slides.add_slide(blank_layout(prs))
    slide_bg(slide)
    header_band(slide, "Build Orchestrator — Idea to Working Product",
                "Describe a product → get a production-ready codebase with agents, tests, and CI/CD")

    # Flow diagram
    steps = [
        ("Plain-language\ndescription", SAGE_GREEN_L),
        ("Domain Detection\n(14 domains)", SAGE_GREEN_L),
        ("Task Decomposition\n(32 task types)", SAGE_GREEN_L),
        ("Agent Assignment\n(Q-learning router)", SAGE_GREEN_L),
        ("Critic Review\n(actor-critic loop)", SAGE_MINT),
        ("HITL Gate\n(founder approves)", SAGE_WARN_BG),
        ("Parallel Build\n(wave execution)", SAGE_GREEN_L),
        ("Anti-Drift\nCheckpoints", SAGE_MINT),
        ("Working\nCodebase", SAGE_GREEN_L),
    ]
    step_w = Inches(1.25)
    step_h = Inches(0.9)
    sgap = Inches(0.1)
    total_steps_w = len(steps) * step_w + (len(steps) - 1) * sgap
    sx_start = (SLIDE_W - total_steps_w) / 2
    sy = Inches(1.8)

    for i, (text, bg) in enumerate(steps):
        x = sx_start + i * (step_w + sgap)
        add_rect(slide, x, sy, step_w, step_h,
                 fill_color=bg, line_color=SAGE_GREEN, line_width=Pt(1))
        add_text_box(slide, text,
                     x + Inches(0.03), sy + Inches(0.05), step_w - Inches(0.06), step_h - Inches(0.1),
                     font_size=9, color=SAGE_GREEN_D, align=PP_ALIGN.CENTER, bold=True)
        if i < len(steps) - 1:
            add_text_box(slide, "→",
                         x + step_w, sy + Inches(0.25), sgap, Inches(0.4),
                         font_size=14, color=SAGE_GREEN, bold=True, align=PP_ALIGN.CENTER)

    # Details below
    top = Inches(3.0)

    # Workforce teams
    add_text_box(slide, "5 Workforce Teams — 20 Agent Roles",
                 Inches(0.5), top, Inches(6), Inches(0.35),
                 font_size=14, bold=True, color=SAGE_GREEN_D)
    teams = [
        "Engineering:  developer, qa_engineer, system_tester, devops, localization",
        "Analysis:     analyst, business_analyst, financial_analyst, data_scientist",
        "Design:       ux_designer, product_manager",
        "Compliance:   regulatory_specialist, legal_advisor, safety_engineer",
        "Operations:   operations_manager, technical_writer, marketing_strategist",
    ]
    top += Inches(0.4)
    for t in teams:
        add_text_box(slide, t, Inches(0.7), top, Inches(6), Inches(0.3),
                     font_size=10, color=SAGE_DARK)
        top += Inches(0.28)

    # Right side: domain detection
    add_text_box(slide, "14 Industry Domains Auto-Detected",
                 Inches(7), Inches(3.0), Inches(6), Inches(0.35),
                 font_size=14, bold=True, color=SAGE_GREEN_D)
    domains = [
        "Medical Device — IEC 62304, ISO 13485",
        "Automotive — ISO 26262, AUTOSAR",
        "Avionics — DO-178C, ARP4754A",
        "Robotics — ISO 10218",
        "FinTech — PCI DSS, SOX",
        "IoT — IEC 62443",
        "ML/AI — model governance",
        "SaaS / Consumer / Enterprise / E-commerce",
        "Healthcare SW / EdTech / Hardware",
    ]
    top = Inches(3.4)
    for d in domains:
        add_text_box(slide, f"· {d}", Inches(7.2), top, Inches(5.5), Inches(0.28),
                     font_size=10, color=SAGE_DARK)
        top += Inches(0.27)

    # HITL levels
    add_text_box(slide, "3 HITL Levels:  minimal (final only)  ·  standard (plan + final)  ·  strict (plan + per-component + final)",
                 Inches(0.5), Inches(5.8), Inches(12), Inches(0.35),
                 font_size=12, bold=True, color=SAGE_GREEN_D)

    add_text_box(slide, "3-tier degradation:  OpenSWE runner → LLM direct → template scaffold.  Always produces a buildable output.",
                 Inches(0.5), Inches(6.2), Inches(12), Inches(0.35),
                 font_size=11, italic=True, color=SAGE_GRAY)

    add_text_box(slide, "Adaptive Router:  Q-learning agent routing — learns which agent performs best per task type. 3+ observations before overriding defaults.",
                 Inches(0.5), Inches(6.5), Inches(12), Inches(0.35),
                 font_size=11, italic=True, color=SAGE_GRAY)

    footer_band(slide)

    # ─────────────────────────────────────────────────────────────────────────
    # SLIDE 8 — Multi-LLM Provider Pool
    # ─────────────────────────────────────────────────────────────────────────
    slide = prs.slides.add_slide(blank_layout(prs))
    slide_bg(slide)
    header_band(slide, "Multi-LLM Provider Support",
                "7 providers. Zero mandatory API keys. Parallel generation with 4 strategies.")

    # Provider table
    providers = [
        ["Gemini CLI", "npm install -g @google/gemini-cli", "Free Google login", "gemini-2.5-flash/pro"],
        ["Claude Code CLI", "npm install -g @anthropic-ai/claude-code", "Free login", "claude-sonnet-4-6/opus-4-5"],
        ["Ollama", "ollama.com → ollama serve → ollama pull", "Fully offline", "llama3.2, mistral, qwen2.5"],
        ["Local GGUF", "pip install llama-cpp-python + model", "Air-gapped, GPU-direct", "Any .gguf model"],
        ["Claude API", "Set ANTHROPIC_API_KEY", "Only paid option", "claude-3.5-sonnet"],
        ["Generic CLI", "Set generic_cli_path in config", "Any CLI tool", "Custom"],
    ]

    table_slide_helper(slide, Inches(0.5), Inches(1.7),
        ["Provider", "Setup", "Auth", "Models"],
        providers,
        [Inches(2.0), Inches(4.5), Inches(2.5), Inches(3.3)])

    # Parallel strategies
    top = Inches(4.7)
    add_text_box(slide, "ProviderPool — Parallel Generation Strategies",
                 Inches(0.5), top, Inches(12), Inches(0.35),
                 font_size=14, bold=True, color=SAGE_GREEN_D)
    strategies = [
        ("Voting", "Send to multiple LLMs, take majority consensus"),
        ("Fastest", "Send to all, return first response"),
        ("Fallback", "Try providers in order, fail over automatically"),
        ("Quality", "Send to all, return longest/richest response"),
    ]
    top += Inches(0.45)
    for name, desc in strategies:
        add_text_box(slide, f"● {name}", Inches(0.7), top, Inches(2), Inches(0.3),
                     font_size=12, bold=True, color=SAGE_GREEN_D)
        add_text_box(slide, desc, Inches(2.7), top, Inches(9), Inches(0.3),
                     font_size=12, color=SAGE_DARK)
        top += Inches(0.35)

    # Teacher-Student
    top += Inches(0.15)
    add_text_box(slide, "Teacher-Student Distillation:  Heavy model (GPT-4/Claude) → confidence scoring → auto-capture training data → student handles 80% of tasks over time",
                 Inches(0.5), top, Inches(12), Inches(0.4),
                 font_size=11, italic=True, color=SAGE_GRAY)

    footer_band(slide)

    # ─────────────────────────────────────────────────────────────────────────
    # SLIDE 9 — Framework Comparison
    # ─────────────────────────────────────────────────────────────────────────
    table_slide(prs,
        "Framework Comparison",
        "SAGE vs LangGraph, CrewAI, AutoGen, Semantic Kernel, Dify, n8n",
        ["Feature", "SAGE", "LangGraph", "CrewAI", "AutoGen", "Dify", "n8n"],
        [
            ["Mandatory HITL approval",     "✅", "⚠️ opt-in", "❌", "❌", "⚠️", "⚠️"],
            ["Immutable audit log",          "✅", "❌", "❌", "❌", "⚠️", "❌"],
            ["Compliance standards (ISO/FDA)","✅", "❌", "❌", "❌", "❌", "❌"],
            ["Compounding memory (RAG)",     "✅", "⚠️", "❌", "⚠️", "❌", "❌"],
            ["YAML-first agent config",      "✅", "❌", "❌", "❌", "✅", "✅"],
            ["Offline / air-gapped",         "✅", "⚠️", "⚠️", "⚠️", "❌", "❌"],
            ["React dashboard included",     "✅", "❌", "❌", "❌", "✅", "✅"],
            ["Multi-tenant isolation",       "✅", "❌", "❌", "❌", "⚠️", "⚠️"],
            ["Build orchestrator (0→N)",     "✅", "❌", "❌", "❌", "❌", "❌"],
            ["Multi-LLM provider pool",      "✅", "⚠️", "⚠️", "⚠️", "❌", "⚠️"],
            ["MCP tool standard",            "✅", "⚠️", "❌", "❌", "❌", "❌"],
            ["SWE agent (open-swe)",         "✅", "⚠️", "❌", "⚠️", "❌", "❌"],
            ["Large community",              "⚠️", "✅", "✅", "✅", "✅", "✅"],
        ],
        col_widths=[Inches(3.2), Inches(1.2), Inches(1.6), Inches(1.3), Inches(1.5), Inches(1.2), Inches(1.2)],
        note="✅ = First-class support  ·  ⚠️ = Partial / DIY  ·  ❌ = Not supported"
    )

    # ─────────────────────────────────────────────────────────────────────────
    # SLIDE 10 — Compliance & Regulated Industries
    # ─────────────────────────────────────────────────────────────────────────
    slide = prs.slides.add_slide(blank_layout(prs))
    slide_bg(slide)
    header_band(slide, "Built for Regulated Industries",
                "The only open-source agent framework with compliance infrastructure built-in")

    # Three columns
    col_w = Inches(3.8)
    col_gap = Inches(0.35)
    cols_start = (SLIDE_W - 3*col_w - 2*col_gap) / 2

    sections = [
        ("Immutable Audit Trail", [
            "SQLite per-solution (.sage/audit_log.db)",
            "trace_id: SHA256 unique per decision",
            "timestamp, actor, action_type, I/O",
            "Append-only — no update/delete path",
            "FDA 21 CFR Part 11 compatible",
            "ISO 13485 clause 4.2.5 ready",
            "Full decision reconstruction",
        ]),
        ("Human Approval Gate", [
            "Mandatory at framework level",
            "Cannot be bypassed — not for demos",
            "5 risk tiers with auto-expiry",
            "Slack two-way approval (Block Kit)",
            "Rejection + feedback stored",
            "trace_id links approval to proposal",
            "Compliance mode: interrupt_before",
        ]),
        ("Compounding Intelligence", [
            "Every rejection → vector store",
            "Future queries retrieve corrections",
            "Improves without model retraining",
            "Satisfies ISO 13485 clause 8.5",
            "Memento principle (arXiv 2508.16153)",
            "Audit log = training signal",
            "Measurable quality improvement",
        ]),
    ]

    for i, (title, items) in enumerate(sections):
        x = cols_start + i * (col_w + col_gap)
        add_rect(slide, x, Inches(1.8), col_w, Inches(4.2),
                 fill_color=WHITE, line_color=SAGE_GREEN, line_width=Pt(1))
        add_text_box(slide, title,
                     x + Inches(0.15), Inches(1.9), col_w - Inches(0.3), Inches(0.4),
                     font_size=14, bold=True, color=SAGE_GREEN_D, align=PP_ALIGN.CENTER)
        top = Inches(2.4)
        for item in items:
            add_text_box(slide, f"· {item}",
                         x + Inches(0.2), top, col_w - Inches(0.4), Inches(0.3),
                         font_size=10, color=SAGE_DARK)
            top += Inches(0.26)

    # Standards bar
    add_rect(slide, Inches(0.4), Inches(6.15), Inches(12.5), Inches(0.7),
             fill_color=SAGE_GREEN_L, line_color=SAGE_GREEN, line_width=Pt(0.5))
    add_text_box(slide, "Standards Supported:  IEC 62304 · ISO 13485 · ISO 14971 · FDA 21 CFR 820 · ISO 26262 · DO-178C · DO-254 · PCI DSS · HIPAA · SOC 2 · ISO 27001 · FERPA · COPPA · IEC 62443 · WCAG 2.1",
                 Inches(0.6), Inches(6.25), Inches(12), Inches(0.5),
                 font_size=10, color=SAGE_GREEN_D)

    footer_band(slide)

    # ─────────────────────────────────────────────────────────────────────────
    # SLIDE 11 — Key Features Grid
    # ─────────────────────────────────────────────────────────────────────────
    slide = prs.slides.add_slide(blank_layout(prs))
    slide_bg(slide)
    header_band(slide, "Key Capabilities", "Everything built, tested, and production-ready")

    features = [
        ("Conversational Onboarding", "Point at a repo or describe your domain\n→ 3 YAML files generated by LLM"),
        ("SWE Agent (open-swe)", "Submit task → autonomous coding →\ntests → PR → founder reviews"),
        ("Visual Workflows", "LangGraph → auto-generated Mermaid\ndiagrams on /workflows page"),
        ("Parallel Task Execution", "Wave scheduling — concurrent agents,\nwait only where dependencies exist"),
        ("Action-Aware Chat", "Chat classifies intent → shows\nconfirmation card → executes on confirm"),
        ("Eval / Benchmarking", "YAML test suites per solution,\nkeyword scoring, SQLite history"),
        ("Knowledge Base CRUD", "Vector store per solution,\nbulk import, semantic search"),
        ("Org Chart + Traceability", "Agent hierarchy, live status,\ndaily task counts, audit drill-down"),
        ("HIL Testing", "Hardware-in-the-loop — flash firmware,\nrun tests, generate regulatory evidence"),
        ("Agent Budgets", "Monthly call ceilings per agent\nrole — declared in project.yaml"),
        ("Git Worktrees", "Isolated worktree per code_diff\nproposal — concurrent, no conflicts"),
        ("Theme System", "Per-solution branding via CSS vars\n— project.yaml → ThemeProvider"),
    ]

    card_w = Inches(3.0)
    card_h = Inches(1.2)
    cols = 4
    gap_x = Inches(0.2)
    gap_y = Inches(0.15)
    total_grid_w = cols * card_w + (cols - 1) * gap_x
    start_x = (SLIDE_W - total_grid_w) / 2
    start_y = Inches(1.7)

    for idx, (title, desc) in enumerate(features):
        r = idx // cols
        c = idx % cols
        x = start_x + c * (card_w + gap_x)
        y = start_y + r * (card_h + gap_y)
        add_rect(slide, x, y, card_w, card_h,
                 fill_color=SAGE_GREEN_L, line_color=SAGE_GREEN, line_width=Pt(0.5))
        add_text_box(slide, title,
                     x + Inches(0.12), y + Inches(0.08), card_w - Inches(0.24), Inches(0.35),
                     font_size=12, bold=True, color=SAGE_GREEN_D)
        add_text_box(slide, desc,
                     x + Inches(0.12), y + Inches(0.45), card_w - Inches(0.24), Inches(0.7),
                     font_size=9, color=SAGE_GRAY)

    footer_band(slide)

    # ─────────────────────────────────────────────────────────────────────────
    # SLIDE 12 — Decision Guide
    # ─────────────────────────────────────────────────────────────────────────
    slide = prs.slides.add_slide(blank_layout(prs))
    slide_bg(slide)
    header_band(slide, "When to Choose SAGE", "Decision guide for teams evaluating agent frameworks")

    decisions = [
        ("Is the project in a regulated industry?", "→ SAGE", SAGE_GREEN_D),
        ("Is auditability required for any AI decision?", "→ SAGE", SAGE_GREEN_D),
        ("Does the team have < 10 engineers?", "→ SAGE", SAGE_GREEN_D),
        ("Is self-hosted / air-gapped required?", "→ SAGE (Ollama)", SAGE_GREEN_D),
        ("Is pure orchestration research the goal?", "→ LangGraph", SAGE_GRAY),
        ("Is speed of prototyping the top priority, no compliance?", "→ CrewAI", SAGE_GRAY),
        ("Is Docker-sandboxed code execution the core feature?", "→ AutoGen", SAGE_GRAY),
        ("Is the team Azure / Microsoft committed?", "→ Semantic Kernel", SAGE_GRAY),
        ("Are non-technical users building workflows?", "→ Dify", SAGE_GRAY),
    ]

    top = Inches(1.8)
    for question, answer, color in decisions:
        is_sage = "SAGE" in answer and "LangGraph" not in answer and "CrewAI" not in answer
        bg = SAGE_GREEN_L if is_sage else WHITE
        add_rect(slide, Inches(0.5), top, Inches(12.3), Inches(0.5),
                 fill_color=bg, line_color=SAGE_GREEN_L, line_width=Pt(0.5))
        add_text_box(slide, question,
                     Inches(0.7), top + Inches(0.07), Inches(8.5), Inches(0.35),
                     font_size=13, color=SAGE_DARK)
        add_text_box(slide, answer,
                     Inches(9.5), top + Inches(0.07), Inches(3), Inches(0.35),
                     font_size=13, bold=True, color=color)
        top += Inches(0.55)

    # Industry recommendations
    top += Inches(0.2)
    add_text_box(slide, "By Industry:  Medical → SAGE  |  Manufacturing → SAGE  |  Legal → SAGE  |  Fintech → SAGE/LangGraph  |  DevOps → LangGraph  |  Marketing → CrewAI/Dify",
                 Inches(0.5), top, Inches(12.3), Inches(0.4),
                 font_size=11, color=SAGE_GRAY)

    footer_band(slide)

    # ─────────────────────────────────────────────────────────────────────────
    # SLIDE 13 — Integration Architecture
    # ─────────────────────────────────────────────────────────────────────────
    table_slide(prs,
        "Integration Architecture — 12+ Phases",
        "Progressive integration strategy — each phase adds capability without breaking existing ones",
        ["Phase", "Feature", "Key Files", "Status"],
        [
            ["0",    "Langfuse observability",       "llm_gateway.py",         "✅ Live"],
            ["1",    "LlamaIndex + LangChain + mem0","vector_store.py, langchain_tools.py", "✅ Live"],
            ["1.5",  "MCP tool registry",            "mcp_registry.py",        "✅ Live"],
            ["2",    "n8n webhook receiver",         "api.py /webhook/n8n",    "✅ Live"],
            ["3",    "LangGraph orchestration",      "langgraph_runner.py",    "✅ Live"],
            ["4",    "AutoGen code agent",           "autogen_runner.py",      "✅ Live"],
            ["5",    "SSE streaming",                "api.py /analyze/stream", "✅ Live"],
            ["6",    "Onboarding wizard",            "onboarding.py",          "✅ Live"],
            ["7",    "Knowledge base CRUD",          "vector_store.py",        "✅ Live"],
            ["8",    "Slack two-way approval",       "slack_approver.py",      "✅ Live"],
            ["9",    "Eval / benchmarking",          "eval_runner.py",         "✅ Live"],
            ["10",   "Multi-tenant isolation",       "tenant.py",             "✅ Live"],
            ["11",   "Temporal durable workflows",   "temporal_runner.py",     "✅ Live"],
            ["12",   "Build Orchestrator (0→1→N)",   "build_orchestrator.py",  "✅ Live"],
        ],
        col_widths=[Inches(0.8), Inches(3.5), Inches(4.5), Inches(1.2)],
    )

    # ─────────────────────────────────────────────────────────────────────────
    # SLIDE 14 — UI Dashboard Tour
    # ─────────────────────────────────────────────────────────────────────────
    slide = prs.slides.add_slide(blank_layout(prs))
    slide_bg(slide)
    header_band(slide, "30-Page React Dashboard",
                "5-area accordion nav · solution switching · per-solution theming · onboarding tour")

    areas = [
        ("Work", "●", SAGE_RED, "Dashboard · Approvals · Queue · Live Console"),
        ("Intelligence", "●", RGBColor(0xA7, 0x8B, 0xFA), "Analyst · Developer · Monitor · Agents · Improvements · Workflows · Goals"),
        ("Knowledge", "●", SAGE_GREEN, "Knowledge · Activity · Audit · Costs"),
        ("Organization", "●", RGBColor(0x6B, 0x7E, 0x80), "Org Graph · Onboarding"),
        ("Admin", "●", SAGE_GRAY, "LLM · YAML Editor · Access Control · Integrations · Settings"),
    ]

    top = Inches(1.8)
    for area, marker, color, pages in areas:
        add_rect(slide, Inches(0.5), top, Inches(12.3), Inches(0.6),
                 fill_color=WHITE, line_color=SAGE_GREEN_L, line_width=Pt(0.5))
        add_text_box(slide, marker, Inches(0.65), top + Inches(0.1), Inches(0.3), Inches(0.4),
                     font_size=14, bold=True, color=color)
        add_text_box(slide, area, Inches(0.95), top + Inches(0.1), Inches(2), Inches(0.4),
                     font_size=14, bold=True, color=SAGE_DARK)
        add_text_box(slide, pages, Inches(3.0), top + Inches(0.12), Inches(9.5), Inches(0.4),
                     font_size=12, color=SAGE_GRAY)
        top += Inches(0.65)

    top += Inches(0.3)
    ui_features = [
        "Solution Rail:  44px icon column — jump between solutions instantly (2-letter avatars)",
        "Stats Strip:  APPROVALS (red) · QUEUED (amber) · AGENTS (green) — live counts, 10s polling",
        "Theme System:  Each solution has its own accent color — auto-applied from project.yaml",
        "Command Palette:  Cmd+K → jump to any page, search agents, switch solutions",
        "Onboarding Tour:  6-stop spotlight tour for new solutions — auto-triggers on first load",
        "Dark Mode:  Toggle via user preferences — full dark theme with proper contrast",
    ]
    for feat in ui_features:
        add_text_box(slide, f"● {feat}",
                     Inches(0.5), top, Inches(12.3), Inches(0.32),
                     font_size=11, color=SAGE_DARK)
        top += Inches(0.33)

    footer_band(slide)

    # ─────────────────────────────────────────────────────────────────────────
    # SLIDE 15 — Quick Start
    # ─────────────────────────────────────────────────────────────────────────
    slide = prs.slides.add_slide(blank_layout(prs))
    slide_bg(slide, WHITE)
    header_band(slide, "Quick Start — 3 Commands", "From git clone to running dashboard in under 5 minutes")

    # Code block
    code_bg = RGBColor(0xF1, 0xF5, 0xF9)  # slate-100 (not black!)
    add_rect(slide, Inches(1), Inches(1.8), Inches(11), Inches(3.5),
             fill_color=code_bg, line_color=SAGE_GREEN, line_width=Pt(1))

    code_lines = [
        "$ git clone https://github.com/Sumanharapanahalli/sage",
        "$ cd sage",
        "$ make venv              # one-time setup, creates .venv",
        "",
        "# Pick your free LLM (no API key needed):",
        "$ ollama serve && ollama pull llama3.2     # fully local",
        "# OR",
        "$ npm install -g @google/gemini-cli && gemini  # free Google login",
        "# OR",
        "$ npm install -g @anthropic-ai/claude-code && claude  # free login",
        "",
        "$ ./launch.sh starter    # starts backend (8000) + frontend (5173)",
        "",
        "# Open http://localhost:5173 — that's it!",
    ]
    top = Inches(1.95)
    for line in code_lines:
        color = SAGE_GREEN_D if line.startswith("$") else SAGE_GRAY
        if line.startswith("#"):
            color = SAGE_LGRAY
        add_text_box(slide, line,
                     Inches(1.3), top, Inches(10.5), Inches(0.28),
                     font_size=11, color=color)
        top += Inches(0.24)

    # Bottom section
    add_text_box(slide, "No API keys. No cloud accounts. Runs on your laptop. MIT licensed.",
                 Inches(1), Inches(5.6), Inches(11), Inches(0.4),
                 font_size=16, bold=True, color=SAGE_GREEN_D, align=PP_ALIGN.CENTER)

    # New solution
    add_rect(slide, Inches(1), Inches(6.1), Inches(5.3), Inches(0.7),
             fill_color=SAGE_GREEN_L, line_color=SAGE_GREEN, line_width=Pt(0.5))
    add_text_box(slide, "New solution from template:",
                 Inches(1.15), Inches(6.15), Inches(5), Inches(0.25),
                 font_size=10, bold=True, color=SAGE_GREEN_D)
    add_text_box(slide, "cp -r solutions/starter solutions/my_domain && make run PROJECT=my_domain",
                 Inches(1.15), Inches(6.4), Inches(5), Inches(0.25),
                 font_size=9, color=SAGE_DARK)

    # Onboarding wizard
    add_rect(slide, Inches(6.6), Inches(6.1), Inches(5.7), Inches(0.7),
             fill_color=SAGE_GREEN_L, line_color=SAGE_GREEN, line_width=Pt(0.5))
    add_text_box(slide, "LLM-powered onboarding wizard:",
                 Inches(6.75), Inches(6.15), Inches(5.4), Inches(0.25),
                 font_size=10, bold=True, color=SAGE_GREEN_D)
    add_text_box(slide, "POST /onboarding/generate → describe your domain → 3 YAML files auto-generated",
                 Inches(6.75), Inches(6.4), Inches(5.4), Inches(0.25),
                 font_size=9, color=SAGE_DARK)

    footer_band(slide)

    # ─────────────────────────────────────────────────────────────────────────
    # SLIDE 16 — Open Source Model
    # ─────────────────────────────────────────────────────────────────────────
    slide = prs.slides.add_slide(blank_layout(prs))
    slide_bg(slide, WHITE)
    header_band(slide, "Open Source Model", "Framework is open. Solutions are private. Community contributions welcome.")

    # Three-column layout
    col_w = Inches(3.8)
    col_gap = Inches(0.35)
    cols_start = (SLIDE_W - 3*col_w - 2*col_gap) / 2

    oss_sections = [
        ("SAGE Framework", SAGE_GREEN, [
            "Open source (MIT)",
            "136 API endpoints",
            "30 UI pages",
            "16+ solution templates",
            "github.com/Sumanharapanahalli/sage",
        ]),
        ("Your Solutions", SAGE_GRAY, [
            "Private — your own repo",
            "Mount via SAGE_SOLUTIONS_DIR",
            ".sage/ runtime data — gitignored",
            "Your IP stays private",
            "Compliance data per-solution",
        ]),
        ("Community", SAGE_GREEN_D, [
            "Star, fork, submit PRs",
            "Build MCP tool servers",
            "Create solution templates",
            "File issues for bugs / features",
            "See CONTRIBUTING.md",
        ]),
    ]

    for i, (title, accent, items) in enumerate(oss_sections):
        x = cols_start + i * (col_w + col_gap)
        add_rect(slide, x, Inches(1.8), col_w, Inches(3.2),
                 fill_color=SAGE_GREEN_L, line_color=accent, line_width=Pt(1.5))
        add_text_box(slide, title,
                     x + Inches(0.15), Inches(1.95), col_w - Inches(0.3), Inches(0.4),
                     font_size=16, bold=True, color=accent, align=PP_ALIGN.CENTER)
        top = Inches(2.5)
        for item in items:
            add_text_box(slide, f"· {item}",
                         x + Inches(0.25), top, col_w - Inches(0.5), Inches(0.3),
                         font_size=12, color=SAGE_DARK)
            top += Inches(0.32)

    # Advantage for regulators
    add_rect(slide, Inches(0.5), Inches(5.3), Inches(12.3), Inches(1.0),
             fill_color=SAGE_GREEN_L, line_color=SAGE_GREEN, line_width=Pt(1))
    add_text_box(slide, "Why Open Source Matters for Regulated Industries",
                 Inches(0.7), Inches(5.4), Inches(12), Inches(0.35),
                 font_size=14, bold=True, color=SAGE_GREEN_D)
    add_text_box(slide, "Auditor transparency — regulators can inspect the entire AI agent toolchain.  No vendor lock-in — fork it, own it.  Community security review — continuous scrutiny reduces vulnerabilities in compliance-critical paths.",
                 Inches(0.7), Inches(5.8), Inches(12), Inches(0.4),
                 font_size=11, color=SAGE_DARK)

    footer_band(slide)

    # ─────────────────────────────────────────────────────────────────────────
    # SLIDE 17 — Closing
    # ─────────────────────────────────────────────────────────────────────────
    slide = prs.slides.add_slide(blank_layout(prs))
    slide_bg(slide, WHITE)
    add_rect(slide, 0, 0, SLIDE_W, Inches(0.06), fill_color=SAGE_GREEN)
    add_rect(slide, 0, Inches(7.44), SLIDE_W, Inches(0.06), fill_color=SAGE_GREEN)

    add_text_box(slide, "SAGE[ai]",
                 Inches(1), Inches(1.5), Inches(11), Inches(1.0),
                 font_size=60, bold=True, color=SAGE_GREEN, align=PP_ALIGN.CENTER)
    add_text_box(slide, "Lean Development × Agentic AI × Human Oversight",
                 Inches(1), Inches(2.7), Inches(11), Inches(0.5),
                 font_size=22, color=SAGE_DARK, align=PP_ALIGN.CENTER)

    add_text_box(slide, "136 endpoints  ·  30 pages  ·  840+ tests  ·  20 agent roles  ·  14 domains  ·  7 LLM providers",
                 Inches(1), Inches(3.8), Inches(11), Inches(0.4),
                 font_size=14, color=SAGE_GRAY, align=PP_ALIGN.CENTER)

    add_text_box(slide, "Open Source (MIT)  ·  Self-Hosted  ·  Air-Gappable  ·  Zero API Keys Required",
                 Inches(1), Inches(4.3), Inches(11), Inches(0.4),
                 font_size=14, color=SAGE_GRAY, align=PP_ALIGN.CENTER)

    add_text_box(slide, "github.com/Sumanharapanahalli/sage",
                 Inches(1), Inches(5.2), Inches(11), Inches(0.5),
                 font_size=20, bold=True, color=SAGE_GREEN, align=PP_ALIGN.CENTER)

    add_text_box(slide, "Star it. Fork it. Build on it.",
                 Inches(1), Inches(5.8), Inches(11), Inches(0.4),
                 font_size=16, italic=True, color=SAGE_GREEN_D, align=PP_ALIGN.CENTER)

    # ── Save ─────────────────────────────────────────────────────────────────
    out_path = os.path.join(os.path.dirname(__file__), "SageAI_Tech_Pitch.pptx")
    prs.save(out_path)
    print(f"✓ Saved: {out_path}")
    return out_path


def table_slide_helper(slide, left, top, headers, rows, col_widths):
    """Inline table on an existing slide."""
    n_cols = len(headers)
    n_rows = len(rows)
    row_h = Inches(0.42)
    tbl_h = row_h * (n_rows + 1)

    table = slide.shapes.add_table(
        n_rows + 1, n_cols, left, top,
        sum(col_widths), tbl_h
    ).table

    for ci, w in enumerate(col_widths):
        table.columns[ci].width = w

    for ci, h in enumerate(headers):
        cell = table.cell(0, ci)
        cell.fill.solid()
        cell.fill.fore_color.rgb = SAGE_GREEN
        p = cell.text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = h
        run.font.bold = True
        run.font.size = Pt(11)
        run.font.color.rgb = WHITE

    for ri, row in enumerate(rows):
        bg = WHITE if ri % 2 == 0 else SAGE_GREEN_L
        for ci, val in enumerate(row):
            cell = table.cell(ri + 1, ci)
            cell.fill.solid()
            cell.fill.fore_color.rgb = bg
            p = cell.text_frame.paragraphs[0]
            p.alignment = PP_ALIGN.LEFT
            run = p.add_run()
            run.text = str(val)
            run.font.size = Pt(10)
            run.font.color.rgb = SAGE_DARK


if __name__ == "__main__":
    build_tech_deck()
