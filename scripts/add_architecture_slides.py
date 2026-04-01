"""
SAGE Tech Pitch — Architecture Deep-Dive Slides
==================================================
Adds 9 detailed technical architecture slides to SageAI_Tech_Pitch.pptx.
Inserts them after slide 4 (Architecture Overview) as a contiguous
"Architecture Deep Dive" section.

Run:  python docs/add_architecture_slides.py

New slides:
  5.  Request Lifecycle (HTTP → Agent → LLM → Audit → Response)
  6.  SAGE Lean Loop — Phase-by-Phase Data Flow
  7.  Approval Pipeline Architecture (Risk → Store → HITL → Executor)
  8.  Memory Architecture (Vector Store + Audit Log Compounding)
  9.  Build Orchestrator Internal Workflow
 10.  LLM Gateway Internals (Singleton, Thread Lock, 7 Providers)
 11.  Solution Isolation Model (.sage/ per solution)
 12.  MCP Tool Architecture (Framework + Solution Layers)
 13.  Integration Architecture — Graceful Degradation
 14.  Sandboxed Execution — 3-Tier Isolation Cascade
 15.  Agentic Patterns — 0→1 Greenfield Build
 16.  Agentic Patterns — 1→N Incremental Refinement
"""

import copy
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

# ── Colour palette (matches existing deck) ────────────────────────────────────
SAGE_DARK  = RGBColor(0x1F, 0x2A, 0x37)
SAGE_GREEN = RGBColor(0x10, 0xB9, 0x81)
SAGE_TEAL  = RGBColor(0x04, 0x78, 0x57)
SAGE_AMBER = RGBColor(0xF5, 0x9E, 0x0B)
SAGE_RED   = RGBColor(0xEF, 0x44, 0x44)
SAGE_BLUE  = RGBColor(0x3B, 0x82, 0xF6)
SAGE_LIGHT = RGBColor(0xFA, 0xFA, 0xFA)
SAGE_MINT  = RGBColor(0xD1, 0xFA, 0xE5)
WARN_BG    = RGBColor(0xFF, 0xF3, 0xCC)
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
GREY       = RGBColor(0x4B, 0x55, 0x63)
LIGHT_GREY = RGBColor(0x9C, 0xA3, 0xAF)

SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)


# ═══════════════════════════════════════════════════════════════════════════════
# Helpers (same pattern as generate_presentations.py)
# ═══════════════════════════════════════════════════════════════════════════════

def blank_layout(prs):
    return prs.slide_layouts[6]


def add_rect(slide, l, t, w, h, fill_color=None, line_color=None):
    shape = slide.shapes.add_shape(1, l, t, w, h)
    shape.line.width = Pt(0)
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    else:
        shape.fill.background()
    if line_color:
        shape.line.color.rgb = line_color
    else:
        shape.line.fill.background()
    return shape


def add_text(slide, text, l, t, w, h, sz=14, bold=False, color=SAGE_DARK,
             align=PP_ALIGN.LEFT, italic=False):
    txBox = slide.shapes.add_textbox(l, t, w, h)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(sz)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return txBox


def add_multiline(slide, lines, l, t, w, line_h=Inches(0.34), sz=12,
                  color=SAGE_DARK, bullet="", bold_first=False):
    """Add multiple lines of text as separate text boxes."""
    for i, line in enumerate(lines):
        txt = f"{bullet}{line}" if bullet else line
        b = bold_first and i == 0
        add_text(slide, txt, l, t + line_h * i, w, line_h, sz=sz, bold=b, color=color)


def slide_bg(slide):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = SAGE_LIGHT


def header_band(slide, title, subtitle=None):
    add_rect(slide, 0, 0, SLIDE_W, Inches(1.45), fill_color=SAGE_DARK)
    add_text(slide, title, Inches(0.45), Inches(0.12), Inches(12.4), Inches(0.8),
             sz=30, bold=True, color=WHITE)
    if subtitle:
        add_text(slide, subtitle, Inches(0.45), Inches(0.88), Inches(12.4), Inches(0.48),
                 sz=16, color=SAGE_TEAL)


def footer_band(slide):
    add_rect(slide, 0, Inches(7.1), SLIDE_W, Inches(0.4), fill_color=SAGE_DARK)
    add_text(slide, "SAGE[ai] — Architecture Deep Dive | github.com/Sumanharapanahalli/sage",
             Inches(0.3), Inches(7.12), Inches(9), Inches(0.3), sz=9, color=GREY)


def diagram_box(slide, text, l, t, w, h, fill=SAGE_DARK, text_color=WHITE, sz=11):
    add_rect(slide, l, t, w, h, fill_color=fill)
    add_text(slide, text, l + Inches(0.08), t + Inches(0.04), w - Inches(0.16),
             h - Inches(0.08), sz=sz, bold=True, color=text_color, align=PP_ALIGN.CENTER)


def arrow_right(slide, l, t, w=Inches(0.5)):
    add_text(slide, "\u2192", l, t, w, Inches(0.3), sz=20, bold=True,
             color=SAGE_GREEN, align=PP_ALIGN.CENTER)


def arrow_down(slide, l, t):
    add_text(slide, "\u2193", l, t, Inches(0.4), Inches(0.3), sz=18, bold=True,
             color=SAGE_GREEN, align=PP_ALIGN.CENTER)


def section_label(slide, text, l, t, w=Inches(2)):
    add_rect(slide, l, t, w, Inches(0.28), fill_color=SAGE_GREEN)
    add_text(slide, text, l + Inches(0.06), t + Inches(0.02), w - Inches(0.12),
             Inches(0.24), sz=10, bold=True, color=WHITE, align=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════════════════════════════
# Slide builders
# ═══════════════════════════════════════════════════════════════════════════════

def slide_request_lifecycle(prs):
    """Slide 5: Request Lifecycle — HTTP to Audit Log."""
    slide = prs.slides.add_slide(blank_layout(prs))
    slide_bg(slide)
    header_band(slide, "Request Lifecycle — End-to-End Trace",
                "Every request traceable from HTTP ingress to audit log via trace_id")

    # 8-stage horizontal pipeline
    stages = [
        ("HTTP\nRequest", "FastAPI\nRate-limited"),
        ("Tenant\nResolve", "X-SAGE-Tenant\nContextVar"),
        ("Route\nHandler", "Pydantic\ntrace_id gen"),
        ("Agent\nDispatch", "Analyst/Dev\n/Planner"),
        ("LLM\nGateway", "Lock + PII\n+ Budget"),
        ("Proposal\nCreate", "Risk class\nSQLite"),
        ("Human\nDecision", "Approve or\nReject + fb"),
        ("Audit +\nCompound", ".sage/db\n+ vector"),
    ]

    box_w = Inches(1.35)
    box_h = Inches(0.65)
    desc_h = Inches(0.45)
    gap = Inches(0.18)
    start_l = Inches(0.3)
    top = Inches(2.0)

    for i, (label, desc) in enumerate(stages):
        l = start_l + i * (box_w + gap)
        # Box
        fill = SAGE_DARK if i not in (5, 6) else (SAGE_AMBER if i == 5 else SAGE_RED)
        diagram_box(slide, label, l, top, box_w, box_h, fill=fill, sz=10)
        # Description below
        add_text(slide, desc, l, top + box_h + Inches(0.05), box_w, desc_h,
                 sz=9, color=GREY, align=PP_ALIGN.CENTER)
        # Arrow between
        if i < len(stages) - 1:
            arrow_right(slide, l + box_w - Inches(0.05), top + Inches(0.18), Inches(0.4))

    # Key implementation details
    top2 = Inches(4.0)
    section_label(slide, "Implementation Details", Inches(0.3), top2, Inches(2.2))
    details = [
        "threading.Lock in LLM Gateway — single-lane inference, thread-safe singleton",
        "Rate limiting: 120 writes/min, 300 reads/min per IP (sliding window middleware)",
        "PII scrub before every LLM call — configurable per solution",
        "Per-agent budget ceilings — monthly call limits defined in project.yaml",
        "Langfuse tracing (opt-in) — every generate() call creates a span",
        "_resolve_db_path() routes audit to per-solution .sage/audit_log.db",
    ]
    add_multiline(slide, details, Inches(0.4), top2 + Inches(0.4),
                  Inches(12.5), sz=11, color=SAGE_DARK, bullet="\u2022  ")

    footer_band(slide)
    return slide


def slide_lean_loop_detail(prs):
    """Slide 6: SAGE Lean Loop — Detailed data flow per phase."""
    slide = prs.slides.add_slide(blank_layout(prs))
    slide_bg(slide)
    header_band(slide, "The SAGE Lean Loop — Phase-by-Phase Data Flow",
                "Every human correction compounds into future agent context")

    # 5-phase circular-ish layout
    phases = [
        ("1. SURFACE", "api.py / webhook / n8n / Slack",
         "HTTP payload, log entry, event", "Structured task + trace_id"),
        ("2. CONTEXTUALIZE", "vector_store.py .search()",
         "Task text + trace_id", "Top-k past decisions + corrections"),
        ("3. PROPOSE", "Agent + llm_gateway.generate()",
         "Task + context + prompts.yaml", "Structured proposal (JSON)"),
        ("4. DECIDE", "ProposalStore + Human review",
         "Proposal with risk class", "Approved/Rejected + feedback"),
        ("5. COMPOUND", "audit_logger + vector_store.add()",
         "Decision + feedback text", "Audit row + vector embedding"),
    ]

    col_w = [Inches(1.8), Inches(2.8), Inches(2.8), Inches(2.8)]
    headers = ["Phase", "Component", "Data In", "Data Out"]
    top = Inches(1.7)
    left = Inches(0.5)
    row_h = Inches(0.55)

    # Header row
    x = left
    for i, (hdr, w) in enumerate(zip(headers, col_w)):
        add_rect(slide, x, top, w, row_h, fill_color=SAGE_GREEN)
        add_text(slide, hdr, x + Inches(0.08), top + Inches(0.1), w - Inches(0.16),
                 Inches(0.35), sz=12, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        x += w

    # Data rows
    for ri, (phase, comp, din, dout) in enumerate(phases):
        y = top + row_h * (ri + 1)
        bg = WHITE if ri % 2 == 0 else SAGE_LIGHT
        vals = [phase, comp, din, dout]
        x = left
        for ci, (val, w) in enumerate(zip(vals, col_w)):
            add_rect(slide, x, y, w, row_h, fill_color=bg)
            c = SAGE_DARK if ci > 0 else SAGE_TEAL
            b = ci == 0
            add_text(slide, val, x + Inches(0.08), y + Inches(0.08), w - Inches(0.16),
                     row_h - Inches(0.16), sz=10, bold=b, color=c)
            x += w

    # Feedback loop annotation
    top3 = Inches(4.8)
    add_rect(slide, Inches(0.4), top3, Inches(12.5), Inches(1.8), fill_color=SAGE_DARK)
    section_label(slide, "COMPOUNDING LOOP", Inches(0.6), top3 + Inches(0.1), Inches(2))
    loop_text = (
        "Phase 5 feeds back into Phase 2 — this IS the Memento principle.\n\n"
        "Human rejects proposal with feedback  \u2192  feedback stored in vector store as document  "
        "\u2192  next similar query retrieves this correction  \u2192  agent generates better proposal  "
        "\u2192  quality improves monotonically without model retraining.\n\n"
        "Every rejection is stored with the human's actual reasoning, not just 'rejected'. "
        "This is the training signal AND the compliance record."
    )
    add_text(slide, loop_text, Inches(0.6), top3 + Inches(0.45), Inches(12.1), Inches(1.3),
             sz=11, color=SAGE_MINT)

    footer_band(slide)
    return slide


def slide_approval_pipeline(prs):
    """Slide 7: Approval Pipeline Architecture."""
    slide = prs.slides.add_slide(blank_layout(prs))
    slide_bg(slide)
    header_band(slide, "Approval Pipeline — Risk Classification to Execution",
                "ProposalStore (SQLite) \u2192 Risk Tiers \u2192 HITL \u2192 ProposalExecutor")

    # Left column: Proposal Model
    section_label(slide, "Proposal Model", Inches(0.4), Inches(1.65))
    model_fields = [
        "trace_id: str (UUID4)",
        "action_type: yaml_edit | code_diff | knowledge_delete | agent_hire",
        "risk_class: RiskClass (5-level enum)",
        "proposed_by: 'AnalystAgent' | 'user:admin' | 'OnboardingWizard'",
        "status: pending | approved | rejected | expired",
        "approved_by / approver_role / approver_email (RBAC)",
        "required_role: 'admin' | 'operator' | None",
    ]
    add_multiline(slide, model_fields, Inches(0.5), Inches(2.0),
                  Inches(5.5), sz=10, color=SAGE_DARK, bullet="\u2022  ", line_h=Inches(0.28))

    # Center column: Risk tiers
    section_label(slide, "Risk Tiers & Expiry", Inches(6.2), Inches(1.65), Inches(2.5))
    tiers = [
        ("INFORMATIONAL", "1 hour", "read-only query"),
        ("EPHEMERAL", "8 hours", "LLM switch, config"),
        ("STATEFUL", "7 days", "knowledge edit, YAML"),
        ("EXTERNAL", "14 days", "GitLab MR, deploy"),
        ("DESTRUCTIVE", "Never", "drop data, delete agent"),
    ]
    y = Inches(2.0)
    for tier, expiry, desc in tiers:
        colors = {
            "INFORMATIONAL": GREY,
            "EPHEMERAL": SAGE_AMBER,
            "STATEFUL": SAGE_GREEN,
            "EXTERNAL": SAGE_BLUE,
            "DESTRUCTIVE": SAGE_RED,
        }
        add_rect(slide, Inches(6.3), y, Inches(1.8), Inches(0.3), fill_color=colors[tier])
        add_text(slide, tier, Inches(6.35), y + Inches(0.02), Inches(1.7), Inches(0.26),
                 sz=9, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        add_text(slide, f"{expiry} — {desc}", Inches(8.2), y + Inches(0.02),
                 Inches(4.5), Inches(0.26), sz=10, color=SAGE_DARK)
        y += Inches(0.38)

    # Bottom: Executor dispatch
    section_label(slide, "ProposalExecutor Dispatch", Inches(0.4), Inches(4.5), Inches(2.8))
    executors = [
        "_execute_yaml_edit()  — validate YAML, write to solution dir, reload config",
        "_execute_code_diff()  — apply diff, run tests, create git commit",
        "_execute_config_switch()  — reload project_config singleton",
        "_execute_llm_switch()  — instantiate new provider, reset usage counters",
        "_execute_knowledge_delete()  — remove from vector store, audit log entry",
    ]
    add_multiline(slide, executors, Inches(0.5), Inches(4.9),
                  Inches(12), sz=10, color=SAGE_DARK, bullet="\u2022  ", line_h=Inches(0.28))

    # Two-tier note
    add_rect(slide, Inches(0.4), Inches(6.4), Inches(12.5), Inches(0.55), fill_color=WARN_BG)
    add_text(slide, ("Two tiers enforced: Framework control ops (config_switch, llm_switch) "
                     "execute immediately — operator action, not agent action. "
                     "Agent proposals (yaml_edit, code_diff, knowledge_delete, agent_hire) "
                     "always require human sign-off. This distinction IS the compliance guarantee."),
             Inches(0.5), Inches(6.42), Inches(12.3), Inches(0.5),
             sz=10, italic=True, color=SAGE_DARK)

    footer_band(slide)
    return slide


def slide_memory_architecture(prs):
    """Slide 8: Memory Architecture — Vector Store + Audit Log."""
    slide = prs.slides.add_slide(blank_layout(prs))
    slide_bg(slide)
    header_band(slide, "Memory Architecture — Compounding Intelligence",
                "ChromaDB semantic search + SQLite audit log. Per-solution isolation.")

    # Left: Vector Store
    section_label(slide, "Vector Store", Inches(0.4), Inches(1.65), Inches(1.8))
    vs_lines = [
        "Backend: ChromaDB + sentence-transformers (HuggingFaceEmbeddings)",
        "Minimal mode (SAGE_MINIMAL=1): keyword-match in-memory, 5MB RAM",
        "Collection per tenant: <tenant>_knowledge",
        "CRUD: GET/POST/DELETE /knowledge/...",
        "Bulk import supported (POST /knowledge/bulk)",
        "Semantic search with configurable top-k",
        "Fallback chain: HuggingFace \u2192 langchain_community \u2192 keyword",
    ]
    add_multiline(slide, vs_lines, Inches(0.5), Inches(2.0),
                  Inches(5.8), sz=10, color=SAGE_DARK, bullet="\u2022  ", line_h=Inches(0.28))

    # Right: Audit Logger
    section_label(slide, "Audit Logger", Inches(6.8), Inches(1.65), Inches(1.8))
    al_lines = [
        "SQLite per-solution: .sage/audit_log.db",
        "Table: compliance_audit_log",
        "Columns: id, timestamp, actor, action_type,",
        "  input_context, output_content, metadata (JSON),",
        "  verification_signature, approved_by,",
        "  approver_role, approver_email, approver_provider",
        "Append-only (no UPDATE/DELETE)",
        "FDA 21 CFR Part 11 compatible",
        "ISO 13485 clause 4.2.5 ready",
    ]
    add_multiline(slide, al_lines, Inches(6.9), Inches(2.0),
                  Inches(5.8), sz=10, color=SAGE_DARK, bullet="\u2022  ", line_h=Inches(0.28))

    # Bottom: Compounding loop diagram
    top3 = Inches(4.8)
    diagram_box(slide, "Human rejects\nwith feedback", Inches(0.5), top3,
                Inches(2), Inches(0.7), fill=SAGE_RED, sz=10)
    arrow_right(slide, Inches(2.5), top3 + Inches(0.2))
    diagram_box(slide, "Feedback stored\nin vector store", Inches(3.2), top3,
                Inches(2.2), Inches(0.7), fill=SAGE_TEAL, sz=10)
    arrow_right(slide, Inches(5.4), top3 + Inches(0.2))
    diagram_box(slide, "Next query\nretrieves correction", Inches(6.1), top3,
                Inches(2.2), Inches(0.7), fill=SAGE_DARK, sz=10)
    arrow_right(slide, Inches(8.3), top3 + Inches(0.2))
    diagram_box(slide, "Agent generates\nbetter proposal", Inches(9.0), top3,
                Inches(2.2), Inches(0.7), fill=SAGE_GREEN, sz=10)
    arrow_right(slide, Inches(11.2), top3 + Inches(0.2))
    diagram_box(slide, "Quality\nimproves", Inches(11.8), top3,
                Inches(1.2), Inches(0.7), fill=SAGE_DARK, sz=10)

    add_text(slide, "Memento principle: behavioral optimization via memory retrieval, "
             "not retraining. Inference costs stay flat. Improvement velocity increases.",
             Inches(0.5), Inches(5.7), Inches(12), Inches(0.4),
             sz=11, italic=True, color=GREY)

    footer_band(slide)
    return slide


def slide_build_orchestrator(prs):
    """Slide 9: Build Orchestrator Internal Workflow."""
    slide = prs.slides.add_slide(blank_layout(prs))
    slide_bg(slide)
    header_band(slide, "Build Orchestrator — Internal State Machine",
                "Idea \u2192 Domain Detection \u2192 Task Decomposition \u2192 Critic \u2192 HITL \u2192 Parallel Build")

    # State machine (vertical)
    states = [
        ("decomposing", "LLM decomposes into tasks", SAGE_DARK),
        ("critic_plan", "CriticAgent scores 0-100", SAGE_TEAL),
        ("awaiting_plan_approval", "HITL gate (if standard/strict)", SAGE_AMBER),
        ("building", "Parallel wave execution", SAGE_GREEN),
        ("critic_code", "Code quality review", SAGE_TEAL),
        ("integrating", "Merge all outputs", SAGE_DARK),
        ("critic_integration", "Final system review", SAGE_TEAL),
        ("awaiting_final_approval", "HITL gate", SAGE_RED),
        ("completed", "Done", SAGE_GREEN),
    ]

    y = Inches(1.7)
    for i, (state, desc, color) in enumerate(states):
        diagram_box(slide, state, Inches(0.5), y, Inches(2.5), Inches(0.38), fill=color, sz=9)
        add_text(slide, desc, Inches(3.2), y + Inches(0.04), Inches(3), Inches(0.3),
                 sz=10, color=SAGE_DARK)
        if i < len(states) - 1:
            arrow_down(slide, Inches(1.5), y + Inches(0.38))
        y += Inches(0.58)

    # Right panel: Key subsystems
    section_label(slide, "AdaptiveRouter (Q-learning)", Inches(7), Inches(1.7), Inches(2.8))
    router_lines = [
        "scores[task_type][agent_role] — EMA updates",
        "LEARNING_RATE = 0.3",
        "MIN_OBSERVATIONS = 3 before override",
        "Exploration: picks random agent 20% of time",
        "GET /build/router/stats for live weights",
    ]
    add_multiline(slide, router_lines, Inches(7.1), Inches(2.05),
                  Inches(5.5), sz=10, color=SAGE_DARK, bullet="\u2022  ", line_h=Inches(0.26))

    section_label(slide, "CriticAgent (Actor-Critic)", Inches(7), Inches(3.5), Inches(2.8))
    critic_lines = [
        "review_plan(): {score, flaws, suggestions, missing}",
        "review_code(): {score, bugs, security_risks}",
        "review_integration(): {score, conflicts, gaps}",
        "Calibrated: MVP 60-70, prod 80+, regulated 85+",
    ]
    add_multiline(slide, critic_lines, Inches(7.1), Inches(3.85),
                  Inches(5.5), sz=10, color=SAGE_DARK, bullet="\u2022  ", line_h=Inches(0.26))

    section_label(slide, "Crash Recovery", Inches(7), Inches(5.0), Inches(2.2))
    recovery_lines = [
        "SQLite checkpoints: build_checkpoints.db",
        "_checkpoint() after each state transition",
        "_restore_runs() on startup",
        "Anti-drift: BUILD_DRIFT_WARNING events",
    ]
    add_multiline(slide, recovery_lines, Inches(7.1), Inches(5.35),
                  Inches(5.5), sz=10, color=SAGE_DARK, bullet="\u2022  ", line_h=Inches(0.26))

    footer_band(slide)
    return slide


def slide_llm_gateway(prs):
    """Slide 10: LLM Gateway Internals."""
    slide = prs.slides.add_slide(blank_layout(prs))
    slide_bg(slide)
    header_band(slide, "LLM Gateway — Thread-Safe Multi-Provider Singleton",
                "7 providers, provider pool, PII scrub, per-agent budgets")

    # Provider hierarchy
    section_label(slide, "Provider Classes (ABC)", Inches(0.4), Inches(1.65), Inches(2.5))
    providers = [
        ("GeminiCLIProvider", "subprocess, -p flag, no API key"),
        ("ClaudeCodeCLIProvider", "subprocess, claude CLI, no key"),
        ("OllamaProvider", "REST to localhost:11434, offline"),
        ("LocalLlamaProvider", "llama-cpp-python, GPU-direct, GGUF"),
        ("ClaudeAPIProvider", "anthropic SDK, only paid option"),
        ("GenericCLIProvider", "{prompt} template substitution"),
    ]
    y = Inches(2.0)
    for name, desc in providers:
        add_text(slide, name, Inches(0.5), y, Inches(2.5), Inches(0.24),
                 sz=10, bold=True, color=SAGE_TEAL)
        add_text(slide, desc, Inches(3.1), y, Inches(3), Inches(0.24),
                 sz=10, color=GREY)
        y += Inches(0.28)

    # generate() pipeline
    section_label(slide, "generate() Pipeline", Inches(0.4), Inches(3.9), Inches(2.2))
    pipeline = [
        "1. PII scrub (pii_filter.scrub_text)",
        "2. Data residency check",
        "3. Tenant + solution budget check",
        "4. Per-agent call ceiling (project.yaml)",
        "5. threading.Lock acquire",
        "6. provider.generate(prompt, system_prompt)",
        "7. Langfuse span (if enabled)",
        "8. Usage counter update + lock release",
    ]
    add_multiline(slide, pipeline, Inches(0.5), Inches(4.25),
                  Inches(5.5), sz=10, color=SAGE_DARK, line_h=Inches(0.26))

    # Right: Provider Pool
    section_label(slide, "ProviderPool — Parallel Strategies", Inches(7), Inches(1.65), Inches(3.2))
    strategies = [
        ("voting", "Send to N providers, majority consensus wins"),
        ("fastest", "First response wins (concurrent.futures)"),
        ("fallback", "Sequential — first success returned"),
        ("quality", "Longest/richest response selected"),
    ]
    y = Inches(2.05)
    for name, desc in strategies:
        diagram_box(slide, name, Inches(7.1), y, Inches(1.2), Inches(0.35),
                    fill=SAGE_TEAL, sz=9)
        add_text(slide, desc, Inches(8.5), y + Inches(0.04), Inches(4.3), Inches(0.3),
                 sz=10, color=SAGE_DARK)
        y += Inches(0.45)

    # Model limits
    section_label(slide, "Model Limits & Budgets", Inches(7), Inches(4.0), Inches(2.5))
    limits = [
        "Gemini free: 500/day (flash), 25/day (pro)",
        "Ollama / Local: unlimited (self-hosted)",
        "Claude API: per-key billing",
        "Per-agent ceilings in project.yaml",
        "Tenant-level budget isolation",
        "Teacher-student distillation: heavy \u2192 student",
    ]
    add_multiline(slide, limits, Inches(7.1), Inches(4.35),
                  Inches(5.5), sz=10, color=SAGE_DARK, bullet="\u2022  ", line_h=Inches(0.26))

    footer_band(slide)
    return slide


def slide_solution_isolation(prs):
    """Slide 11: Solution Isolation Model."""
    slide = prs.slides.add_slide(blank_layout(prs))
    slide_bg(slide)
    header_band(slide, "Solution Isolation — .sage/ Per Solution + Multi-Tenant",
                "Zero data overlap. Runtime state lives with the solution, not the framework.")

    # Left: Directory structure
    section_label(slide, "Filesystem Layout", Inches(0.4), Inches(1.65), Inches(2))
    fs_lines = [
        "your-private-repo/",
        "  board_games/",
        "    project.yaml         \u2190 committed",
        "    prompts.yaml         \u2190 committed",
        "    tasks.yaml           \u2190 committed",
        "    .sage/               \u2190 auto-created, gitignored",
        "      audit_log.db       \u2190 proposals, approvals, audit",
        "      chroma_db/         \u2190 vector knowledge store",
        "  medtech_team/",
        "    project.yaml  ...",
        "    .sage/               \u2190 SEPARATE, ISOLATED",
        "      audit_log.db       \u2190 different DB, different data",
    ]
    add_multiline(slide, fs_lines, Inches(0.5), Inches(2.0),
                  Inches(6), sz=10, color=SAGE_DARK, line_h=Inches(0.28))

    # Right: Tenant context
    section_label(slide, "Multi-Tenant Context", Inches(7), Inches(1.65), Inches(2.5))
    tenant_lines = [
        "X-SAGE-Tenant header (optional)",
        "\u2192  tenant.py ContextVar (per-request scoped)",
        "\u2192  No header? Falls back to active solution name",
        "",
        "Scopes isolated per tenant:",
        "  \u2022  Vector collection: <tenant>_knowledge",
        "  \u2022  Audit log: tenant_id field in metadata",
        "  \u2022  Task queue: tagged with tenant",
        "",
        "SAGE_SOLUTIONS_DIR env var:",
        "  \u2022  Mount private solutions repo",
        "  \u2022  Framework repo has zero user data",
    ]
    add_multiline(slide, tenant_lines, Inches(7.1), Inches(2.0),
                  Inches(5.5), sz=10, color=SAGE_DARK, line_h=Inches(0.28))

    # Bottom: Key guarantees
    add_rect(slide, Inches(0.4), Inches(5.8), Inches(12.5), Inches(1.0), fill_color=SAGE_DARK)
    guarantees = [
        "Moving or archiving a solution takes its entire history with it",
        "The SAGE framework repo contains no user data, ever",
        "Regulated industries: .sage/audit_log.db IS the per-solution compliance record",
        "Two solutions on the same SAGE instance have zero data overlap",
    ]
    add_multiline(slide, guarantees, Inches(0.6), Inches(5.9),
                  Inches(12), sz=10, color=SAGE_MINT, bullet="\u2713  ", line_h=Inches(0.22))

    footer_band(slide)
    return slide


def slide_mcp_architecture(prs):
    """Slide 12: MCP Tool Architecture."""
    slide = prs.slides.add_slide(blank_layout(prs))
    slide_bg(slide)
    header_band(slide, "MCP Tool Architecture — Framework + Solution Layers",
                "Provider-agnostic: agents call invoke() as Python functions, works on any LLM")

    # Two-layer diagram
    section_label(slide, "Framework Tools (always loaded)", Inches(0.4), Inches(1.65), Inches(3.2))
    fw_tools = [
        ("filesystem_tools.py", "read/write/list/search — sandboxed to solution dir"),
        ("browser_tools.py", "Playwright — scrape, screenshot, interact"),
        ("sqlite_tools.py", "query audit_log.db — read-only enforced"),
    ]
    y = Inches(2.0)
    for name, desc in fw_tools:
        add_text(slide, name, Inches(0.5), y, Inches(2.5), Inches(0.24),
                 sz=10, bold=True, color=SAGE_TEAL)
        add_text(slide, desc, Inches(3.1), y, Inches(4), Inches(0.24), sz=10, color=SAGE_DARK)
        y += Inches(0.28)

    section_label(slide, "Solution Tools (per active solution)", Inches(0.4), Inches(3.0), Inches(3.2))
    sol_tools = [
        ("hardware_tools.py", "openocd, kicad, jlink, serial, CAN bus"),
        ("domain_tools.py", "Custom per-domain tools (flash_firmware, drc_check, etc.)"),
        ("Override", "Same-name tool in solution replaces framework version"),
    ]
    y = Inches(3.35)
    for name, desc in sol_tools:
        add_text(slide, name, Inches(0.5), y, Inches(2.5), Inches(0.24),
                 sz=10, bold=True, color=SAGE_TEAL)
        add_text(slide, desc, Inches(3.1), y, Inches(4), Inches(0.24), sz=10, color=SAGE_DARK)
        y += Inches(0.28)

    # Right: MCPRegistry
    section_label(slide, "MCPRegistry Class", Inches(7), Inches(1.65), Inches(2.5))
    registry_lines = [
        "_servers: dict[str, module]",
        "_tool_map: dict[str, callable]",
        "_loaded_solution: str",
        "",
        "load(force=False):",
        "  1. Discover src/mcp_servers/*.py",
        "  2. Discover solutions/<name>/mcp_servers/*.py",
        "  3. Solution tools override framework",
        "",
        "list_tools() \u2192 [{name, description, server}]",
        "invoke(tool, args, trace_id) \u2192 {result}",
        "as_react_tools() \u2192 dict for agent loops",
        "",
        "Every invocation audit-logged with trace_id",
        "FastMCP v2 support (providers \u2192 FunctionTool.fn)",
    ]
    add_multiline(slide, registry_lines, Inches(7.1), Inches(2.0),
                  Inches(5.5), sz=10, color=SAGE_DARK, line_h=Inches(0.26))

    # Bottom flow
    top3 = Inches(5.3)
    diagram_box(slide, "Agent calls\ninvoke('tool', args)", Inches(0.5), top3,
                Inches(2.5), Inches(0.6), fill=SAGE_DARK, sz=9)
    arrow_right(slide, Inches(3.0), top3 + Inches(0.15))
    diagram_box(slide, "MCPRegistry\nlookup + dispatch", Inches(3.7), top3,
                Inches(2.5), Inches(0.6), fill=SAGE_TEAL, sz=9)
    arrow_right(slide, Inches(6.2), top3 + Inches(0.15))
    diagram_box(slide, "Python function\nexecutes", Inches(6.9), top3,
                Inches(2), Inches(0.6), fill=SAGE_GREEN, sz=9)
    arrow_right(slide, Inches(8.9), top3 + Inches(0.15))
    diagram_box(slide, "Audit log\n+ result", Inches(9.5), top3,
                Inches(2), Inches(0.6), fill=SAGE_DARK, sz=9)

    add_text(slide, "Works identically on Gemini, Ollama, Claude, local GGUF — "
             "LLM never knows about MCP protocol",
             Inches(0.5), Inches(6.1), Inches(12), Inches(0.3),
             sz=11, italic=True, color=GREY)

    footer_band(slide)
    return slide


def slide_integration_architecture(prs):
    """Slide 13: Integration Architecture — Graceful Degradation."""
    slide = prs.slides.add_slide(blank_layout(prs))
    slide_bg(slide)
    header_band(slide, "Integration Architecture — Graceful Degradation Pattern",
                "Every integration is optional. SAGE runs fully functional with zero integrations.")

    # 4-quadrant layout
    q_w = Inches(6.0)
    q_h = Inches(2.4)
    gap = Inches(0.3)
    l1 = Inches(0.4)
    l2 = l1 + q_w + gap
    t1 = Inches(1.65)
    t2 = t1 + q_h + Inches(0.15)

    # Q1: Orchestration
    add_rect(slide, l1, t1, q_w, q_h, fill_color=WHITE, line_color=SAGE_GREEN)
    section_label(slide, "Orchestration Engines", l1 + Inches(0.1), t1 + Inches(0.1), Inches(2.5))
    q1_lines = [
        "LangGraph: StateGraph workflows, interrupt_before for HITL",
        "  \u2192  SQLite checkpoints (SqliteSaver), MemorySaver fallback",
        "  \u2192  solutions/<name>/workflows/*.py",
        "Temporal: durable workflows via TEMPORAL_HOST",
        "  \u2192  LangGraph fallback if unavailable",
        "AutoGen: Docker-sandboxed code execution",
    ]
    add_multiline(slide, q1_lines, l1 + Inches(0.15), t1 + Inches(0.45),
                  q_w - Inches(0.3), sz=10, color=SAGE_DARK, line_h=Inches(0.26))

    # Q2: External Communication
    add_rect(slide, l2, t1, q_w, q_h, fill_color=WHITE, line_color=SAGE_BLUE)
    section_label(slide, "External Communication", l2 + Inches(0.1), t1 + Inches(0.1), Inches(2.5))
    q2_lines = [
        "Slack: Block Kit approval cards, /webhook/slack callback",
        "  \u2192  Two-way: propose in Slack, approve in Slack",
        "n8n: Webhook receiver (N8N_WEBHOOK_SECRET)",
        "SSE Streaming: /analyze/stream, /agent/stream",
        "  \u2192  Real-time event streams to dashboard",
    ]
    add_multiline(slide, q2_lines, l2 + Inches(0.15), t1 + Inches(0.45),
                  q_w - Inches(0.3), sz=10, color=SAGE_DARK, line_h=Inches(0.26))

    # Q3: Coding Agents
    add_rect(slide, l1, t2, q_w, q_h, fill_color=WHITE, line_color=SAGE_AMBER)
    section_label(slide, "Coding Agents", l1 + Inches(0.1), t2 + Inches(0.1), Inches(2))
    q3_lines = [
        "OpenSWE: autonomous repo explore \u2192 implement \u2192 test \u2192 PR",
        "  \u2192  3-tier degradation: OpenSWE \u2192 LLM direct \u2192 template",
        "LangChain Tools: tool loader per solution",
        "  \u2192  langchain_tools.py discovers solution tool configs",
        "MCP Tools: framework + solution layers (see MCP slide)",
    ]
    add_multiline(slide, q3_lines, l1 + Inches(0.15), t2 + Inches(0.45),
                  q_w - Inches(0.3), sz=10, color=SAGE_DARK, line_h=Inches(0.26))

    # Q4: Observability
    add_rect(slide, l2, t2, q_w, q_h, fill_color=WHITE, line_color=SAGE_RED)
    section_label(slide, "Observability", l2 + Inches(0.1), t2 + Inches(0.1), Inches(2))
    q4_lines = [
        "Langfuse: trace per generate() call, opt-in",
        "  \u2192  observability.langfuse_enabled: true in config",
        "Built-in audit log: always active, no opt-in",
        "  \u2192  Per-solution .sage/audit_log.db",
        "Eval suite: YAML test cases, keyword scoring, history",
    ]
    add_multiline(slide, q4_lines, l2 + Inches(0.15), t2 + Inches(0.45),
                  q_w - Inches(0.3), sz=10, color=SAGE_DARK, line_h=Inches(0.26))

    # Bottom pattern
    add_text(slide, ("Pattern: every integration import is wrapped in try/except ImportError. "
                     "If langgraph is not installed, LangGraphRunner returns error dicts. "
                     "SAGE never crashes due to a missing optional dependency."),
             Inches(0.5), Inches(6.7), Inches(12), Inches(0.35),
             sz=10, italic=True, color=GREY)

    footer_band(slide)
    return slide


def slide_sandboxed_execution(prs):
    """Slide 14: Sandboxed Execution — 3-Tier Cascade."""
    slide = prs.slides.add_slide(blank_layout(prs))
    slide_bg(slide)
    header_band(slide, "Sandboxed Execution — 3-Tier Isolation Cascade",
                "OpenShell (container) → SandboxRunner (local clone) → OpenSWE (autonomous coding)")

    # 3-tier cascade diagram
    tiers = [
        ("Tier 1: OpenShell", "NVIDIA container sandbox",
         [
             "SSH-based execution inside container",
             "YAML security policies (commands, files, network)",
             "GPU-accelerated (ML training, inference)",
             "Full process + filesystem isolation",
             "Graceful degradation if unavailable",
         ], SAGE_DARK),
        ("Tier 2: SandboxRunner", "Local repo isolation",
         [
             "Clones solution repo → temp working dir",
             "Creates isolated branch per execution",
             "Restricts file ops to sandbox directory",
             "Primitives: execute(), read_file(), write_file()",
             "Configurable retain for debugging",
         ], SAGE_TEAL),
        ("Tier 3: OpenSWE", "Autonomous coding agent",
         [
             "Repo explore → implement → test → PR",
             "Internal 3-tier: SWE agent → LangGraph → LLM direct",
             "Unified output: {files_changed, tests_passed, pr_url}",
             "Integrates with build orchestrator wave execution",
         ], SAGE_GREEN),
    ]

    for i, (title, subtitle, bullets, color) in enumerate(tiers):
        l = Inches(0.4) + i * Inches(4.3)
        w = Inches(4.0)
        # Header box
        diagram_box(slide, title, l, Inches(1.7), w, Inches(0.45), fill=color, sz=11)
        add_text(slide, subtitle, l + Inches(0.1), Inches(2.2), w - Inches(0.2),
                 Inches(0.3), sz=10, italic=True, color=GREY, align=PP_ALIGN.CENTER)
        # Bullets
        add_multiline(slide, bullets, l + Inches(0.15), Inches(2.6),
                      w - Inches(0.3), sz=10, color=SAGE_DARK, bullet="• ", line_h=Inches(0.26))

        # Arrow between tiers
        if i < 2:
            arrow_l = l + w + Inches(0.02)
            add_text(slide, "fallback →", arrow_l, Inches(1.78), Inches(0.8), Inches(0.3),
                     sz=9, bold=True, color=SAGE_AMBER, align=PP_ALIGN.CENTER)

    # Bottom: cascade flow
    top3 = Inches(4.8)
    add_rect(slide, Inches(0.4), top3, Inches(12.5), Inches(2.0), fill_color=SAGE_DARK)
    section_label(slide, "Cascade Logic (build_orchestrator.py)", Inches(0.6), top3 + Inches(0.1), Inches(3.5))

    cascade_lines = [
        "1. Try OpenShell: if NVIDIA container available + task needs isolation → SSH exec in container",
        "2. Fallback to SandboxRunner: clone repo, isolate branch, execute locally",
        "3. Fallback to OpenSWE: autonomous coding agent (explore → implement → test → PR)",
        "4. Each tier returns same format: {success, output, files_changed} — orchestrator is tier-agnostic",
        "",
        "Key files: openshell_runner.py (Tier 1), sandbox_runner.py (Tier 2), openswe_runner.py (Tier 3)",
    ]
    add_multiline(slide, cascade_lines, Inches(0.6), top3 + Inches(0.45),
                  Inches(12), sz=10, color=SAGE_MINT, bullet="", line_h=Inches(0.26))

    footer_band(slide)
    return slide


def slide_agentic_patterns_0to1(prs):
    """Slide 15: Agentic Patterns — 0→1 Greenfield Build."""
    slide = prs.slides.add_slide(blank_layout(prs))
    slide_bg(slide)
    header_band(slide, "Agentic Patterns — 0→1 Greenfield Build",
                "Idea → Domain Detection → Decomposition → Critic → HITL → Wave Build → Ship")

    # Left: Pipeline stages
    section_label(slide, "0→1 Pipeline Stages", Inches(0.4), Inches(1.65), Inches(2.5))
    stages = [
        ("Description Input", "Plain-language product idea", SAGE_DARK),
        ("Domain Detection", "13 DOMAIN_RULES — auto-selects compliance, toolchains", SAGE_TEAL),
        ("Workforce Assembly", "19 agents, 5 teams from WORKFORCE_REGISTRY", SAGE_DARK),
        ("Hierarchical Decomposition", "LLM → 32 task types → dependency graph", SAGE_TEAL),
        ("Critic Reviews Plan", "CriticAgent scores 0-100, flaws, suggestions", SAGE_AMBER),
        ("HITL Plan Approval", "Human reviews decomposed plan", SAGE_RED),
        ("Wave Execution", "Independent tasks in parallel, deps sequential", SAGE_GREEN),
        ("Critic Reviews Code", "Per-task quality review", SAGE_AMBER),
        ("Integration Merge", "All outputs combined", SAGE_DARK),
        ("Critic Reviews Integration", "System-level review", SAGE_AMBER),
        ("HITL Final Approval", "Human signs off on completed build", SAGE_RED),
    ]

    y = Inches(2.0)
    for label, desc, color in stages:
        diagram_box(slide, label, Inches(0.5), y, Inches(2.8), Inches(0.32), fill=color, sz=8)
        add_text(slide, desc, Inches(3.5), y + Inches(0.02), Inches(3.5), Inches(0.28),
                 sz=9, color=SAGE_DARK)
        y += Inches(0.4)

    # Right: Patterns used
    section_label(slide, "Patterns in 0→1", Inches(7.2), Inches(1.65), Inches(2))
    patterns = [
        ("ReAct", "Per-task: observe → think → act → observe"),
        ("Hierarchical Decomp", "Description → task graph (32 types)"),
        ("Wave Parallel", "_compute_waves() groups independent work"),
        ("Adaptive Router", "Q-learning: scores[task][agent], ε=0.2"),
        ("Actor-Critic", "CriticAgent at 3 checkpoints (plan/code/integration)"),
        ("HITL Gates", "2 mandatory stops: post-plan, post-integration"),
        ("Anti-Drift", "Checkpoint after each state, BUILD_DRIFT_WARNING"),
        ("Iterative Refine", "Score < threshold → retry with critic feedback"),
    ]
    y = Inches(2.0)
    for name, desc in patterns:
        add_text(slide, name, Inches(7.3), y, Inches(2), Inches(0.24),
                 sz=9, bold=True, color=SAGE_TEAL)
        add_text(slide, desc, Inches(9.4), y, Inches(3.5), Inches(0.24),
                 sz=9, color=SAGE_DARK)
        y += Inches(0.3)

    # Bottom note
    add_rect(slide, Inches(0.4), Inches(6.4), Inches(12.5), Inches(0.55), fill_color=WARN_BG)
    add_text(slide, ("0→1 starts cold: AdaptiveRouter uses uniform scores, domain detected fresh, "
                     "full workforce assembled. Every decision during the build trains the router "
                     "and critic for future 1→N refinements."),
             Inches(0.5), Inches(6.42), Inches(12.3), Inches(0.5),
             sz=10, italic=True, color=SAGE_DARK)

    footer_band(slide)
    return slide


def slide_agentic_patterns_1toN(prs):
    """Slide 16: Agentic Patterns — 1→N Refinement."""
    slide = prs.slides.add_slide(blank_layout(prs))
    slide_bg(slide)
    header_band(slide, "Agentic Patterns — 1→N Incremental Refinement",
                "Same orchestrator, scoped to changes. Router warm. Critic calibrated. Intelligence compounds.")

    # Left: What's different in 1→N
    section_label(slide, "How 1→N Differs from 0→1", Inches(0.4), Inches(1.65), Inches(3))
    diffs = [
        ("Domain", "Already known — skips detection"),
        ("Workforce", "Stable — router has learned agent strengths"),
        ("Decomposition", "Scoped to the change, not the whole product"),
        ("Router Scores", "Warm — compounded from prior builds (EMA)"),
        ("Critic Calibration", "Higher baseline — knows existing quality"),
        ("Anti-Drift", "Compares against established baseline, not blank"),
    ]
    y = Inches(2.0)
    for label, desc in diffs:
        add_text(slide, label, Inches(0.5), y, Inches(1.8), Inches(0.24),
                 sz=10, bold=True, color=SAGE_TEAL)
        add_text(slide, desc, Inches(2.4), y, Inches(4.5), Inches(0.24),
                 sz=10, color=SAGE_DARK)
        y += Inches(0.3)

    # Right: Compounding intelligence diagram
    section_label(slide, "Compounding Intelligence Loop", Inches(7.2), Inches(1.65), Inches(3))
    loop_boxes = [
        ("0→1 Build", "Router starts uniform\nCritic calibrates", SAGE_DARK),
        ("Human Feedback", "Approvals + rejections\nstored with reasoning", SAGE_RED),
        ("Vector Store", "Prior decisions retrieved\nfor next context", SAGE_TEAL),
        ("1→N Refinement", "Router warm, critic tuned\nFaster + higher quality", SAGE_GREEN),
    ]
    y = Inches(2.0)
    for label, desc, color in loop_boxes:
        diagram_box(slide, label, Inches(7.3), y, Inches(2.2), Inches(0.55), fill=color, sz=9)
        add_text(slide, desc, Inches(9.7), y + Inches(0.04), Inches(3), Inches(0.48),
                 sz=9, color=SAGE_DARK)
        if label != "1→N Refinement":
            arrow_down(slide, Inches(8.3), y + Inches(0.55))
        y += Inches(0.75)

    # Feedback arrow back to top
    add_text(slide, "↻ compounds back", Inches(10.5), Inches(4.85), Inches(2), Inches(0.3),
             sz=9, bold=True, color=SAGE_GREEN)

    # Bottom: Triggers for 1→N
    top3 = Inches(4.3)
    section_label(slide, "1→N Triggers", Inches(0.4), top3, Inches(1.8))
    triggers = [
        "Feature request (from Improvements page, scope: solution)",
        "Bug report (from monitoring agent or human)",
        "Critic-identified improvement (score below threshold on re-eval)",
        "Human correction during approval (rejection feedback → improvement task)",
        "External event (Slack webhook, n8n trigger, GitHub issue)",
    ]
    add_multiline(slide, triggers, Inches(0.5), top3 + Inches(0.35),
                  Inches(6.5), sz=10, color=SAGE_DARK, bullet="• ", line_h=Inches(0.26))

    # Bottom summary
    add_rect(slide, Inches(0.4), Inches(6.2), Inches(12.5), Inches(0.7), fill_color=SAGE_DARK)
    add_text(slide, ("The same 8 patterns power both 0→1 and 1→N. The difference is state: "
                     "0→1 starts cold, 1→N starts warm. Every build makes the next one better — "
                     "this is the Memento principle applied to the entire product lifecycle."),
             Inches(0.6), Inches(6.3), Inches(12.1), Inches(0.55),
             sz=11, color=SAGE_MINT)

    footer_band(slide)
    return slide


# ═══════════════════════════════════════════════════════════════════════════════
# Main — Insert slides into existing deck
# ═══════════════════════════════════════════════════════════════════════════════

def main():
    script_dir = os.path.dirname(os.path.abspath(__file__))
    input_path = os.path.join(script_dir, "SageAI_Tech_Pitch.pptx")
    output_path = os.path.join(script_dir, "SageAI_Tech_Pitch.pptx")

    if not os.path.isfile(input_path):
        print(f"ERROR: {input_path} not found. Run generate_presentations.py first.")
        return

    # Build the new slides in a fresh presentation (same dimensions)
    new_prs = Presentation()
    new_prs.slide_width = SLIDE_W
    new_prs.slide_height = SLIDE_H

    builders = [
        slide_request_lifecycle,
        slide_lean_loop_detail,
        slide_approval_pipeline,
        slide_memory_architecture,
        slide_build_orchestrator,
        slide_llm_gateway,
        slide_solution_isolation,
        slide_mcp_architecture,
        slide_integration_architecture,
        slide_sandboxed_execution,
        slide_agentic_patterns_0to1,
        slide_agentic_patterns_1toN,
    ]

    for builder in builders:
        builder(new_prs)
        print(f"  Built: {builder.__doc__.strip().split(chr(10))[0]}")

    # Save as standalone deep-dive deck
    deepdive_path = os.path.join(script_dir, "SageAI_Architecture_DeepDive.pptx")
    new_prs.save(deepdive_path)
    print(f"\n  [OK] Architecture Deep Dive saved -> {deepdive_path}")
    print(f"       ({len(new_prs.slides)} slides)")
    print(f"\n  Share both decks:")
    print(f"    1. SageAI_Tech_Pitch.pptx           — overview (17 slides)")
    print(f"    2. SageAI_Architecture_DeepDive.pptx — deep dive (12 slides)")
    print(f"\n  Or merge them: copy deep-dive slides after slide 4 of the tech pitch.")


if __name__ == "__main__":
    main()
