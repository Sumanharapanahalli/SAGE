#!/usr/bin/env python3
"""
Update SAGE PowerPoint presentations with new features.
Run with: C:/sandbox/SAGE/.venv/Scripts/python C:/sandbox/SAGE/docs/update_presentations.py
"""

import sys
import io
import copy

# Force UTF-8 output
sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding='utf-8', errors='replace')

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from lxml import etree

TECH_PITCH_PATH = 'C:/sandbox/SAGE/docs/SageAI_Tech_Pitch.pptx'
BUSINESS_CASE_PATH = 'C:/sandbox/SAGE/docs/SageAI_Business_Case.pptx'
INVESTOR_PITCH_PATH = 'C:/sandbox/SAGE/docs/SageAI_Investor_Pitch.pptx'

changes_log = []

def log(msg):
    print(msg)
    changes_log.append(msg)

# ─────────────────────────────────────────────
# Helper: find shape by exact text match
# ─────────────────────────────────────────────
def find_shape_by_text(slide, target_text):
    for shape in slide.shapes:
        if shape.has_text_frame:
            if shape.text_frame.text.strip() == target_text.strip():
                return shape
    return None

def find_shape_containing(slide, substring):
    for shape in slide.shapes:
        if shape.has_text_frame:
            if substring in shape.text_frame.text:
                return shape
    return None

# ─────────────────────────────────────────────
# Helper: replace text in a shape preserving runs
# ─────────────────────────────────────────────
def replace_text_in_shape(shape, old_text, new_text):
    """Replace old_text with new_text across all paragraphs in a shape."""
    if not shape.has_text_frame:
        return False
    full_text = shape.text_frame.text
    if old_text not in full_text:
        return False

    # For single-paragraph shapes with one run, replace directly
    tf = shape.text_frame
    for para in tf.paragraphs:
        for run in para.runs:
            if old_text in run.text:
                run.text = run.text.replace(old_text, new_text)
                return True

    # If not found in a single run, try full paragraph text replacement
    for para in tf.paragraphs:
        para_text = ''.join(r.text for r in para.runs)
        if old_text in para_text:
            # Put new text in first run, clear others
            if para.runs:
                para.runs[0].text = para_text.replace(old_text, new_text)
                for run in para.runs[1:]:
                    run.text = ''
                return True
    return False

def set_shape_text(shape, new_text):
    """Set all text in a shape's text frame to new_text, preserving first run's formatting."""
    if not shape.has_text_frame:
        return
    tf = shape.text_frame
    if tf.paragraphs and tf.paragraphs[0].runs:
        # Save font from first run
        first_run = tf.paragraphs[0].runs[0]
        # Clear all paragraphs except first
        for para in tf.paragraphs[1:]:
            p_elem = para._p
            p_elem.getparent().remove(p_elem)
        # Set text in first paragraph first run
        tf.paragraphs[0].runs[0].text = new_text
        for run in tf.paragraphs[0].runs[1:]:
            run.text = ''
    elif tf.paragraphs:
        tf.paragraphs[0].text = new_text

# ─────────────────────────────────────────────
# Helper: copy a slide
# ─────────────────────────────────────────────
def copy_slide(prs, slide_index):
    """Copy a slide and append it to the end of the presentation."""
    template = prs.slides[slide_index]
    blank_slide_layout = prs.slide_layouts[6]  # blank layout
    copied_slide = prs.slides.add_slide(blank_slide_layout)

    # Copy all shapes from template
    for shape in template.shapes:
        el = shape.element
        newel = copy.deepcopy(el)
        copied_slide.shapes._spTree.insert(2, newel)

    return copied_slide

def move_slide_to_index(prs, current_index, target_index):
    """
    Move slide from current_index to target_index in the presentation.
    Uses XML manipulation on the slide ID list.
    """
    xml_slides = prs.slides._sldIdLst
    slides = list(xml_slides)

    # current_index is the last slide (just added), move to target_index
    slide_elem = slides[current_index]
    xml_slides.remove(slide_elem)
    xml_slides.insert(target_index, slide_elem)

# ─────────────────────────────────────────────
# Clear all text from copied slide and set new content
# ─────────────────────────────────────────────
def clear_slide_text(slide):
    """Clear text from all text boxes in a slide."""
    for shape in slide.shapes:
        if shape.has_text_frame:
            for para in shape.text_frame.paragraphs:
                for run in para.runs:
                    run.text = ''

def get_shapes_with_text(slide):
    """Get shapes that have non-empty text."""
    result = []
    for shape in slide.shapes:
        if shape.has_text_frame and shape.text_frame.text.strip():
            result.append((shape.name, shape.text_frame.text.strip()[:80]))
    return result

# ─────────────────────────────────────────────
# NEW SLIDE: Composio Integration
# ─────────────────────────────────────────────
def build_composio_slide(prs, template_slide_index=14):
    """
    Create a new slide for Composio integration based on slide 15 (index 14) template.
    Slide 15 has: title, subtitle, 6 content blocks (rect + text label + text desc), footer.
    """
    new_slide = copy_slide(prs, template_slide_index)

    # Map the template shapes to understand structure:
    # Shape[0]=bg rect, Shape[1]=title TextBox 2, Shape[2]=subtitle TextBox 3,
    # Shape[3]=bg rect 4, then pairs: (rect, label, desc) × 6, last=footer
    # Template slide 15 layout: title at [1], subtitle at [2],
    # content pairs at indices 4-21, footer at [22]

    shapes = list(new_slide.shapes)

    # Title (TextBox 2 equivalent)
    title_shape = None
    subtitle_shape = None
    footer_shape = None

    for shape in shapes:
        if not shape.has_text_frame:
            continue
        t = shape.text_frame.text.strip()
        if 'Real-Time Token Streaming' in t:
            title_shape = shape
        elif 'Server-Sent Events for progressive' in t:
            subtitle_shape = shape
        elif 'SAGE[ai] — Confidential' in t and 'March 2026' in t:
            footer_shape = shape

    if title_shape:
        # Replace title
        for para in title_shape.text_frame.paragraphs:
            for run in para.runs:
                if 'Real-Time Token Streaming' in run.text:
                    run.text = run.text.replace(
                        'Real-Time Token Streaming (Phase 5)',
                        'Composio Integration — 100+ Tools in One Package'
                    )
                    break
        log("  [Composio slide] Set title: 'Composio Integration — 100+ Tools in One Package'")

    if subtitle_shape:
        for para in subtitle_shape.text_frame.paragraphs:
            for run in para.runs:
                if run.text.strip():
                    run.text = 'Connect any SaaS tool to SAGE agents — OAuth handled, no per-integration credential code'
                    break
        log("  [Composio slide] Set subtitle")

    # Content blocks — find the 6 label+description pairs
    # In slide 15 structure: blocks go FastAPI StreamingResponse, generate_stream(),
    # Claude API provider, CLI providers, Ollama provider, Two endpoints
    content_data = [
        ('FastAPI StreamingResponse', 'text/event-stream content-type; browser-native SSE protocol',
         '100+ Apps', 'GitHub, GitLab, Jira, Linear, Slack, Notion, Confluence, Salesforce, HubSpot, Zendesk, Stripe, Google Workspace, and 90+ more'),
        ('generate_stream() method', 'Yields chunks from LLMGateway — same thread-safe singleton, streaming path',
         'OAuth Handled', 'Composio manages auth flows — SAGE never stores third-party credentials'),
        ('Claude API provider', 'Native token-by-token streaming via Anthropic SDK stream context',
         'LangChain Native', 'Uses existing langchain_tools.py — composio:github in project.yaml integrations'),
        ('CLI providers (Gemini, etc)', 'Full generation runs then output is split into 4-word chunks — progressive feel',
         'Web UI', 'New Integrations page: connect apps via OAuth, see active tools per solution'),
        ('Ollama provider', 'stream=True to local REST API — true token-level streaming offline',
         'HITL Connect', 'POST /integrations/composio/connect creates a proposal — approve on Dashboard'),
        ('Two endpoints', 'POST /analyze/stream · POST /agent/stream — both support X-SAGE-Tenant',
         'Zero Code', 'pip install composio-langchain + COMPOSIO_API_KEY — no per-tool code'),
    ]

    replaced_labels = 0
    for old_label, old_desc, new_label, new_desc in content_data:
        # Find and replace label
        for shape in new_slide.shapes:
            if shape.has_text_frame and shape.text_frame.text.strip() == old_label:
                if shape.text_frame.paragraphs and shape.text_frame.paragraphs[0].runs:
                    shape.text_frame.paragraphs[0].runs[0].text = new_label
                    for run in shape.text_frame.paragraphs[0].runs[1:]:
                        run.text = ''
                    replaced_labels += 1
                    break
        # Find and replace description
        for shape in new_slide.shapes:
            if shape.has_text_frame and shape.text_frame.text.strip() == old_desc:
                if shape.text_frame.paragraphs and shape.text_frame.paragraphs[0].runs:
                    shape.text_frame.paragraphs[0].runs[0].text = new_desc
                    for run in shape.text_frame.paragraphs[0].runs[1:]:
                        run.text = ''
                    break

    log(f"  [Composio slide] Replaced {replaced_labels} content block labels and descriptions")

    if footer_shape:
        log("  [Composio slide] Footer preserved: 'SAGE[ai] — Confidential | March 2026'")

    return new_slide


# ─────────────────────────────────────────────
# NEW SLIDE: Hire Agent & Solution Theming
# ─────────────────────────────────────────────
def build_hire_agent_slide(prs, template_slide_index=14):
    """
    Create a new slide for Hire Agent & Solution Theming.
    We use slide 15 (index 14) as template again.
    """
    new_slide = copy_slide(prs, template_slide_index)
    shapes = list(new_slide.shapes)

    title_shape = None
    subtitle_shape = None
    footer_shape = None

    for shape in shapes:
        if not shape.has_text_frame:
            continue
        t = shape.text_frame.text.strip()
        if 'Real-Time Token Streaming' in t:
            title_shape = shape
        elif 'Server-Sent Events for progressive' in t:
            subtitle_shape = shape
        elif 'SAGE[ai] — Confidential' in t and 'March 2026' in t:
            footer_shape = shape

    if title_shape:
        for para in title_shape.text_frame.paragraphs:
            for run in para.runs:
                if 'Real-Time Token Streaming' in run.text:
                    run.text = run.text.replace(
                        'Real-Time Token Streaming (Phase 5)',
                        'Hire Agent & Solution Theming — Runtime Customisation'
                    )
                    break
        log("  [Hire Agent slide] Set title: 'Hire Agent & Solution Theming — Runtime Customisation'")

    if subtitle_shape:
        for para in subtitle_shape.text_frame.paragraphs:
            for run in para.runs:
                if run.text.strip():
                    run.text = 'Extend agents and brand the UI without touching code or YAML files'
                    break
        log("  [Hire Agent slide] Set subtitle")

    # Content blocks for Hire Agent (left) and Solution Theming (right)
    # We reuse the 6 label+description slots, 3 per column conceptually
    # Left column: Hire Agent items, Right column: Theming items
    content_data = [
        ('FastAPI StreamingResponse', 'text/event-stream content-type; browser-native SSE protocol',
         'Hire Agent (Web UI)', 'From the web UI, click Hire Agent — choose icon, name, description, system prompt'),
        ('generate_stream() method', 'Yields chunks from LLMGateway — same thread-safe singleton, streaming path',
         'HITL Approval', 'Creates HITL proposal — approve on Dashboard — role appears in Agents page immediately'),
        ('Claude API provider', 'Native token-by-token streaming via Anthropic SDK stream context',
         'API Backed', 'Backed by POST /agents/hire — optionally add task types (written to tasks.yaml)'),
        ('CLI providers (Gemini, etc)', 'Full generation runs then output is split into 4-word chunks — progressive feel',
         'Solution Theming', 'theme: block in project.yaml — sidebar_bg, accent, badge_bg CSS colour strings'),
        ('Ollama provider', 'stream=True to local REST API — true token-level streaming offline',
         'CSS Variables', 'ThemeProvider writes CSS variables to document.documentElement — switches instantly'),
        ('Two endpoints', 'POST /analyze/stream · POST /agent/stream — both support X-SAGE-Tenant',
         'Example Themes', 'Medical blue, game studio purple, startup green — no Tailwind dependency'),
    ]

    replaced_labels = 0
    for old_label, old_desc, new_label, new_desc in content_data:
        for shape in new_slide.shapes:
            if shape.has_text_frame and shape.text_frame.text.strip() == old_label:
                if shape.text_frame.paragraphs and shape.text_frame.paragraphs[0].runs:
                    shape.text_frame.paragraphs[0].runs[0].text = new_label
                    for run in shape.text_frame.paragraphs[0].runs[1:]:
                        run.text = ''
                    replaced_labels += 1
                    break
        for shape in new_slide.shapes:
            if shape.has_text_frame and shape.text_frame.text.strip() == old_desc:
                if shape.text_frame.paragraphs and shape.text_frame.paragraphs[0].runs:
                    shape.text_frame.paragraphs[0].runs[0].text = new_desc
                    for run in shape.text_frame.paragraphs[0].runs[1:]:
                        run.text = ''
                    break

    log(f"  [Hire Agent slide] Replaced {replaced_labels} content block labels and descriptions")

    if footer_shape:
        log("  [Hire Agent slide] Footer preserved: 'SAGE[ai] — Confidential | March 2026'")

    return new_slide


# ═══════════════════════════════════════════════════════════
#  UPDATE: SageAI_Tech_Pitch.pptx
# ═══════════════════════════════════════════════════════════
def update_tech_pitch():
    log("\n" + "="*60)
    log("UPDATING: SageAI_Tech_Pitch.pptx")
    log("="*60)

    prs = Presentation(TECH_PITCH_PATH)

    # ── Slide 11 (index 10): Developer Experience ──
    slide11 = prs.slides[10]
    log("\nSlide 11 - Developer Experience:")

    # Find "New agent role: add YAML block..." shape and append new line
    shape_agent_role = find_shape_containing(slide11, 'New agent role: add YAML block to prompts.yaml')
    if shape_agent_role:
        # Add new text after existing
        tf = shape_agent_role.text_frame
        current_text = tf.text.strip()
        # Replace the text in runs
        for para in tf.paragraphs:
            for run in para.runs:
                if 'New agent role: add YAML block to prompts.yaml' in run.text:
                    run.text = run.text.replace(
                        'New agent role: add YAML block to prompts.yaml + entry to tasks.yaml — no Python required',
                        'New agent role: add YAML block to prompts.yaml + entry to tasks.yaml — no Python required — or use Hire Agent in the web UI, no YAML editing required'
                    )
                    log("  [OK] Updated 'New agent role' text to add Hire Agent mention")
                    break
    else:
        log("  [WARN] Could not find 'New agent role' shape on slide 11")

    # Find "Hot-reload YAML" shape and append Solution Theming info
    shape_hotreload = find_shape_containing(slide11, 'Hot-reload YAML via /edit-solution-yaml skill')
    if shape_hotreload:
        tf = shape_hotreload.text_frame
        for para in tf.paragraphs:
            for run in para.runs:
                if 'Hot-reload YAML via /edit-solution-yaml skill' in run.text:
                    run.text = run.text.replace(
                        'Hot-reload YAML via /edit-solution-yaml skill — backend reloads without restart',
                        'Hot-reload YAML via /edit-solution-yaml skill — backend reloads without restart'
                    )
                    break
        # We need to find or add the theming line — check if there's a shape after it
        # Look for empty shapes or add new info to existing shape
        # The shape TextBox 19 has the hot-reload text; theming info should go in a new note
        # Since the slide layout may not have room, we update the shape text itself
        for para in tf.paragraphs:
            for run in para.runs:
                if 'Hot-reload YAML via /edit-solution-yaml skill — backend reloads without restart' in run.text:
                    run.text = 'Hot-reload YAML via /edit-solution-yaml skill — backend reloads without restart | Solution theming: add theme: block to project.yaml — sidebar, buttons, badges update instantly'
                    log("  [OK] Updated 'Hot-reload YAML' shape to add Solution Theming mention")
                    break
    else:
        log("  [WARN] Could not find 'Hot-reload YAML' shape on slide 11")

    # ── Slide 16 (index 15): Onboarding Wizard ──
    slide16 = prs.slides[15]
    log("\nSlide 16 - Onboarding Wizard:")

    # Find step 1 shape: "POST /onboarding/generate..."
    shape_step1_desc = find_shape_containing(slide16, 'POST /onboarding/generate with a plain-language')
    if shape_step1_desc:
        tf = shape_step1_desc.text_frame
        for para in tf.paragraphs:
            for run in para.runs:
                if 'POST /onboarding/generate with a plain-language description' in run.text:
                    run.text = run.text.replace(
                        'POST /onboarding/generate with a plain-language description + compliance standards',
                        'POST /onboarding/generate with a plain-language description + compliance standards | Web UI path: visit /onboarding — guided multi-turn conversation, LLM extracts domain info, Generate Solution button appears when ready'
                    )
                    log("  [OK] Updated step 1 (DESCRIBE) to add web UI path note")
                    break
    else:
        log("  [WARN] Could not find step 1 description shape on slide 16")

    # ── Add new slides after Slide 21 (index 20) ──
    # Current slide count = 23 (indices 0-22)
    # We need to insert after index 20 (Temporal Durable Workflows, slide 21)
    # Strategy: add slides at end, then reorder

    log("\nAdding new slides (Composio, Hire Agent) after slide 21:")

    original_slide_count = len(prs.slides)

    # Add Composio slide (copies slide 15/index 14 as template)
    composio_slide = build_composio_slide(prs, template_slide_index=14)
    composio_idx_current = len(prs.slides) - 1  # just appended at end

    # Add Hire Agent slide
    hire_slide = build_hire_agent_slide(prs, template_slide_index=14)
    hire_idx_current = len(prs.slides) - 1  # just appended at end

    # Now reorder: we want Composio at position 21 (after index 20 = Temporal)
    # and Hire Agent at position 22 (after Composio)
    # Current positions: composio at index (original_count), hire at (original_count+1)
    # We want them at indices 21 and 22 respectively
    # After insertion the other slides shift: Slide 22 (Platform Completeness) moves to 24, etc.

    # Move composio_slide to index 21
    move_slide_to_index(prs, len(prs.slides) - 2, 21)
    log(f"  [OK] Composio slide inserted at position 22 (index 21)")

    # After moving composio, hire_slide is now at len-1
    move_slide_to_index(prs, len(prs.slides) - 1, 22)
    log(f"  [OK] Hire Agent slide inserted at position 23 (index 22)")

    # ── Slide 22 now becomes Slide 24 (Platform Completeness) ──
    # After inserting 2 slides, slide 22 (old) = index 21 → now index 23
    # Old slide 22 was "Platform Completeness — 12 Integration Phases"
    # It's now at index 23
    log("\nSlide 24 (was Slide 22) - Platform Completeness (footer update check):")
    slide_platform = prs.slides[23]
    # Check subtitle — update if it says old count
    for shape in slide_platform.shapes:
        if shape.has_text_frame:
            t = shape.text_frame.text
            if 'Platform Completeness' in t and '12 Integration Phases' in t:
                for para in shape.text_frame.paragraphs:
                    for run in para.runs:
                        if '12 Integration Phases' in run.text:
                            run.text = run.text.replace(
                                'Platform Completeness — 12 Integration Phases',
                                'Platform Completeness — 12+ Integration Phases'
                            )
                            log("  [OK] Updated Platform Completeness title to '12+ Integration Phases'")
                            break

    prs.save(TECH_PITCH_PATH)
    log(f"\n[SAVED] {TECH_PITCH_PATH}")


# ═══════════════════════════════════════════════════════════
#  UPDATE: SageAI_Business_Case.pptx
# ═══════════════════════════════════════════════════════════
def update_business_case():
    log("\n" + "="*60)
    log("UPDATING: SageAI_Business_Case.pptx")
    log("="*60)

    prs = Presentation(BUSINESS_CASE_PATH)

    # ── Slide 3 (index 2): SAGE[ai] Always-On AI Engineering Partner ──
    slide3 = prs.slides[2]
    log("\nSlide 3 - Always-On AI Engineering Partner (bullet list update):")

    # Current 7 bullets are in shapes TextBox 4 through TextBox 10
    # New 8 bullets to replace items 7 (TextBox 10: 'Web dashboard...zero CLI required')
    # and add one new item

    old_bullets = [
        'Analyzes error logs in under 60 seconds  (vs 45-60 min manual)',
        'AI code reviews with multi-step ReAct reasoning loop',
        'Monitors Teams, Metabase, GitLab — real-time event detection',
        'Every decision logged to immutable ISO 13485 audit trail',
        'Human-in-the-loop: AI advises, humans decide — always',
        'Learns from every correction via vector memory (RAG)',
        'Web dashboard for all stakeholders — zero CLI required',
    ]

    new_bullets = [
        'Analyzes error logs in under 60 seconds  (vs 45-60 min manual)',
        'AI code reviews with multi-step ReAct reasoning loop',
        'Monitors Teams, Metabase, GitLab — real-time event detection',
        'Every decision logged to immutable ISO 13485 audit trail',
        'Human-in-the-loop: AI advises, humans decide — always',
        'Learns from every correction via vector memory (RAG)',
        '100+ tool integrations via Composio — GitHub, Jira, Salesforce, Slack and more',
        'Web dashboard for all stakeholders — hire new agent roles, connect tools, approve proposals',
    ]

    # Find each old bullet shape and update it
    updated_count = 0
    for old_text in old_bullets:
        shape = find_shape_containing(slide3, old_text)
        if shape:
            for para in shape.text_frame.paragraphs:
                for run in para.runs:
                    if old_text in run.text:
                        # Which new text corresponds?
                        idx = old_bullets.index(old_text)
                        new_text = new_bullets[idx]
                        if old_text != new_text:
                            run.text = run.text.replace(old_text, new_text)
                            log(f"  [OK] Updated bullet: '{old_text[:50]}...' → '{new_text[:50]}...'")
                        updated_count += 1
                        break
        else:
            log(f"  [WARN] Could not find bullet shape: '{old_text[:60]}'")

    # The 8th bullet needs to be added — find the shape for old bullet 7 (last one)
    # Old: 'Web dashboard for all stakeholders — zero CLI required'
    # New: 'Web dashboard for all stakeholders — hire new agent roles, connect tools, approve proposals'
    # That was handled above (index 6 → new_bullets[6] and [7])
    # But we only have 7 shape slots. We need to either:
    # (a) Update the 7th shape to the new 7th text (Composio line)
    # (b) The 8th new bullet needs to go somewhere

    # Let's check: we replaced bullet[6] with new_bullets[6] (Composio line) above
    # But new_bullets[7] (Web dashboard updated) has no shape yet
    # Find the shape that now has the Composio text (was old[6]) and check if there's
    # a shape after it, or we need to clone one

    # Actually looking at the slide: shapes TextBox 4 through TextBox 10 = 7 bullets
    # We need to update the last one AND find the shape for old[6] was index 6 in old_bullets
    # Old[6] = 'Web dashboard...' → replaced with new[6] = '100+ tool integrations...' ✓
    # New[7] = 'Web dashboard...' — this needs to go somewhere

    # Find the shape that was updated to Composio text, and check the next sibling
    # Actually the simplest approach: find 'Web dashboard for all stakeholders — zero CLI required'
    # was already replaced above to Composio. We need to find if there's another empty
    # shape or we need to use the rect area.

    # Since there's no 8th shape, let's find all bullet shapes by looking at shape positions
    # and copy the last one's formatting for the new 8th bullet.
    # For now, let's find the shape that now has Composio text and look for any unused shape

    # Find the shape with '100+ tool integrations' (just set) to get its position
    composio_shape = find_shape_containing(slide3, '100+ tool integrations via Composio')
    web_dashboard_old_shape = find_shape_containing(slide3, 'Web dashboard for all stakeholders — zero CLI required')

    if web_dashboard_old_shape:
        # This shape still has old text (not yet updated since old[6] was replaced with new[6]=composio)
        # Wait — old[6] = 'Web dashboard...zero CLI required' and new[6] = '100+ tool integrations...'
        # So this shape should now have Composio text. Let me re-check the logic.
        # Actually: replace old_bullets[6] = 'Web dashboard...zero CLI' with new_bullets[6] = '100+ tool integrations'
        # So 'Web dashboard...zero CLI' shape → now has Composio text
        # new_bullets[7] = 'Web dashboard...hire new agent roles...' needs a NEW shape

        # We need to ADD a new shape for bullet 8
        # Clone the Composio shape (which was old bullet 7 shape)
        pass

    # Find the current last bullet shape (should now have Composio text)
    current_last_bullet = find_shape_containing(slide3, '100+ tool integrations via Composio')

    if current_last_bullet:
        # Clone this shape and position it below
        from pptx.util import Emu
        el = current_last_bullet.element
        newel = copy.deepcopy(el)
        slide3.shapes._spTree.append(newel)

        # Find the newly added shape (last one appended)
        new_shape = None
        for s in slide3.shapes:
            if s.has_text_frame and '100+ tool integrations via Composio' in s.text_frame.text:
                if s != current_last_bullet:
                    new_shape = s
                    break

        if new_shape is None:
            # If same text, get all matching and pick the second
            matches = [s for s in slide3.shapes if s.has_text_frame and '100+ tool integrations' in s.text_frame.text]
            if len(matches) > 1:
                new_shape = matches[-1]

        if new_shape:
            # Set text to 8th bullet
            eighth_text = 'Web dashboard for all stakeholders — hire new agent roles, connect tools, approve proposals'
            for para in new_shape.text_frame.paragraphs:
                for run in para.runs:
                    if run.text.strip():
                        run.text = eighth_text
                        break

            # Reposition: move it below the 7th bullet
            # Get position of the 7th bullet (current_last_bullet)
            orig_top = current_last_bullet.top
            orig_height = current_last_bullet.height
            new_shape.top = orig_top + orig_height + Emu(76200)  # ~60000 EMU gap
            log("  [OK] Added 8th bullet: 'Web dashboard...hire new agent roles...'")
        else:
            log("  [WARN] Could not find newly cloned shape for 8th bullet")
    else:
        log("  [WARN] Could not find Composio shape to clone for 8th bullet")

    # ── Slide 8 (index 7): What SAGE Does Today — subtitle update ──
    slide8 = prs.slides[7]
    log("\nSlide 8 - What SAGE Does Today (subtitle update):")

    shape_subtitle = find_shape_containing(slide8, 'Production-ready capabilities across 12 integration phases')
    if shape_subtitle:
        tf = shape_subtitle.text_frame
        for para in tf.paragraphs:
            for run in para.runs:
                if 'Production-ready capabilities across 12 integration phases' in run.text:
                    run.text = run.text.replace(
                        'Production-ready capabilities across 12 integration phases — available now',
                        'Production-ready capabilities — agents, tools, approvals, and 100+ integrations — available now'
                    )
                    log("  [OK] Updated subtitle: 'Production-ready capabilities — agents, tools, approvals, and 100+ integrations — available now'")
                    break
    else:
        log("  [WARN] Could not find subtitle shape on slide 8")

    prs.save(BUSINESS_CASE_PATH)
    log(f"\n[SAVED] {BUSINESS_CASE_PATH}")


# ═══════════════════════════════════════════════════════════
#  UPDATE: SageAI_Investor_Pitch.pptx
# ═══════════════════════════════════════════════════════════
def update_investor_pitch():
    log("\n" + "="*60)
    log("UPDATING: SageAI_Investor_Pitch.pptx")
    log("="*60)

    prs = Presentation(INVESTOR_PITCH_PATH)

    # ── Slide 7 (index 6): Five Defensible Technology Moats — add 6th moat ──
    slide7 = prs.slides[6]
    log("\nSlide 7 - Five Defensible Technology Moats (add 6th moat):")

    # Current moats are in pairs: (Rectangle, TextBox label, TextBox description)
    # Last moat: 'Provider Agnostic' / 'A single config change swaps LLM providers...'
    # We need to add: 'Integration Network Effect' after it

    # Find the 'Provider Agnostic' label shape
    provider_agnostic_shape = find_shape_containing(slide7, 'Provider Agnostic')
    provider_agnostic_desc = find_shape_containing(slide7, 'A single config change swaps LLM providers')

    if provider_agnostic_shape and provider_agnostic_desc:
        # Clone the rect + label + desc pattern from Provider Agnostic moat
        # First find the rectangle associated with Provider Agnostic
        # The structure is: Rectangle 16, TextBox 17 (label), TextBox 18 (desc)

        # Clone the label shape for new moat label
        from pptx.util import Emu

        # Get positions
        label_top = provider_agnostic_shape.top
        label_height = provider_agnostic_shape.height
        desc_top = provider_agnostic_desc.top

        # Find the associated rectangle (should be just above the label)
        moat_rect = None
        for shape in slide7.shapes:
            if not shape.has_text_frame:
                # Check if this rectangle is near Provider Agnostic label
                if abs(shape.top - label_top) < Emu(200000) and abs(shape.left - provider_agnostic_shape.left) < Emu(200000):
                    moat_rect = shape
                    break

        # Calculate new positions (below existing moat)
        new_top_offset = provider_agnostic_shape.top - provider_agnostic_desc.top
        # Estimate row height from existing moats
        # Moat rows seem to be about 600000 EMU tall (rough estimate)
        # Let's use the gap between consecutive label tops

        # Find all moat labels to compute spacing
        moat_labels = []
        moat_names = ['Compliance by Design', 'Compounding Intelligence', 'YAML-First Architecture',
                      'No API Key Lock-In', 'Provider Agnostic']
        for name in moat_names:
            s = find_shape_containing(slide7, name)
            if s:
                moat_labels.append(s)

        if len(moat_labels) >= 2:
            row_height = moat_labels[1].top - moat_labels[0].top
        else:
            row_height = Emu(600000)

        new_label_top = provider_agnostic_shape.top + row_height
        new_desc_top = provider_agnostic_desc.top + row_height

        # Clone rectangle if found
        if moat_rect:
            rect_el = moat_rect.element
            new_rect_el = copy.deepcopy(rect_el)
            slide7.shapes._spTree.append(new_rect_el)
            # Find and reposition
            for s in slide7.shapes:
                if not s.has_text_frame:
                    if s.element is new_rect_el:
                        s.top = moat_rect.top + row_height
                        break

        # Clone label shape
        label_el = provider_agnostic_shape.element
        new_label_el = copy.deepcopy(label_el)
        slide7.shapes._spTree.append(new_label_el)

        # Find new label shape and update
        new_label_shape = None
        for s in slide7.shapes:
            if s.has_text_frame and 'Provider Agnostic' in s.text_frame.text and s.element is new_label_el:
                new_label_shape = s
                break

        if new_label_shape is None:
            # Try finding by element identity more carefully
            all_provider_shapes = [s for s in slide7.shapes if s.has_text_frame and 'Provider Agnostic' in s.text_frame.text]
            if len(all_provider_shapes) > 1:
                new_label_shape = all_provider_shapes[-1]

        if new_label_shape:
            new_label_shape.top = new_label_top
            for para in new_label_shape.text_frame.paragraphs:
                for run in para.runs:
                    if run.text.strip():
                        run.text = 'Integration Network Effect'
                        break
            log("  [OK] Added 6th moat label: 'Integration Network Effect'")
        else:
            log("  [WARN] Could not find new label shape after cloning")

        # Clone description shape
        desc_el = provider_agnostic_desc.element
        new_desc_el = copy.deepcopy(desc_el)
        slide7.shapes._spTree.append(new_desc_el)

        new_desc_shape = None
        for s in slide7.shapes:
            if s.has_text_frame and 'A single config change swaps LLM providers' in s.text_frame.text and s.element is new_desc_el:
                new_desc_shape = s
                break

        if new_desc_shape is None:
            all_provider_desc = [s for s in slide7.shapes if s.has_text_frame and 'A single config change swaps LLM providers' in s.text_frame.text]
            if len(all_provider_desc) > 1:
                new_desc_shape = all_provider_desc[-1]

        new_moat_text = '100+ pre-built tool integrations via Composio. Each new connector increases platform value for all customers. Solution marketplace (APM vision) creates community flywheel — share YAML configs across industries.'

        if new_desc_shape:
            new_desc_shape.top = new_desc_top
            for para in new_desc_shape.text_frame.paragraphs:
                for run in para.runs:
                    if run.text.strip():
                        run.text = new_moat_text
                        for r2 in para.runs[1:]:
                            r2.text = ''
                        break
            log("  [OK] Added 6th moat description: 'Integration Network Effect — 100+ pre-built...'")
        else:
            log("  [WARN] Could not find new description shape after cloning")
    else:
        log("  [WARN] Could not find 'Provider Agnostic' shapes on slide 7")

    # ── Slide 11 (index 10): Traction ──
    slide11 = prs.slides[10]
    log("\nSlide 11 - Traction (metrics update):")

    # Find "383+" shape
    shape_383 = find_shape_containing(slide11, '383+')
    if shape_383:
        log("  [OK] Found '383+' metric shape — keeping as-is (still accurate)")
    else:
        log("  [WARN] Could not find '383+' metric shape on slide 11")

    # Find "12 integrations live:" text shape
    shape_integrations = find_shape_containing(slide11, '12 integrations live:')
    if shape_integrations:
        tf = shape_integrations.text_frame
        for para in tf.paragraphs:
            for run in para.runs:
                if '12 integrations live:' in run.text:
                    run.text = run.text.replace(
                        '12 integrations live: SSE streaming, onboarding wizard, Slack, LangGraph, Temporal, eval, multi-tenant',
                        'New capabilities live: Composio (100+ tool integrations), Hire Agent (web UI role creation), conversational onboarding wizard, solution theming system'
                    )
                    log("  [OK] Updated '12 integrations live:' to new capabilities text")
                    break
    else:
        log("  [WARN] Could not find '12 integrations live:' shape on slide 11")

    # Find "Pilot in progress" — keep as-is but verify it's there
    shape_pilot = find_shape_containing(slide11, 'Pilot in progress')
    if shape_pilot:
        log("  [OK] Found 'Pilot in progress' shape — keeping as-is")
    else:
        log("  [WARN] Could not find 'Pilot in progress' shape on slide 11")

    prs.save(INVESTOR_PITCH_PATH)
    log(f"\n[SAVED] {INVESTOR_PITCH_PATH}")


# ═══════════════════════════════════════════════════════════
#  MAIN
# ═══════════════════════════════════════════════════════════
if __name__ == '__main__':
    print("Starting SAGE PowerPoint update script...")
    print()

    try:
        update_tech_pitch()
    except Exception as e:
        print(f"\n[ERROR] Tech Pitch update failed: {e}")
        import traceback
        traceback.print_exc()

    try:
        update_business_case()
    except Exception as e:
        print(f"\n[ERROR] Business Case update failed: {e}")
        import traceback
        traceback.print_exc()

    try:
        update_investor_pitch()
    except Exception as e:
        print(f"\n[ERROR] Investor Pitch update failed: {e}")
        import traceback
        traceback.print_exc()

    print("\n" + "="*60)
    print("SCRIPT COMPLETE")
    print("="*60)
    print(f"\nTotal changes logged: {len(changes_log)}")
