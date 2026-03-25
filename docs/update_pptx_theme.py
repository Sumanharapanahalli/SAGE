"""
Update Build Orchestrator slides in SAGE pptx files from dark theme to white+green.
Only touches slides we added (identified by title content).
"""
from pptx import Presentation
from pptx.util import Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.dml import MSO_THEME_COLOR
import os

# Color mapping: old dark -> new light
# Shape/background fills
FILL_MAP = {
    RGBColor(0x0D, 0x1B, 0x2A): RGBColor(0xEC, 0xFD, 0xF5),  # dark navy header -> light green
    RGBColor(0x1E, 0x29, 0x3B): RGBColor(0xEC, 0xFD, 0xF5),  # dark card bg -> light green
    RGBColor(0x00, 0x78, 0xD4): RGBColor(0x10, 0xB9, 0x81),  # blue table header -> emerald
}

# Text color mapping
TEXT_MAP = {
    RGBColor(0xFF, 0xFF, 0xFF): RGBColor(0x06, 0x5F, 0x46),  # white title -> dark green
    RGBColor(0x00, 0xB4, 0xD8): RGBColor(0x4B, 0x55, 0x63),  # cyan subtitle -> medium gray
    RGBColor(0x38, 0xBD, 0xF8): RGBColor(0x05, 0x96, 0x69),  # blue accent -> emerald
    RGBColor(0xCB, 0xD5, 0xE1): RGBColor(0x37, 0x41, 0x51),  # light body -> dark gray
    RGBColor(0x94, 0xA3, 0xB8): RGBColor(0x6B, 0x72, 0x80),  # gray notes -> medium gray
    RGBColor(0x64, 0x74, 0x87): RGBColor(0x6B, 0x72, 0x80),  # footer -> medium gray
}

# Table cell text: 0D1B2A (dark navy text) -> 374151 (dark gray, readable on light bg)
TABLE_TEXT_MAP = {
    RGBColor(0x0D, 0x1B, 0x2A): RGBColor(0x37, 0x41, 0x51),
}

# Table cell fill map
TABLE_FILL_MAP = {
    RGBColor(0x00, 0x78, 0xD4): RGBColor(0x10, 0xB9, 0x81),  # blue header -> emerald
    RGBColor(0xF0, 0xF4, 0xF8): RGBColor(0xEC, 0xFD, 0xF5),  # light gray stripe -> light green
    # White cells stay white
}

# Slide background -> white
BG_WHITE = RGBColor(0xFF, 0xFF, 0xFF)


def is_target_slide(slide):
    """Check if slide is one we added (Build Orchestrator / Agentic Patterns / Platform Completeness)."""
    keywords = [
        "Build Orchestrator",
        "Agentic Patterns",
        "Platform Completeness",
    ]
    for shape in slide.shapes:
        if shape.has_text_frame:
            text = shape.text_frame.text
            for kw in keywords:
                if kw in text:
                    return True
    return False


def remap_color(rgb, color_map):
    """Return mapped color or None."""
    if rgb is None:
        return None
    for old, new in color_map.items():
        if rgb == old:
            return new
    return None


def update_shape_fill(shape, fill_map):
    """Update shape fill color if it matches."""
    try:
        fill = shape.fill
        if fill.type is not None:
            try:
                rgb = fill.fore_color.rgb
                new_rgb = remap_color(rgb, fill_map)
                if new_rgb:
                    fill.solid()
                    fill.fore_color.rgb = new_rgb
            except (AttributeError, TypeError):
                pass
    except (AttributeError, TypeError):
        pass


def update_text_colors(shape, text_map):
    """Update text run colors in a shape."""
    if not shape.has_text_frame:
        return
    for para in shape.text_frame.paragraphs:
        for run in para.runs:
            try:
                rgb = run.font.color.rgb
                new_rgb = remap_color(rgb, text_map)
                if new_rgb:
                    run.font.color.rgb = new_rgb
            except (AttributeError, TypeError):
                pass


def update_table(shape, table_fill_map, table_text_map_, text_map):
    """Update table cell fills and text colors."""
    if not hasattr(shape, 'table'):
        return
    tbl = shape.table
    for row_idx in range(len(tbl.rows)):
        for col_idx in range(len(tbl.columns)):
            cell = tbl.cell(row_idx, col_idx)
            # Cell fill
            try:
                rgb = cell.fill.fore_color.rgb
                new_rgb = remap_color(rgb, table_fill_map)
                if new_rgb:
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = new_rgb
            except (AttributeError, TypeError):
                pass
            # Cell text
            for para in cell.text_frame.paragraphs:
                for run in para.runs:
                    try:
                        rgb = run.font.color.rgb
                        # Check table-specific map first, then general text map
                        new_rgb = remap_color(rgb, table_text_map_)
                        if new_rgb is None:
                            new_rgb = remap_color(rgb, text_map)
                        if new_rgb:
                            run.font.color.rgb = new_rgb
                    except (AttributeError, TypeError):
                        pass


def update_slide_background(slide):
    """Set slide background to white."""
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = BG_WHITE


def process_file(filepath):
    prs = Presentation(filepath)
    updated_count = 0

    for slide in prs.slides:
        if not is_target_slide(slide):
            continue

        updated_count += 1
        print(f"  Updating: {[s.text_frame.text[:60] for s in slide.shapes if s.has_text_frame][:1]}")

        # Set slide background to white
        update_slide_background(slide)

        for shape in slide.shapes:
            # Update shape fills
            update_shape_fill(shape, FILL_MAP)
            # Update text colors
            update_text_colors(shape, TEXT_MAP)
            # Update tables
            if hasattr(shape, 'table'):
                update_table(shape, TABLE_FILL_MAP, TABLE_TEXT_MAP, TEXT_MAP)

    if updated_count > 0:
        prs.save(filepath)
        print(f"  Saved {filepath} ({updated_count} slides updated)")
    else:
        print(f"  No target slides found in {filepath}")


if __name__ == "__main__":
    docs_dir = os.path.dirname(os.path.abspath(__file__))
    files = [
        os.path.join(docs_dir, "SageAI_Tech_Pitch.pptx"),
        os.path.join(docs_dir, "SageAI_Business_Case.pptx"),
        os.path.join(docs_dir, "SageAI_Investor_Pitch.pptx"),
    ]

    for f in files:
        print(f"\n=== Processing {os.path.basename(f)} ===")
        if os.path.exists(f):
            process_file(f)
        else:
            print(f"  File not found: {f}")
