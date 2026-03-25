#!/usr/bin/env python3
"""Add Build Orchestrator slides to all three SAGE presentation decks."""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
import copy
from lxml import etree

# ── Style constants — White + Emerald Green (NO blue, NO black) ─────
BG_DARK    = RGBColor(0x1F, 0x2A, 0x37)   # gray-800 charcoal (NOT black)
CARD_BG    = RGBColor(0xEC, 0xFD, 0xF5)   # emerald-50 light green cards
WHITE      = RGBColor(0xff, 0xff, 0xff)
SUBTITLE_C = RGBColor(0x04, 0x78, 0x57)   # emerald-700 dark accent
ACCENT     = RGBColor(0x10, 0xB9, 0x81)   # emerald-500 primary accent
BODY_C     = RGBColor(0x1F, 0x2A, 0x37)   # gray-800 body text
FOOTER_C   = RGBColor(0x4B, 0x55, 0x63)   # gray-600 footer text
MUTED      = RGBColor(0x6B, 0x72, 0x80)   # gray-500 muted text

TITLE_SZ   = 381000    # 30pt
SUBTITLE_SZ = 203200   # 16pt
SECTION_SZ = 190500    # 15pt
BODY_SZ    = 165100    # 13pt
SMALL_SZ   = 152400    # 12pt
FOOTER_SZ  = 114300    # 9pt
TINY_SZ    = 139700    # 11pt

FOOTER_TEXT = "SAGE[ai] \u2014 Open Source (MIT) | github.com/Sumanharapanahalli/sage | March 2026"

# Slide dimensions (from existing files)
SLIDE_W = 12188952
SLIDE_H = 6858000

# Common layout positions (from existing slides)
HEADER_H   = 1325880
TITLE_LEFT = 411480
TITLE_TOP  = 109728
TITLE_W    = 11338560
TITLE_H    = 731520
SUB_TOP    = 804672
SUB_H      = 438912
CONTENT_LEFT = 457200
CONTENT_W  = 11247120
FOOTER_BAR_TOP = 6492240
FOOTER_BAR_H = 365760
FOOTER_TXT_LEFT = 274320
FOOTER_TXT_TOP = 6510528
FOOTER_TXT_W = 8229600
FOOTER_TXT_H = 274320


def add_run(para, text, size, bold=False, color=WHITE, font_name="Calibri"):
    """Add a styled run to a paragraph."""
    run = para.add_run()
    run.text = text
    run.font.size = size
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = font_name
    return run


def add_header_bar(slide):
    """Add dark header rectangle at top."""
    rect = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, HEADER_H)
    rect.fill.solid()
    rect.fill.fore_color.rgb = BG_DARK
    rect.line.fill.background()
    return rect


def add_footer(slide):
    """Add footer bar and text."""
    rect = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, FOOTER_BAR_TOP, SLIDE_W, FOOTER_BAR_H)
    rect.fill.solid()
    rect.fill.fore_color.rgb = BG_DARK
    rect.line.fill.background()

    tb = slide.shapes.add_textbox(FOOTER_TXT_LEFT, FOOTER_TXT_TOP, FOOTER_TXT_W, FOOTER_TXT_H)
    para = tb.text_frame.paragraphs[0]
    add_run(para, FOOTER_TEXT, FOOTER_SZ, bold=False, color=FOOTER_C)
    return rect


def add_title(slide, title_text, subtitle_text):
    """Add title and subtitle text boxes."""
    tb = slide.shapes.add_textbox(TITLE_LEFT, TITLE_TOP, TITLE_W, TITLE_H)
    para = tb.text_frame.paragraphs[0]
    add_run(para, title_text, TITLE_SZ, bold=True, color=WHITE)

    tb2 = slide.shapes.add_textbox(TITLE_LEFT, SUB_TOP, TITLE_W, SUB_H)
    para2 = tb2.text_frame.paragraphs[0]
    add_run(para2, subtitle_text, SUBTITLE_SZ, bold=False, color=SUBTITLE_C)


def add_card(slide, left, top, width, height, header_text, body_text):
    """Add a card with dark background, header, and body text."""
    rect = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    rect.fill.solid()
    rect.fill.fore_color.rgb = CARD_BG
    rect.line.fill.background()

    # Header
    padding = Emu(91440)  # ~0.1"
    tb = slide.shapes.add_textbox(left + padding, top + padding, width - 2*padding, Emu(274320))
    para = tb.text_frame.paragraphs[0]
    add_run(para, header_text, SECTION_SZ, bold=True, color=ACCENT)

    # Body
    tb2 = slide.shapes.add_textbox(left + padding, top + Emu(320040), width - 2*padding, height - Emu(365760))
    tb2.text_frame.word_wrap = True
    para2 = tb2.text_frame.paragraphs[0]
    add_run(para2, body_text, BODY_SZ, bold=False, color=BODY_C)


def add_bottom_note(slide, text, top=None):
    """Add a bottom note above the footer."""
    if top is None:
        top = FOOTER_BAR_TOP - Emu(365760)
    tb = slide.shapes.add_textbox(CONTENT_LEFT, top, CONTENT_W, Emu(274320))
    tb.text_frame.word_wrap = True
    para = tb.text_frame.paragraphs[0]
    add_run(para, text, TINY_SZ, bold=False, color=MUTED)


def move_slide_to_position(prs, slide_index_from, target_position):
    """Move a slide from one position to another in the slide list."""
    sldIdLst = prs.slides._sldIdLst
    items = list(sldIdLst)
    item = items[slide_index_from]
    sldIdLst.remove(item)
    items_after = list(sldIdLst)
    if target_position >= len(items_after):
        sldIdLst.append(item)
    else:
        ref = items_after[target_position]
        sldIdLst.insert(sldIdLst.index(ref), item)


# ═══════════════════════════════════════════════════════════════════════
# 1. TECH PITCH — Add 2 slides before last 2
# ═══════════════════════════════════════════════════════════════════════
def update_tech_pitch():
    path = "docs/SageAI_Tech_Pitch.pptx"
    prs = Presentation(path)
    layout = prs.slide_layouts[6]  # Blank
    orig_count = len(prs.slides)

    # ── Slide A: Build Orchestrator — 0→1→N Pipeline ──
    slide_a = prs.slides.add_slide(layout)
    add_header_bar(slide_a)
    add_title(slide_a, "Build Orchestrator \u2014 0\u21921\u2192N Pipeline",
              "From plain-English description to working product")

    # 5 pipeline step cards arranged horizontally
    card_top = Emu(1508760)
    card_h = Emu(2743200)  # tall enough for text
    gap = Emu(91440)
    total_w = CONTENT_W
    card_w = (total_w - 4 * gap) // 5

    steps = [
        ("1  DECOMPOSE", "PlannerAgent breaks idea into parallel tasks (9 types: BACKEND, FRONTEND, TESTS, etc.)"),
        ("2  CRITIC REVIEW", "Actor-critic loop: Builder\u2194Critic iterates until score \u2265 threshold (default 70/100)"),
        ("3  EXECUTE", "Wave-based parallel execution via ReAct pattern (Thought\u2192Action\u2192Observation\u2192Status)"),
        ("4  INTEGRATE", "Merge results, run tests, generate diff \u2014 critic reviews integration"),
        ("5  APPROVE", "Human sees artifact + critic report + revision history \u2192 approve or reject"),
    ]

    for i, (header, body) in enumerate(steps):
        left = CONTENT_LEFT + i * (card_w + gap)
        add_card(slide_a, left, card_top, card_w, card_h, header, body)

    add_bottom_note(slide_a, "3-tier degradation: External Open SWE \u2192 LangGraph workflow \u2192 ReAct LLM fallback",
                    top=Emu(4480560))
    add_footer(slide_a)

    # ── Slide B: Agentic Patterns for Built Products ──
    slide_b = prs.slides.add_slide(layout)
    add_header_bar(slide_b)
    add_title(slide_b, "Agentic Patterns for Built Products",
              "Products built by SAGE can themselves use any agentic architecture")

    # 2 columns of 5 patterns each
    col_w = (CONTENT_W - Emu(182880)) // 2
    col1_left = CONTENT_LEFT
    col2_left = CONTENT_LEFT + col_w + Emu(182880)

    left_patterns = [
        "Single-Agent",
        "Multi-Agent Sequential",
        "Multi-Agent Parallel",
        "Review & Critique",
        "Multi-Agent Coordinator",
    ]
    right_patterns = [
        "Hierarchical Decomposition",
        "Multi-Agent Swarm",
        "ReAct (Reason+Act)",
        "Human-in-the-Loop",
        "Iterative Refinement",
    ]

    row_h = Emu(457200)
    start_top = Emu(1508760)

    for i, pattern in enumerate(left_patterns):
        top = start_top + i * (row_h + Emu(91440))
        rect = slide_b.shapes.add_shape(MSO_SHAPE.RECTANGLE, col1_left, top, col_w, row_h)
        rect.fill.solid()
        rect.fill.fore_color.rgb = CARD_BG
        rect.line.fill.background()
        tb = slide_b.shapes.add_textbox(col1_left + Emu(137160), top + Emu(109728), col_w - Emu(274320), row_h - Emu(137160))
        para = tb.text_frame.paragraphs[0]
        add_run(para, f"{i+1}.  ", BODY_SZ, bold=True, color=ACCENT)
        add_run(para, pattern, BODY_SZ, bold=False, color=BODY_C)

    for i, pattern in enumerate(right_patterns):
        top = start_top + i * (row_h + Emu(91440))
        rect = slide_b.shapes.add_shape(MSO_SHAPE.RECTANGLE, col2_left, top, col_w, row_h)
        rect.fill.solid()
        rect.fill.fore_color.rgb = CARD_BG
        rect.line.fill.background()
        tb = slide_b.shapes.add_textbox(col2_left + Emu(137160), top + Emu(109728), col_w - Emu(274320), row_h - Emu(137160))
        para = tb.text_frame.paragraphs[0]
        add_run(para, f"{i+6}.  ", BODY_SZ, bold=True, color=ACCENT)
        add_run(para, pattern, BODY_SZ, bold=False, color=BODY_C)

    add_bottom_note(slide_b, "3 HITL levels: minimal (2 gates) \u00b7 standard (3 gates) \u00b7 strict (every stage)",
                    top=Emu(4480560 + 457200))
    add_footer(slide_b)

    # Move the 2 new slides (at end) to before the last 2 original slides
    # Original: [0..26] (27 slides), new are at [27] and [28]
    # Want them at positions [25] and [26], pushing original [25],[26] to [27],[28]
    move_slide_to_position(prs, orig_count, orig_count - 2)      # slide A
    move_slide_to_position(prs, orig_count + 1, orig_count - 1)  # slide B

    prs.save(path)
    print(f"Tech Pitch: saved with {len(prs.slides)} slides")


# ═══════════════════════════════════════════════════════════════════════
# 2. BUSINESS CASE — Add 1 slide before last slide
# ═══════════════════════════════════════════════════════════════════════
def update_business_case():
    path = "docs/SageAI_Business_Case.pptx"
    prs = Presentation(path)
    layout = prs.slide_layouts[6]
    orig_count = len(prs.slides)

    slide = prs.slides.add_slide(layout)
    add_header_bar(slide)
    add_title(slide, "Build Orchestrator \u2014 Product Factory",
              "Describe a product in plain English \u2192 get a working codebase")

    # 3 value prop cards
    card_top = Emu(1508760)
    card_h = Emu(3200400)
    gap = Emu(182880)
    card_w = (CONTENT_W - 2 * gap) // 3

    props = [
        ("10x Faster", "Parallel agent execution with dependency-aware wave scheduling"),
        ("Quality Built In", "Critic agent reviews every artifact before human sees it (actor-critic pattern)"),
        ("Human Always in Control", "3 configurable HITL levels: minimal, standard, strict \u2014 compliance guaranteed"),
    ]

    for i, (header, body) in enumerate(props):
        left = CONTENT_LEFT + i * (card_w + gap)
        add_card(slide, left, card_top, card_w, card_h, header, body)

    add_bottom_note(slide, "Supports 10 agentic AI patterns (Google Cloud reference) \u00b7 9 task types \u00b7 3-tier degradation",
                    top=Emu(4937760))
    add_footer(slide)

    # Move to before last original slide
    move_slide_to_position(prs, orig_count, orig_count - 1)

    prs.save(path)
    print(f"Business Case: saved with {len(prs.slides)} slides")


# ═══════════════════════════════════════════════════════════════════════
# 3. INVESTOR PITCH — Add 1 slide before last slide
# ═══════════════════════════════════════════════════════════════════════
def update_investor_pitch():
    path = "docs/SageAI_Investor_Pitch.pptx"
    prs = Presentation(path)
    layout = prs.slide_layouts[6]
    orig_count = len(prs.slides)

    slide = prs.slides.add_slide(layout)
    add_header_bar(slide)
    add_title(slide, "Build Orchestrator \u2014 The Meta-Agent Factory",
              "SAGE doesn\u2019t just assist development \u2014 it builds entire products autonomously")

    # Pipeline flow
    flow_top = Emu(1508760)
    flow_tb = slide.shapes.add_textbox(CONTENT_LEFT, flow_top, CONTENT_W, Emu(457200))
    flow_tb.text_frame.word_wrap = True
    flow_para = flow_tb.text_frame.paragraphs[0]
    flow_para.alignment = PP_ALIGN.CENTER
    steps = ["Idea", "Decompose", "Critic", "Approve", "Build", "Critic", "Approve", "Ship"]
    for i, step in enumerate(steps):
        add_run(flow_para, step, SECTION_SZ, bold=True, color=ACCENT)
        if i < len(steps) - 1:
            add_run(flow_para, "  \u2192  ", SECTION_SZ, bold=False, color=MUTED)

    # 3 differentiation cards
    card_top = Emu(2194560)
    card_h = Emu(2743200)
    gap = Emu(182880)
    card_w = (CONTENT_W - 2 * gap) // 3

    diffs = [
        ("Actor-Critic Pattern", "Every output stress-tested before humans see it. Score improves across iterations."),
        ("10 Agentic Architectures", "Built products can use any Google Cloud reference pattern \u2014 not just chat bots."),
        ("Compounding Intelligence", "Critic feedback stored in vector memory \u2014 quality improves with every build."),
    ]

    for i, (header, body) in enumerate(diffs):
        left = CONTENT_LEFT + i * (card_w + gap)
        add_card(slide, left, card_top, card_w, card_h, header, body)

    add_bottom_note(slide, "One human + SAGE = billion-dollar company infrastructure",
                    top=Emu(5166360))
    add_footer(slide)

    # Move to before last original slide
    move_slide_to_position(prs, orig_count, orig_count - 1)

    prs.save(path)
    print(f"Investor Pitch: saved with {len(prs.slides)} slides")


# ═══════════════════════════════════════════════════════════════════════
# Run all updates
# ═══════════════════════════════════════════════════════════════════════
if __name__ == "__main__":
    update_tech_pitch()
    update_business_case()
    update_investor_pitch()

    # Verify
    print("\n=== Verification ===")
    for path in ["docs/SageAI_Tech_Pitch.pptx", "docs/SageAI_Business_Case.pptx", "docs/SageAI_Investor_Pitch.pptx"]:
        prs = Presentation(path)
        print(f"\n{path}: {len(prs.slides)} slides")
        for i, slide in enumerate(prs.slides):
            texts = []
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        t = para.text.strip()
                        if t and t != FOOTER_TEXT:
                            texts.append(t)
                            break
            title = texts[0] if texts else "(no text)"
            marker = " <<<NEW" if "Build Orchestrator" in title or "Agentic Patterns" in title else ""
            print(f"  Slide {i+1}: {title[:80]}{marker}")
