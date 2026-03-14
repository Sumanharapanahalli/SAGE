"""
Generate / update SAGE presentations:
  1. SageAI_Business_Case.pptx  — existing deck, refresh + add cost disclaimer
  2. SageAI_Tech_Pitch.pptx     — new engineering / technical pitch deck
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.util import Inches, Pt
import copy, os

# ── colour palette ────────────────────────────────────────────────────────────
SAGE_DARK   = RGBColor(0x0D, 0x1B, 0x2A)   # near-black navy
SAGE_BLUE   = RGBColor(0x00, 0x78, 0xD4)   # Microsoft-style blue
SAGE_TEAL   = RGBColor(0x00, 0xB4, 0xD8)   # accent teal
SAGE_AMBER  = RGBColor(0xFF, 0xB7, 0x00)   # warning / highlight amber
SAGE_GREEN  = RGBColor(0x21, 0xC5, 0x5D)   # success green
SAGE_LIGHT  = RGBColor(0xF0, 0xF4, 0xF8)   # slide background
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
GREY        = RGBColor(0x64, 0x74, 0x87)

DISCLAIMER = (
    "⚠️  COST DISCLAIMER — All financial figures, FTE savings, and ROI estimates "
    "on this slide are based on internal assumptions and industry benchmarks. "
    "They must be validated and replaced with actual measured data before any "
    "business decision is made."
)

SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)


# ═══════════════════════════════════════════════════════════════════════════════
# Helper functions
# ═══════════════════════════════════════════════════════════════════════════════

def new_prs():
    prs = Presentation()
    prs.slide_width  = SLIDE_W
    prs.slide_height = SLIDE_H
    return prs


def blank_layout(prs):
    return prs.slide_layouts[6]   # completely blank


def add_rect(slide, l, t, w, h, fill_color=None, line_color=None, line_width=Pt(0)):
    from pptx.util import Pt
    shape = slide.shapes.add_shape(
        1, l, t, w, h   # MSO_SHAPE_TYPE.RECTANGLE = 1
    )
    shape.line.width = line_width
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    else:
        shape.fill.background()
    if line_color:
        shape.line.color.rgb = line_color
    else:
        shape.line.fill.background()
    return shape


def add_text_box(slide, text, l, t, w, h,
                 font_size=18, bold=False, color=SAGE_DARK,
                 align=PP_ALIGN.LEFT, wrap=True, italic=False):
    txBox = slide.shapes.add_textbox(l, t, w, h)
    tf = txBox.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return txBox


def slide_bg(slide, color=SAGE_LIGHT):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def header_band(slide, title, subtitle=None, bg=SAGE_DARK, title_color=WHITE, sub_color=SAGE_TEAL):
    add_rect(slide, 0, 0, SLIDE_W, Inches(1.45), fill_color=bg)
    add_text_box(slide, title,
                 Inches(0.45), Inches(0.12), Inches(12.4), Inches(0.8),
                 font_size=30, bold=True, color=title_color, align=PP_ALIGN.LEFT)
    if subtitle:
        add_text_box(slide, subtitle,
                     Inches(0.45), Inches(0.88), Inches(12.4), Inches(0.48),
                     font_size=16, bold=False, color=sub_color, align=PP_ALIGN.LEFT)


def footer_band(slide, text="SAGE[ai] — Confidential | March 2026", include_disclaimer=False):
    add_rect(slide, 0, Inches(7.1), SLIDE_W, Inches(0.4), fill_color=SAGE_DARK)
    add_text_box(slide, text,
                 Inches(0.3), Inches(7.12), Inches(9), Inches(0.3),
                 font_size=9, color=GREY, align=PP_ALIGN.LEFT)
    if include_disclaimer:
        add_rect(slide, 0, Inches(6.55), SLIDE_W, Inches(0.58),
                 fill_color=RGBColor(0xFF, 0xF3, 0xCC))
        add_text_box(slide, DISCLAIMER,
                     Inches(0.3), Inches(6.57), Inches(12.7), Inches(0.54),
                     font_size=8.5, italic=True, color=RGBColor(0x7B, 0x4F, 0x00),
                     align=PP_ALIGN.LEFT)


def bullet_slide(prs, title, subtitle, bullets, include_disclaimer=False):
    slide = prs.slides.add_slide(blank_layout(prs))
    slide_bg(slide)
    header_band(slide, title, subtitle)

    top = Inches(1.6)
    for b in bullets:
        add_text_box(slide, b,
                     Inches(0.55), top, Inches(12.2), Inches(0.42),
                     font_size=15, color=SAGE_DARK)
        top += Inches(0.5)

    footer_band(slide, include_disclaimer=include_disclaimer)
    return slide


def table_slide(prs, title, subtitle, headers, rows, col_widths=None,
                include_disclaimer=False, note=None):
    slide = prs.slides.add_slide(blank_layout(prs))
    slide_bg(slide)
    header_band(slide, title, subtitle)

    n_cols = len(headers)
    n_rows = len(rows)
    if col_widths is None:
        col_widths = [Inches(12.3 / n_cols)] * n_cols

    tbl_top  = Inches(1.65)
    tbl_left = Inches(0.5)
    row_h    = Inches(0.52)
    tbl_h    = row_h * (n_rows + 1)

    table = slide.shapes.add_table(
        n_rows + 1, n_cols, tbl_left, tbl_top,
        sum(col_widths), tbl_h
    ).table

    # set col widths
    for ci, w in enumerate(col_widths):
        table.columns[ci].width = w

    # header row
    for ci, h in enumerate(headers):
        cell = table.cell(0, ci)
        cell.fill.solid()
        cell.fill.fore_color.rgb = SAGE_BLUE
        p = cell.text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = h
        run.font.bold = True
        run.font.size = Pt(13)
        run.font.color.rgb = WHITE

    # data rows
    for ri, row in enumerate(rows):
        bg = SAGE_LIGHT if ri % 2 == 0 else WHITE
        for ci, val in enumerate(row):
            cell = table.cell(ri + 1, ci)
            cell.fill.solid()
            cell.fill.fore_color.rgb = bg
            p = cell.text_frame.paragraphs[0]
            p.alignment = PP_ALIGN.LEFT
            run = p.add_run()
            run.text = str(val)
            run.font.size = Pt(12)
            run.font.color.rgb = SAGE_DARK

    if note:
        top_pos = tbl_top + tbl_h + Inches(0.1)
        add_text_box(slide, note,
                     Inches(0.5), top_pos, Inches(12.3), Inches(0.4),
                     font_size=11, italic=True, color=GREY)

    footer_band(slide, include_disclaimer=include_disclaimer)
    return slide


# ═══════════════════════════════════════════════════════════════════════════════
#  1.  BUSINESS CASE DECK  (updated)
# ═══════════════════════════════════════════════════════════════════════════════

def build_business_case():
    prs = new_prs()

    # ── Slide 1 — Cover ────────────────────────────────────────────────────────
    slide = prs.slides.add_slide(blank_layout(prs))
    add_rect(slide, 0, 0, SLIDE_W, SLIDE_H, fill_color=SAGE_DARK)
    add_rect(slide, 0, Inches(5.2), SLIDE_W, Inches(0.08), fill_color=SAGE_TEAL)

    add_text_box(slide, "SAGE[ai]",
                 Inches(1), Inches(1.4), Inches(11), Inches(1.4),
                 font_size=64, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text_box(slide, "Autonomous Manufacturing Intelligence",
                 Inches(1), Inches(2.85), Inches(11), Inches(0.7),
                 font_size=26, bold=False, color=SAGE_TEAL, align=PP_ALIGN.CENTER)
    add_text_box(slide, "Transforming Medical Device Production with Agentic AI",
                 Inches(1), Inches(3.55), Inches(11), Inches(0.55),
                 font_size=18, bold=False, color=GREY, align=PP_ALIGN.CENTER)
    add_text_box(slide, "Business Case Presentation — March 2026",
                 Inches(1), Inches(5.5), Inches(11), Inches(0.45),
                 font_size=14, color=GREY, align=PP_ALIGN.CENTER)
    add_text_box(slide, "CONFIDENTIAL — Internal Business Case",
                 Inches(1), Inches(6.1), Inches(11), Inches(0.4),
                 font_size=12, italic=True, color=SAGE_AMBER, align=PP_ALIGN.CENTER)

    # ── Slide 2 — Problem ──────────────────────────────────────────────────────
    bullet_slide(prs,
        "Current State: Manual Bottlenecks at Scale",
        "The hidden cost of running without agentic AI",
        [
            "⏱  45–60 min per error log analysis (manual triage)",
            "🔁  3–5 day MR review backlog — engineers blocked waiting",
            "📋  8+ hours/week compliance reporting — pure administrative overhead",
            "🧠  Expert knowledge lost permanently when engineers leave",
            "⚠️   Error detection relies on human vigilance 24 / 7 — unsustainable",
            "",
            "▶  Result: ~68 hrs/week of preventable manual work — every single week",
        ]
    )

    # ── Slide 3 — Solution ─────────────────────────────────────────────────────
    bullet_slide(prs,
        "SAGE[ai]: Always-On AI Engineering Partner",
        "Replace manual toil with autonomous, auditable intelligence",
        [
            "✅  Analyzes error logs in < 60 seconds  (vs 45–60 min manual)",
            "✅  AI code reviews with multi-step ReAct reasoning loop",
            "✅  Monitors Teams, Metabase, GitLab — real-time event detection",
            "✅  Every decision logged to immutable ISO 13485 audit trail",
            "✅  Human-in-the-loop: AI advises, humans decide — always",
            "✅  Learns from every correction via vector memory (RAG)",
            "✅  Web dashboard for all stakeholders — zero CLI required",
        ]
    )

    # ── Slide 4 — Industry Evidence ────────────────────────────────────────────
    table_slide(prs,
        "Companies Already Winning with Agentic AI",
        "Validated industry outcomes — this is the direction of travel",
        ["Company", "Implementation", "Result"],
        [
            ["Siemens",    "AI predictive maintenance agents",       "30% less unplanned downtime, €1.5B saved/yr"],
            ["BMW",        "AI quality inspection",                  "99.5% defect detection, 30% QC cost reduction"],
            ["Bosch",      "AI-assisted code review",                "40% faster releases, 60% fewer defect escapes"],
            ["Amazon",     "CodeGuru automated review",              "50% reduction in production incidents"],
            ["Microsoft",  "GitHub Copilot",                        "55% faster code completion, 46% more PRs merged/day"],
            ["Medtronic",  "AI in quality management",              "35% shorter CAPA cycle, faster FDA submissions"],
            ["Stryker",    "AI regulatory documentation",           "40% reduction in submission preparation time"],
            ["J&J",        "AI manufacturing analytics",            "20% yield improvement, significant waste reduction"],
        ],
        col_widths=[Inches(1.7), Inches(4.5), Inches(6.1)],
        note="Sources: company annual reports, press releases, and published case studies (2023–2025)."
    )

    # ── Slide 5 — Lean Alignment ───────────────────────────────────────────────
    table_slide(prs,
        "SAGE[ai] Is Built on Lean Principles",
        "Lean development + agentic AI — the natural pairing",
        ["Lean Principle", "How SAGE[ai] Delivers"],
        [
            ["Eliminate Waste (Muda)",         "Removes 60+ hrs/week of repetitive analysis"],
            ["Continuous Improvement (Kaizen)", "Learns from every rejection via RAG memory"],
            ["Error-Proofing (Poka-yoke)",      "Human-in-the-loop prevents AI errors reaching production"],
            ["Visual Management",               "Real-time dashboard for all stakeholders"],
            ["Single Piece Flow",               "Single-lane task queue: deterministic, auditable"],
            ["Respect for People",              "Amplifies engineers; never replaces human judgment"],
        ],
        col_widths=[Inches(4.0), Inches(8.3)]
    )

    # ── Slide 6 — ROI  (with disclaimer) ──────────────────────────────────────
    table_slide(prs,
        "Return on Investment: Month 2 Payback",
        "Conservative baseline — replace with your actual measured numbers",
        ["Activity", "Before SAGE[ai]", "After SAGE[ai]", "Savings"],
        [
            ["Error log analysis",    "45 min × 10/day × 2 eng = 15 hrs/day",  "< 5 min total",        "~93% reduction"],
            ["MR code review",        "3 hrs × 15 MRs/week = 45 hrs/week",     "~4 hrs/week",          "~91% reduction"],
            ["Compliance reporting",  "8 hrs/week manual",                      "0 hrs (auto-gen)",     "~100% reduction"],
            ["Knowledge capture",     "Lost on attrition",                      "Stored in vector mem", "Permanent"],
            ["TOTAL SAVINGS",         "68 hrs/week ≈ 1.7 FTE",                 "—",                    "~€120K/yr*"],
        ],
        col_widths=[Inches(2.8), Inches(3.8), Inches(2.8), Inches(2.9)],
        include_disclaimer=True,
        note="* All figures are ASSUMPTION-BASED estimates. Replace with actual team data before presenting externally."
    )

    # ── Slide 7 — Architecture ─────────────────────────────────────────────────
    slide = prs.slides.add_slide(blank_layout(prs))
    slide_bg(slide)
    header_band(slide,
                "Simple, Auditable, Secure Architecture",
                "End-to-end data flow — one direction, fully logged")

    # flow boxes
    boxes = [
        ("Event\nSources",    Inches(0.3),  SAGE_DARK),
        ("Monitor\nAgent",    Inches(2.3),  SAGE_BLUE),
        ("Task\nQueue",       Inches(4.3),  SAGE_BLUE),
        ("Analyst /\nDev Agent", Inches(6.3), SAGE_BLUE),
        ("Human\nGate",       Inches(8.6),  SAGE_AMBER),
        ("Audit\nTrail",      Inches(10.9), SAGE_GREEN),
    ]
    for label, left, color in boxes:
        add_rect(slide, left, Inches(1.8), Inches(1.85), Inches(1.2),
                 fill_color=color)
        add_text_box(slide, label,
                     left, Inches(1.85), Inches(1.85), Inches(1.1),
                     font_size=12, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        if left < Inches(10.9):
            add_text_box(slide, "→",
                         left + Inches(1.85), Inches(2.2), Inches(0.4), Inches(0.5),
                         font_size=22, bold=True, color=SAGE_DARK, align=PP_ALIGN.CENTER)

    details = [
        ("Event Sources",    "Teams channels · Metabase dashboards · GitLab issues & MRs · Error log uploads"),
        ("AI Agents",        "AnalystAgent (log triage) · DeveloperAgent (code review + MR creation) · PlannerAgent (orchestration) · ReAct multi-step loop"),
        ("Compliance Layer", "Immutable SQLite audit log · ISO 13485 trace IDs · FDA 21 CFR Part 11 · Human approval on every proposal"),
    ]
    top = Inches(3.3)
    for title, body in details:
        add_text_box(slide, f"▸  {title}",
                     Inches(0.5), top, Inches(12.3), Inches(0.35),
                     font_size=13, bold=True, color=SAGE_BLUE)
        add_text_box(slide, body,
                     Inches(0.9), top + Inches(0.33), Inches(12.0), Inches(0.35),
                     font_size=12, color=SAGE_DARK)
        top += Inches(0.75)

    footer_band(slide)

    # ── Slide 8 — Capabilities ─────────────────────────────────────────────────
    table_slide(prs,
        "What SAGE[ai] Does Today",
        "Production-ready capabilities across 12 integration phases — available now",
        ["#", "Capability", "Description"],
        [
            ["🔍", "Log Analysis",        "AI triage in < 60 s — severity RED / AMBER / GREEN with root-cause reasoning"],
            ["🤖", "Code Review",         "ReAct multi-step reasoning + CI pipeline status check"],
            ["📋", "MR Creation",         "Auto-draft merge request from GitLab issue, branch naming, description"],
            ["👁",  "24 / 7 Monitor",     "Teams, Metabase, GitLab event detection — zero polling lag"],
            ["📡", "SSE Streaming",       "Real-time token-by-token output to dashboard — all providers supported"],
            ["🚀", "Onboarding Wizard",   "Plain-language description → full YAML solution in < 60 seconds"],
            ["📚", "Knowledge Base CRUD", "List / add / delete / bulk-import / semantic search the RAG vector store"],
            ["💬", "Slack Approvals",     "Block Kit proposals with Approve/Reject buttons — HMAC-verified callbacks"],
            ["🧪", "Eval & Benchmarking", "YAML test cases · 0–100 scoring · SQLite history · regression tracking"],
            ["🏢", "Multi-Tenant",        "X-SAGE-Tenant header → per-team knowledge collection isolation"],
            ["⏱",  "Durable Workflows",  "LangGraph (interrupt/resume) + Temporal — workflows survive server restart"],
            ["📊", "Audit Trail",         "Every decision traceable, ISO 13485 compliant, UUID per event"],
            ["🌐", "Web Dashboard",       "Dashboard · Analyst · Developer · Audit · Monitor — no CLI needed"],
        ],
        col_widths=[Inches(0.5), Inches(2.2), Inches(9.6)]
    )

    # ── Slide 9 — Compliance ───────────────────────────────────────────────────
    slide = prs.slides.add_slide(blank_layout(prs))
    slide_bg(slide)
    header_band(slide,
                "Built for Regulated Medical Device Environments",
                "Standards compliance is built-in, not bolted-on")

    standards = [
        ("ISO 13485:2016", "Quality Management System"),
        ("ISO 14971:2019", "Risk Management — 7 risks identified & controlled"),
        ("IEC 62304:2006", "Medical Device Software Lifecycle"),
        ("FDA 21 CFR Part 11", "Electronic Records & Signatures"),
        ("FDA Cybersecurity Guidance 2023", "Threat modelling, SBOM, update policy"),
    ]
    top = Inches(1.7)
    for std, desc in standards:
        add_text_box(slide, f"✅  {std}  —  {desc}",
                     Inches(0.6), top, Inches(12.0), Inches(0.4),
                     font_size=15, color=SAGE_DARK)
        top += Inches(0.5)

    add_text_box(slide, "How We Comply",
                 Inches(0.6), top + Inches(0.1), Inches(12.0), Inches(0.4),
                 font_size=16, bold=True, color=SAGE_BLUE)
    top += Inches(0.55)
    how = [
        "Immutable append-only audit log (no deletes, ever)",
        "UUID trace ID on every AI decision",
        "Human approval gate on every proposal — enforced in code",
        "Full DHF: SRS · RTM · V&V Plan · SOUP Inventory",
        "Air-gapped LLM option — no cloud dependency required",
    ]
    for h in how:
        add_text_box(slide, f"  •  {h}",
                     Inches(0.9), top, Inches(12.0), Inches(0.38),
                     font_size=13, color=SAGE_DARK)
        top += Inches(0.42)

    footer_band(slide)

    # ── Slide 10 — Improvement Loop ────────────────────────────────────────────
    slide = prs.slides.add_slide(blank_layout(prs))
    slide_bg(slide)
    header_band(slide,
                "SAGE[ai] Gets Better Every Day — From Your Own Team",
                "Compounding intelligence through the RAG feedback loop")

    steps = [
        ("1", "USER SUBMITS REQUEST",    "Click 'Request Improvement' on any module — logged instantly",           SAGE_BLUE),
        ("2", "AI PLANS",                "Planner Agent decomposes into subtasks, queued for implementation",      SAGE_TEAL),
        ("3", "IMPLEMENTED & LEARNED",   "Change deployed; RAG memory updated — next analysis is smarter",         SAGE_GREEN),
    ]
    left = Inches(0.6)
    for num, title, body, color in steps:
        add_rect(slide, left, Inches(1.75), Inches(3.7), Inches(3.5), fill_color=color)
        add_text_box(slide, num,
                     left + Inches(0.15), Inches(1.9), Inches(0.6), Inches(0.6),
                     font_size=28, bold=True, color=WHITE)
        add_text_box(slide, title,
                     left + Inches(0.15), Inches(2.55), Inches(3.4), Inches(0.5),
                     font_size=14, bold=True, color=WHITE)
        add_text_box(slide, body,
                     left + Inches(0.15), Inches(3.1), Inches(3.4), Inches(1.8),
                     font_size=12, color=WHITE)
        left += Inches(4.2)

    add_text_box(slide,
                 "The system improves itself through the same AI pipeline it provides to engineers.",
                 Inches(0.6), Inches(5.5), Inches(12.0), Inches(0.55),
                 font_size=15, italic=True, bold=True, color=SAGE_DARK, align=PP_ALIGN.CENTER)
    footer_band(slide)

    # ── Slide 11 — Roadmap ─────────────────────────────────────────────────────
    table_slide(prs,
        "From Pilot to Production: 3-Month Plan",
        "12 integration phases complete — pilot deployment is next",
        ["Phase", "Timeline", "Deliverable", "Status"],
        [
            ["Phases 0–4: Core",   "Done",    "Agents · API · Web UI · Approval gate · Audit log",          "✅ Complete"],
            ["Phase 5: Streaming", "Done",    "SSE token streaming for all providers",                       "✅ Complete"],
            ["Phase 6: Onboard",   "Done",    "Plain-language → YAML solution wizard",                       "✅ Complete"],
            ["Phases 7–11",        "Done",    "Knowledge CRUD · Slack · Eval · Multi-tenant · Temporal",     "✅ Complete"],
            ["Pilot Deploy",       "Month 1", "Production deploy, engineer training, KPI baseline",          "🔵 Next"],
            ["Iterate",            "Month 2", "KPI tracking · RAG feedback loop · first improvements",       "🔵 Planned"],
            ["Scale",              "Month 3", "Multi-team rollout · Spira integration · exec dashboards",    "🔵 Planned"],
        ],
        col_widths=[Inches(2.0), Inches(1.4), Inches(6.9), Inches(2.0)]
    )

    # ── Slide 12 — Next Steps ──────────────────────────────────────────────────
    slide = prs.slides.add_slide(blank_layout(prs))
    slide_bg(slide)
    header_band(slide,
                "Next Steps: 3-Month Pilot Proposal",
                "What we need vs what you get")

    add_text_box(slide, "What We Need",
                 Inches(0.5), Inches(1.65), Inches(6.0), Inches(0.45),
                 font_size=16, bold=True, color=SAGE_BLUE)
    needs = [
        "Production deployment approval (internal network)",
        "2-day engineering team training session",
        "Access to GitLab, Metabase, Teams API credentials",
        "Dedicated server or VM  (8-core CPU · 32 GB RAM · optional GPU)",
    ]
    top = Inches(2.15)
    for n in needs:
        add_text_box(slide, f"✅  {n}",
                     Inches(0.7), top, Inches(5.8), Inches(0.42),
                     font_size=13, color=SAGE_DARK)
        top += Inches(0.48)

    add_text_box(slide, "What You Get",
                 Inches(6.9), Inches(1.65), Inches(6.0), Inches(0.45),
                 font_size=16, bold=True, color=SAGE_GREEN)
    gets = [
        "📈  KPI dashboard by Week 2",
        "⏱   60+ hours/week reclaimed from manual work",
        "📋  ISO 13485 audit trail from day 1",
        "🤖  AI that learns from your engineers' expertise",
        "💰  ROI positive within 2 months*",
    ]
    top = Inches(2.15)
    for g in gets:
        add_text_box(slide, g,
                     Inches(7.1), top, Inches(5.8), Inches(0.42),
                     font_size=13, color=SAGE_DARK)
        top += Inches(0.48)

    add_rect(slide, Inches(6.7), Inches(1.6), Inches(0.04), Inches(3.8),
             fill_color=SAGE_TEAL)

    footer_band(slide, include_disclaimer=True)

    # ── Slide 13 — Closing ─────────────────────────────────────────────────────
    slide = prs.slides.add_slide(blank_layout(prs))
    add_rect(slide, 0, 0, SLIDE_W, SLIDE_H, fill_color=SAGE_DARK)
    add_rect(slide, 0, Inches(3.4), SLIDE_W, Inches(0.06), fill_color=SAGE_TEAL)

    add_text_box(slide,
                 "SAGE[ai] is not a replacement for your engineers\n—\nit's their most productive teammate.",
                 Inches(0.8), Inches(1.8), Inches(11.5), Inches(1.8),
                 font_size=30, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    add_text_box(slide, "Let's talk about your pilot.",
                 Inches(0.8), Inches(3.7), Inches(11.5), Inches(0.65),
                 font_size=22, color=SAGE_TEAL, align=PP_ALIGN.CENTER)

    add_text_box(slide,
                 "* All financial estimates are assumption-based. Validate with actual data before business decisions.",
                 Inches(1.5), Inches(6.0), Inches(10.0), Inches(0.5),
                 font_size=10, italic=True, color=GREY, align=PP_ALIGN.CENTER)

    out = os.path.join(os.path.dirname(__file__), "SageAI_Business_Case.pptx")
    prs.save(out)
    print(f"  ✅  Business Case saved → {out}")


# ═══════════════════════════════════════════════════════════════════════════════
#  2.  TECH PITCH DECK  (new)
# ═══════════════════════════════════════════════════════════════════════════════

def build_tech_pitch():
    prs = new_prs()

    # ── Slide 1 — Cover ────────────────────────────────────────────────────────
    slide = prs.slides.add_slide(blank_layout(prs))
    add_rect(slide, 0, 0, SLIDE_W, SLIDE_H, fill_color=SAGE_DARK)
    add_rect(slide, 0, Inches(5.4), SLIDE_W, Inches(0.06), fill_color=SAGE_TEAL)
    add_rect(slide, Inches(0), Inches(0), Inches(0.35), SLIDE_H, fill_color=SAGE_BLUE)

    add_text_box(slide, "SAGE Framework",
                 Inches(1), Inches(1.0), Inches(11), Inches(1.2),
                 font_size=58, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text_box(slide, "Smart Agentic-Guided Empowerment",
                 Inches(1), Inches(2.3), Inches(11), Inches(0.65),
                 font_size=24, bold=False, color=SAGE_TEAL, align=PP_ALIGN.CENTER)
    add_text_box(slide,
                 "A modular autonomous AI agent framework\nbuilt on lean development methodology at machine speed",
                 Inches(1.5), Inches(3.1), Inches(10), Inches(0.9),
                 font_size=17, italic=True, color=GREY, align=PP_ALIGN.CENTER)
    add_text_box(slide, "Technical Platform Overview — March 2026",
                 Inches(1), Inches(5.7), Inches(11), Inches(0.45),
                 font_size=14, color=GREY, align=PP_ALIGN.CENTER)

    # ── Slide 2 — The Core Thesis ──────────────────────────────────────────────
    slide = prs.slides.add_slide(blank_layout(prs))
    slide_bg(slide)
    header_band(slide,
                "The Core Thesis",
                "Lean development is the natural pairing for agentic AI")

    points = [
        ("Eliminate Waste", "Agents handle the repetitive — humans handle the judgment"),
        ("Shorten Feedback Loops", "Surface → Propose → Decide → Compound in seconds, not days"),
        ("Amplify Human Judgment", "Every proposal waits for a human gate — compliance guaranteed"),
        ("Compound Intelligence", "Every correction feeds the vector store — the system gets smarter with use"),
    ]
    top = Inches(1.75)
    for title, body in points:
        add_rect(slide, Inches(0.45), top, Inches(0.12), Inches(0.8), fill_color=SAGE_BLUE)
        add_text_box(slide, title,
                     Inches(0.75), top, Inches(11.5), Inches(0.38),
                     font_size=15, bold=True, color=SAGE_BLUE)
        add_text_box(slide, body,
                     Inches(0.75), top + Inches(0.35), Inches(11.5), Inches(0.38),
                     font_size=13, color=SAGE_DARK)
        top += Inches(1.05)

    footer_band(slide)

    # ── Slide 3 — The SAGE Lean Loop ───────────────────────────────────────────
    slide = prs.slides.add_slide(blank_layout(prs))
    slide_bg(slide)
    header_band(slide,
                "The SAGE Lean Loop",
                "Every task follows this five-phase cycle — no phase is skippable")

    phases = [
        ("1\nSURFACE",      "Agent detects or receives signal\n(log, webhook, trigger)",               SAGE_DARK),
        ("2\nCONTEXTUALIZE","Vector memory searched;\nprior decisions retrieved",                      SAGE_BLUE),
        ("3\nPROPOSE",      "LLM generates action proposal\nwith trace_id",                            SAGE_TEAL),
        ("4\nDECIDE",       "Human reviews, approves\nor rejects with feedback",                       SAGE_AMBER),
        ("5\nCOMPOUND",     "Feedback ingested into vector store;\naudit log updated",                 SAGE_GREEN),
    ]
    left = Inches(0.25)
    for label, body, color in phases:
        add_rect(slide, left, Inches(1.8), Inches(2.4), Inches(3.8), fill_color=color)
        add_text_box(slide, label,
                     left + Inches(0.1), Inches(1.9), Inches(2.2), Inches(0.9),
                     font_size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        add_text_box(slide, body,
                     left + Inches(0.1), Inches(2.9), Inches(2.2), Inches(2.0),
                     font_size=12, color=WHITE, align=PP_ALIGN.CENTER)
        if left < Inches(10):
            add_text_box(slide, "→",
                         left + Inches(2.4), Inches(3.2), Inches(0.3), Inches(0.5),
                         font_size=22, bold=True, color=SAGE_DARK, align=PP_ALIGN.CENTER)
        left += Inches(2.68)

    add_text_box(slide,
                 "Phase 5 feeds Phase 2 for every future task — this is compounding intelligence.",
                 Inches(0.5), Inches(5.85), Inches(12.3), Inches(0.45),
                 font_size=13, italic=True, bold=True, color=SAGE_BLUE, align=PP_ALIGN.CENTER)
    footer_band(slide)

    # ── Slide 4 — Architecture Overview ───────────────────────────────────────
    slide = prs.slides.add_slide(blank_layout(prs))
    slide_bg(slide)
    header_band(slide,
                "Architecture Overview",
                "Framework is domain-blind — solutions are YAML configs, not code")

    layers = [
        ("solutions/<name>/",  "3 YAML files · tests · tools — fully replaceable per domain",    SAGE_GREEN,   Inches(1.6)),
        ("src/core/",          "LLM gateway · queue · project loader · memory — the brain",       SAGE_BLUE,    Inches(2.35)),
        ("src/agents/",        "Analyst · Developer · Monitor · Planner · Universal — the workers", SAGE_BLUE,  Inches(3.1)),
        ("src/interface/",     "FastAPI — the only public interface (UI never calls agents directly)", SAGE_TEAL, Inches(3.85)),
        ("src/memory/",        "Audit log (compliance + training signal) · vector store (RAG)",   SAGE_TEAL,    Inches(4.6)),
        ("web/src/",           "React 18 + TypeScript — reads from the API door only",            SAGE_AMBER,   Inches(5.35)),
    ]

    for path, desc, color, top in layers:
        add_rect(slide, Inches(0.4), top, Inches(2.5), Inches(0.6), fill_color=color)
        add_text_box(slide, path,
                     Inches(0.45), top + Inches(0.05), Inches(2.4), Inches(0.5),
                     font_size=12, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        add_text_box(slide, desc,
                     Inches(3.1), top + Inches(0.1), Inches(9.8), Inches(0.42),
                     font_size=13, color=SAGE_DARK)

    footer_band(slide)

    # ── Slide 5 — YAML-First Design ────────────────────────────────────────────
    slide = prs.slides.add_slide(blank_layout(prs))
    slide_bg(slide)
    header_band(slide,
                "YAML-First Design — Solutions Are Config, Not Code",
                "A new domain requires zero Python changes")

    files = [
        ("project.yaml",  "What this domain IS",  "Declarative agent manifest: project name, active modules,\nintegrations, team structure, compliance standards"),
        ("prompts.yaml",  "How agents THINK",      "Role definitions + system prompts per agent role.\nAdd a new role = add a YAML block. No Python file."),
        ("tasks.yaml",    "What agents CAN DO",    "Task type registry: routing rules, output schemas,\nverification steps, approval requirements."),
    ]
    left = Inches(0.4)
    for fname, tagline, body in files:
        add_rect(slide, left, Inches(1.75), Inches(4.0), Inches(4.5), fill_color=SAGE_DARK)
        add_text_box(slide, fname,
                     left + Inches(0.15), Inches(1.85), Inches(3.7), Inches(0.55),
                     font_size=18, bold=True, color=SAGE_TEAL, align=PP_ALIGN.CENTER)
        add_text_box(slide, tagline,
                     left + Inches(0.15), Inches(2.5), Inches(3.7), Inches(0.4),
                     font_size=13, bold=True, italic=True, color=SAGE_AMBER, align=PP_ALIGN.CENTER)
        add_text_box(slide, body,
                     left + Inches(0.15), Inches(3.0), Inches(3.7), Inches(2.8),
                     font_size=12, color=WHITE, align=PP_ALIGN.LEFT)
        left += Inches(4.47)

    add_text_box(slide,
                 "Separation of concerns is sacred — the framework knows nothing specific about any industry.",
                 Inches(0.5), Inches(6.45), Inches(12.3), Inches(0.45),
                 font_size=13, italic=True, color=SAGE_BLUE, align=PP_ALIGN.CENTER)
    footer_band(slide)

    # ── Slide 6 — Agent Architecture ──────────────────────────────────────────
    table_slide(prs,
        "Agent Architecture",
        "Role-based, not function-based — roles are YAML, not Python classes",
        ["Agent", "Role", "Key Capabilities"],
        [
            ["AnalystAgent",   "Signal intelligence",    "Log triage · severity scoring · root-cause reasoning · ReAct loop"],
            ["DeveloperAgent", "Code intelligence",      "Code review · MR creation · pipeline status · diff analysis"],
            ["MonitorAgent",   "Event surveillance",     "Teams · Metabase · GitLab polling · real-time event routing"],
            ["PlannerAgent",   "Task orchestration",     "Backlog decomposition · wave scheduling · dependency resolution"],
            ["UniversalAgent", "Routing & fallback",     "Dispatches to the correct role; handles new YAML-defined roles"],
        ],
        col_widths=[Inches(2.0), Inches(2.2), Inches(8.1)]
    )

    # ── Slide 7 — Memory Architecture ─────────────────────────────────────────
    slide = prs.slides.add_slide(blank_layout(prs))
    slide_bg(slide)
    header_band(slide,
                "Two-Layer Memory Architecture",
                "Compliance record + compounding intelligence in one design")

    layers_mem = [
        ("SQLite Audit Log",
         "Compliance record & training signal",
         [
             "Append-only — no deletes, ever",
             "UUID trace_id on every AI decision",
             "Full event: input · output · approval · feedback",
             "Queryable for KPI dashboards",
             "FDA 21 CFR Part 11 / ISO 13485 compliant",
         ],
         SAGE_DARK, Inches(0.4)),
        ("ChromaDB Vector Store",
         "Compounding intelligence via RAG",
         [
             "Prior analyses stored as embeddings",
             "Searched before every LLM call",
             "Updated after every human correction",
             "No model retraining — behavioral improvement via retrieval",
             "Memento principle: policy improves, LLM stays frozen",
         ],
         SAGE_BLUE, Inches(6.9)),
    ]

    for title, subtitle, bullets, color, left in layers_mem:
        add_rect(slide, left, Inches(1.75), Inches(5.9), Inches(4.8), fill_color=color)
        add_text_box(slide, title,
                     left + Inches(0.2), Inches(1.85), Inches(5.5), Inches(0.5),
                     font_size=18, bold=True, color=WHITE)
        add_text_box(slide, subtitle,
                     left + Inches(0.2), Inches(2.4), Inches(5.5), Inches(0.4),
                     font_size=13, italic=True, color=SAGE_TEAL)
        top = Inches(2.9)
        for b in bullets:
            add_text_box(slide, f"•  {b}",
                         left + Inches(0.3), top, Inches(5.3), Inches(0.38),
                         font_size=12, color=WHITE)
            top += Inches(0.45)

    add_text_box(slide, "←  divider  →",
                 Inches(6.15), Inches(3.8), Inches(0.6), Inches(0.4),
                 font_size=10, color=GREY, align=PP_ALIGN.CENTER)
    footer_band(slide)

    # ── Slide 8 — Integrations ─────────────────────────────────────────────────
    table_slide(prs,
        "Integration Surface",
        "Everything your team already uses — all live, no API keys required",
        ["Category", "Tools / Protocols", "Status"],
        [
            ["Source Control",    "GitLab (MRs · Issues · Pipelines · CI status)",            "✅ Live"],
            ["Communication",     "Microsoft Teams · Slack (Block Kit approvals · callbacks)", "✅ Live"],
            ["Analytics",         "Metabase (dashboard event detection)",                      "✅ Live"],
            ["LLM Providers",     "Gemini CLI · Claude Code · Ollama · local · GenericCLI",   "✅ No API key"],
            ["Automation",        "n8n webhook → SAGE task routing (HMAC-verified)",           "✅ Live"],
            ["Workflow Engine",   "LangGraph (interrupt/resume) · Temporal (durable)",         "✅ Live"],
            ["Code Execution",    "AutoGen plan → approve → execute in Docker sandbox",        "✅ Live"],
            ["Compliance",        "Spira TM · REST webhook adapters",                          "✅ Live"],
            ["Multi-Tenant",      "X-SAGE-Tenant header → per-team collection isolation",     "✅ Live"],
            ["MCP Servers",       "Model Context Protocol — plug-in tool ecosystem",           "✅ Ready"],
        ],
        col_widths=[Inches(2.3), Inches(6.8), Inches(3.2)]
    )

    # ── Slide 9 — LLM Gateway ─────────────────────────────────────────────────
    slide = prs.slides.add_slide(blank_layout(prs))
    slide_bg(slide)
    header_band(slide,
                "LLM Gateway Design",
                "Provider-agnostic, thread-safe, single-lane inference")

    props = [
        ("No API Key Needed",    "5 providers require zero credential management — login once or run fully local"),
        ("Gemini CLI",           "npm install -g @google/gemini-cli · browser login once · gemini-2.5-flash"),
        ("Claude Code CLI",      "npm install -g @anthropic-ai/claude-code · Anthropic auth once · sonnet-4-6"),
        ("Ollama (local)",       "ollama serve + ollama pull llama3.2 — fully offline, CUDA-accelerated option"),
        ("Generic CLI",          "Any CLI tool via configurable path + {prompt} arg — Aider, LM Studio, etc."),
        ("Thread-Safe + Stream", "threading.Lock on generate() · SSE streaming via generate_stream() for all providers"),
    ]
    top = Inches(1.75)
    for title, body in props:
        add_rect(slide, Inches(0.4), top, Inches(3.5), Inches(0.7), fill_color=SAGE_BLUE)
        add_text_box(slide, title,
                     Inches(0.45), top + Inches(0.1), Inches(3.4), Inches(0.5),
                     font_size=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        add_text_box(slide, body,
                     Inches(4.1), top + Inches(0.12), Inches(8.8), Inches(0.5),
                     font_size=13, color=SAGE_DARK)
        top += Inches(0.85)

    footer_band(slide)

    # ── Slide 10 — Security & Compliance ──────────────────────────────────────
    table_slide(prs,
        "Security & Compliance Architecture",
        "Designed for regulated environments — not retrofitted",
        ["Control", "Implementation", "Standard"],
        [
            ["Audit Log Integrity",  "Append-only SQLite; no UPDATE/DELETE exposed",           "ISO 13485 · FDA 21 CFR Part 11"],
            ["Decision Traceability","UUID trace_id on every LLM call + human action",         "IEC 62304 · ISO 14971"],
            ["Human Gate",           "Approval endpoint required before any action executes",  "ISO 13485 §8.5 CAPA"],
            ["Access Control",       "JWT auth · role-based permissions (roadmap Phase 7)",   "FDA Cybersecurity 2023"],
            ["Data Residency",       "Air-gapped Ollama option; no cloud egress required",    "GDPR · HIPAA ready"],
            ["SOUP Inventory",       "All third-party libs documented; pinned versions",       "IEC 62304 §8.1.2"],
        ],
        col_widths=[Inches(2.5), Inches(5.8), Inches(4.0)]
    )

    # ── Slide 11 — Dev Experience ─────────────────────────────────────────────
    slide = prs.slides.add_slide(blank_layout(prs))
    slide_bg(slide)
    header_band(slide,
                "Developer Experience",
                "From zero to running agent in minutes")

    cmds = [
        ("Setup",       "make venv              # Create .venv, install all deps"),
        ("Low-RAM",     "make venv-minimal      # Skip ChromaDB/embeddings"),
        ("Run backend", "make run PROJECT=medtech  # FastAPI on :8000"),
        ("Run UI",      "make ui                # Vite frontend on :5173"),
        ("Test",        "make test              # Framework unit tests"),
        ("Test all",    "make test-all          # Framework + all solutions"),
    ]

    add_rect(slide, Inches(0.4), Inches(1.7), Inches(12.4), Inches(3.9),
             fill_color=RGBColor(0x1E, 0x1E, 0x2E))

    top = Inches(1.85)
    for label, cmd in cmds:
        add_text_box(slide, f"# {label}",
                     Inches(0.7), top, Inches(1.6), Inches(0.42),
                     font_size=12, color=GREY)
        add_text_box(slide, cmd,
                     Inches(2.3), top, Inches(10.2), Inches(0.42),
                     font_size=12, color=SAGE_GREEN)
        top += Inches(0.52)

    extras = [
        "New solution in < 5 min: copy medtech/ · edit 3 YAML files · run make run PROJECT=<name>",
        "New agent role: add YAML block to prompts.yaml + entry to tasks.yaml — no Python required",
        "Hot-reload YAML via /edit-solution-yaml skill — backend reloads without restart",
    ]
    top = Inches(5.85)
    for e in extras:
        add_text_box(slide, f"▸  {e}",
                     Inches(0.5), top, Inches(12.3), Inches(0.35),
                     font_size=12, color=SAGE_DARK)
        top += Inches(0.42)

    footer_band(slide)

    # ── Slide 12 — Performance & Scalability ──────────────────────────────────
    table_slide(prs,
        "Performance Characteristics",
        "Benchmarks are assumption-based — validate with your workload",
        ["Metric", "Baseline Target*", "Notes"],
        [
            ["Log analysis latency",    "< 60 seconds end-to-end",      "LLM call + RAG retrieval + audit write"],
            ["Code review latency",     "< 3 minutes per MR",           "ReAct loop: up to 5 reasoning steps"],
            ["Task queue throughput",   "~10–20 tasks/hour single-lane", "threading.Lock serializes LLM calls"],
            ["Vector search",           "< 200 ms for top-5 results",   "ChromaDB in-process, local embeddings"],
            ["Audit log write",         "< 5 ms per event",             "SQLite WAL mode, append-only"],
            ["API response (non-LLM)",  "< 50 ms p99",                  "FastAPI + uvicorn, no DB round-trip"],
            ["Min hardware (pilot)",    "8-core CPU · 32 GB RAM",       "GPU optional — Ollama benefits from CUDA"],
        ],
        col_widths=[Inches(3.0), Inches(3.3), Inches(6.0)],
        include_disclaimer=True,
        note="* All latency targets are assumption-based estimates. Replace with measured values from your environment."
    )

    # ── Slide 13 — Two Backlogs ────────────────────────────────────────────────
    slide = prs.slides.add_slide(blank_layout(prs))
    slide_bg(slide)
    header_band(slide,
                "Two Backlogs — One Platform",
                "SAGE serves both solution teams and the open-source community")

    panels = [
        ("Solution Backlog",
         "scope: \"solution\"",
         [
             "Features to build inside your application",
             "Owned by the builder's team",
             "Full SAGE workflow: Log → AI plan → Approve → Implement",
             "Planner Agent decomposes and prioritises",
             "Visible in: Improvements → Solution Backlog tab",
         ],
         SAGE_BLUE, Inches(0.4)),
        ("SAGE Framework Ideas",
         "scope: \"sage\"",
         [
             "Improvements to the SAGE platform itself",
             "Community contributions",
             "Raised as GitHub Issues on the SAGE repo",
             "Open-source community picks them up",
             "Visible in: Improvements → SAGE Framework Ideas tab",
         ],
         SAGE_DARK, Inches(6.9)),
    ]

    for title, tag, bullets, color, left in panels:
        add_rect(slide, left, Inches(1.75), Inches(5.9), Inches(4.8), fill_color=color)
        add_text_box(slide, title,
                     left + Inches(0.2), Inches(1.85), Inches(5.5), Inches(0.5),
                     font_size=18, bold=True, color=WHITE)
        add_text_box(slide, tag,
                     left + Inches(0.2), Inches(2.4), Inches(5.5), Inches(0.38),
                     font_size=12, italic=True, color=SAGE_TEAL)
        top = Inches(2.9)
        for b in bullets:
            add_text_box(slide, f"•  {b}",
                         left + Inches(0.3), top, Inches(5.3), Inches(0.38),
                         font_size=12, color=WHITE)
            top += Inches(0.45)

    add_text_box(slide,
                 "Never mix the two backlogs — the scope field on every FeatureRequest enforces this in code.",
                 Inches(0.5), Inches(6.82), Inches(12.3), Inches(0.38),
                 font_size=11, italic=True, bold=True, color=SAGE_DARK, align=PP_ALIGN.CENTER)
    footer_band(slide)

    # ── Slide 14 — Deployment Options ─────────────────────────────────────────
    table_slide(prs,
        "Deployment Options",
        "From laptop to air-gapped production server",
        ["Mode", "LLM", "Storage", "Best For"],
        [
            ["Local dev",          "Ollama (local)",     "SQLite + ChromaDB local",   "Developer machine, no internet needed"],
            ["Cloud (SaaS LLM)",   "OpenAI / Azure",     "SQLite + ChromaDB local",   "Small team, fastest setup"],
            ["On-premise",         "Azure OpenAI / Ollama", "SQLite + ChromaDB server", "Regulated, data-sovereignty required"],
            ["Air-gapped",         "Ollama only",        "SQLite + ChromaDB server",  "Classified / no-network environments"],
        ],
        col_widths=[Inches(2.0), Inches(2.8), Inches(3.5), Inches(4.0)]
    )

    # ── Slide 15 — SSE Streaming ───────────────────────────────────────────────
    slide = prs.slides.add_slide(blank_layout(prs))
    slide_bg(slide)
    header_band(slide,
                "Real-Time Token Streaming (Phase 5)",
                "Server-Sent Events for progressive AI output — every provider, no API key required")

    stream_points = [
        ("FastAPI StreamingResponse", "text/event-stream content-type; browser-native SSE protocol"),
        ("generate_stream() method",  "Yields chunks from LLMGateway — same thread-safe singleton, streaming path"),
        ("Claude API provider",        "Native token-by-token streaming via Anthropic SDK stream context"),
        ("CLI providers (Gemini, etc)","Full generation runs then output is split into 4-word chunks — progressive feel"),
        ("Ollama provider",            "stream=True to local REST API — true token-level streaming offline"),
        ("Two endpoints",             "POST /analyze/stream · POST /agent/stream — both support X-SAGE-Tenant"),
    ]
    top = Inches(1.75)
    for title, body in stream_points:
        add_rect(slide, Inches(0.4), top, Inches(3.7), Inches(0.65), fill_color=SAGE_TEAL)
        add_text_box(slide, title,
                     Inches(0.45), top + Inches(0.08), Inches(3.6), Inches(0.5),
                     font_size=12, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        add_text_box(slide, body,
                     Inches(4.3), top + Inches(0.1), Inches(8.7), Inches(0.5),
                     font_size=13, color=SAGE_DARK)
        top += Inches(0.82)
    footer_band(slide)

    # ── Slide 16 — Onboarding Wizard ──────────────────────────────────────────
    slide = prs.slides.add_slide(blank_layout(prs))
    slide_bg(slide)
    header_band(slide,
                "Onboarding Wizard — New Solution in 60 Seconds (Phase 6)",
                "Plain language → production-ready YAML configs via AI generation")

    steps_ob = [
        ("1  DESCRIBE", "POST /onboarding/generate with a plain-language description + compliance standards", SAGE_DARK),
        ("2  GENERATE", "LLM produces project.yaml · prompts.yaml · tasks.yaml tailored to your domain",      SAGE_BLUE),
        ("3  SCAFFOLD",  "solutions/<name>/ created with workflows/ · mcp_servers/ · evals/ subdirs",         SAGE_TEAL),
        ("4  RUN",       "make run PROJECT=<name> — agents immediately operational with domain context",       SAGE_GREEN),
    ]
    left = Inches(0.25)
    for label, body, color in steps_ob:
        add_rect(slide, left, Inches(1.85), Inches(3.0), Inches(3.4), fill_color=color)
        add_text_box(slide, label,
                     left + Inches(0.1), Inches(1.95), Inches(2.8), Inches(0.55),
                     font_size=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        add_text_box(slide, body,
                     left + Inches(0.1), Inches(2.6), Inches(2.8), Inches(2.0),
                     font_size=11, color=WHITE, align=PP_ALIGN.CENTER)
        if left < Inches(9.5):
            add_text_box(slide, "→",
                         left + Inches(3.0), Inches(3.2), Inches(0.25), Inches(0.5),
                         font_size=22, bold=True, color=SAGE_DARK, align=PP_ALIGN.CENTER)
        left += Inches(3.27)

    add_text_box(slide,
                 "Templates available at GET /onboarding/templates — or describe any domain from scratch.",
                 Inches(0.5), Inches(5.55), Inches(12.3), Inches(0.45),
                 font_size=13, italic=True, color=SAGE_BLUE, align=PP_ALIGN.CENTER)
    footer_band(slide)

    # ── Slide 17 — Knowledge Base CRUD ────────────────────────────────────────
    table_slide(prs,
        "Knowledge Base CRUD API (Phase 7)",
        "Full lifecycle management for the RAG vector store — list, add, delete, bulk import, search",
        ["Endpoint", "Method", "Description"],
        [
            ["GET  /knowledge/entries",     "List",         "Paginated entries from ChromaDB with metadata"],
            ["POST /knowledge/add",         "Add",          "Single entry with text + arbitrary metadata dict"],
            ["DELETE /knowledge/entry/{id}","Delete",       "Remove by entry UUID — audit-logged"],
            ["POST /knowledge/import",      "Bulk import",  "JSON array of entries — batch ingest from external knowledge base"],
            ["POST /knowledge/search",      "Semantic",     "query + k → top-k cosine-similarity results"],
        ],
        col_widths=[Inches(3.5), Inches(1.5), Inches(7.3)]
    )

    # ── Slide 18 — Slack Two-Way Approval ─────────────────────────────────────
    slide = prs.slides.add_slide(blank_layout(prs))
    slide_bg(slide)
    header_band(slide,
                "Slack Two-Way Approval Loop (Phase 8)",
                "Proposals delivered as Block Kit messages — approve or reject from Slack without leaving the app")

    slack_flow = [
        ("Agent proposes",      "SAGE generates proposal → POST /slack/send-proposal"),
        ("Block Kit message",   "Approve + Reject buttons sent to #sage-approvals channel"),
        ("Engineer clicks",     "Slack POSTs interactive payload to POST /webhook/slack"),
        ("HMAC verification",   "X-Slack-Signature checked — replays rejected, 5-min window enforced"),
        ("Audit log updated",   "Decision + Slack user ID written to SQLite audit trail"),
        ("Vector store updated","Rejection reason ingested for Phase 5 COMPOUND — system learns"),
    ]
    top = Inches(1.75)
    for label, body in slack_flow:
        add_rect(slide, Inches(0.4), top, Inches(0.08), Inches(0.6), fill_color=SAGE_TEAL)
        add_text_box(slide, label,
                     Inches(0.65), top, Inches(3.3), Inches(0.35),
                     font_size=13, bold=True, color=SAGE_BLUE)
        add_text_box(slide, body,
                     Inches(0.65), top + Inches(0.33), Inches(12.0), Inches(0.35),
                     font_size=12, color=SAGE_DARK)
        top += Inches(0.75)
    footer_band(slide)

    # ── Slide 19 — Eval & Benchmarking ────────────────────────────────────────
    slide = prs.slides.add_slide(blank_layout(prs))
    slide_bg(slide)
    header_band(slide,
                "Eval & Benchmarking Framework (Phase 9)",
                "YAML-defined test cases — score agent quality, track regressions over time")

    eval_panels = [
        ("Eval YAML Format",
         "Define test cases in solutions/<name>/evals/<suite>.yaml",
         [
             "prompt: the input to send to the agent",
             "expected_keywords: list of terms that must appear",
             "max_response_length: optional brevity constraint",
             "Scoring: 70 pts keyword coverage + 30 pts length",
         ],
         SAGE_DARK, Inches(0.4)),
        ("API + Persistence",
         "Run suites on demand, store history in SQLite",
         [
             "GET  /eval/suites — list available suites",
             "POST /eval/run   — run one or all suites",
             "GET  /eval/history — trend data per suite",
             "Scores 0–100; track regressions across LLM upgrades",
         ],
         SAGE_BLUE, Inches(6.9)),
    ]
    for title, subtitle, bullets, color, left in eval_panels:
        add_rect(slide, left, Inches(1.75), Inches(5.9), Inches(4.8), fill_color=color)
        add_text_box(slide, title,
                     left + Inches(0.2), Inches(1.85), Inches(5.5), Inches(0.5),
                     font_size=17, bold=True, color=WHITE)
        add_text_box(slide, subtitle,
                     left + Inches(0.2), Inches(2.4), Inches(5.5), Inches(0.4),
                     font_size=12, italic=True, color=SAGE_TEAL)
        top = Inches(2.9)
        for b in bullets:
            add_text_box(slide, f"•  {b}",
                         left + Inches(0.3), top, Inches(5.3), Inches(0.38),
                         font_size=12, color=WHITE)
            top += Inches(0.48)
    footer_band(slide)

    # ── Slide 20 — Multi-Tenant Isolation ─────────────────────────────────────
    slide = prs.slides.add_slide(blank_layout(prs))
    slide_bg(slide)
    header_band(slide,
                "Multi-Tenant Isolation (Phase 10)",
                "Per-team knowledge namespacing via a single request header")

    mt_points = [
        ("X-SAGE-Tenant header",  "Add to any request: X-SAGE-Tenant: firmware_team — zero auth change needed"),
        ("TenantMiddleware",       "BaseHTTPMiddleware extracts header → sets ContextVar before request reaches agent"),
        ("ContextVar isolation",   "Python contextvars — each async request has its own tenant context; no cross-talk"),
        ("Collection namespacing", "tenant_scoped_collection() → <tenant>_knowledge ChromaDB collection"),
        ("GET /tenant/context",   "Debug endpoint shows resolved tenant for current request headers"),
        ("Backward compatible",   "Requests without header use global collection — no breaking change"),
    ]
    top = Inches(1.75)
    for title, body in mt_points:
        add_rect(slide, Inches(0.4), top, Inches(3.5), Inches(0.65), fill_color=SAGE_BLUE)
        add_text_box(slide, title,
                     Inches(0.45), top + Inches(0.08), Inches(3.4), Inches(0.5),
                     font_size=12, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        add_text_box(slide, body,
                     Inches(4.1), top + Inches(0.1), Inches(8.8), Inches(0.5),
                     font_size=13, color=SAGE_DARK)
        top += Inches(0.82)
    footer_band(slide)

    # ── Slide 21 — Temporal Durable Workflows ─────────────────────────────────
    table_slide(prs,
        "Temporal Durable Workflows (Phase 11)",
        "Production-grade workflow orchestration with automatic retry, state persistence, and LangGraph fallback",
        ["Capability", "Implementation", "Benefit"],
        [
            ["Durable execution",   "Temporal SDK — workflow survives server restart",          "Zero lost work"],
            ["Automatic retry",     "Temporal retry policies per activity",                    "Self-healing pipelines"],
            ["Human pause",         "LangGraph interrupt_before — pause at any node for review","Compliance gate in workflow"],
            ["Resume after approval","POST /workflow/resume · POST /temporal/workflow/resume",  "Async human-in-the-loop"],
            ["Graceful fallback",   "Temporal unavailable → LangGraph runner transparently",    "No hard dependency on server"],
            ["Status tracking",     "GET /workflow/status/{id} · GET /temporal/workflow/list",  "Full observability"],
        ],
        col_widths=[Inches(2.5), Inches(5.3), Inches(4.5)]
    )

    # ── Slide 22 — Phase Summary ───────────────────────────────────────────────
    table_slide(prs,
        "Platform Completeness — 12 Integration Phases",
        "From core agent loop to enterprise-grade orchestration — all live",
        ["Phase", "Feature", "Status"],
        [
            ["0–4",  "Core agents · API · Web UI · Approval gate · Audit log",       "✅ Complete"],
            ["5",    "SSE streaming — real-time token output in dashboard",           "✅ Complete"],
            ["6",    "Onboarding wizard — plain-language → YAML solution in 60 s",   "✅ Complete"],
            ["7",    "Knowledge base CRUD — list / add / delete / bulk / search",    "✅ Complete"],
            ["8",    "Slack two-way approval — Block Kit buttons, HMAC verified",     "✅ Complete"],
            ["9",    "Eval & benchmarking — YAML test cases, 0–100 scoring, history","✅ Complete"],
            ["10",   "Multi-tenant isolation — X-SAGE-Tenant header namespacing",    "✅ Complete"],
            ["11",   "Temporal durable workflows + LangGraph fallback",               "✅ Complete"],
        ],
        col_widths=[Inches(0.9), Inches(8.3), Inches(3.1)]
    )

    # ── Slide 23 — Closing ─────────────────────────────────────────────────────
    slide = prs.slides.add_slide(blank_layout(prs))
    add_rect(slide, 0, 0, SLIDE_W, SLIDE_H, fill_color=SAGE_DARK)
    add_rect(slide, 0, Inches(4.0), SLIDE_W, Inches(0.06), fill_color=SAGE_TEAL)
    add_rect(slide, Inches(0), Inches(0), Inches(0.35), SLIDE_H, fill_color=SAGE_BLUE)

    add_text_box(slide, "SAGE Framework",
                 Inches(0.8), Inches(0.8), Inches(11.5), Inches(0.9),
                 font_size=44, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text_box(slide,
                 "Lean development at machine speed.\nYAML-first. Agent-native. Human-in-the-loop always.",
                 Inches(0.8), Inches(1.75), Inches(11.5), Inches(1.0),
                 font_size=20, color=SAGE_TEAL, align=PP_ALIGN.CENTER)

    links = [
        "📂  github.com/your-org/sage",
        "📖  /docs — SETUP.md · USER_GUIDE.md · ADDING_A_PROJECT.md",
        "🧪  make test-all — full test suite in < 2 min",
        "🚀  make run PROJECT=medtech — running in 60 seconds",
    ]
    top = Inches(4.3)
    for l in links:
        add_text_box(slide, l,
                     Inches(2.0), top, Inches(9.3), Inches(0.45),
                     font_size=14, color=WHITE, align=PP_ALIGN.LEFT)
        top += Inches(0.52)

    add_text_box(slide,
                 "⚠️  All performance figures and cost estimates are assumption-based — validate before decisions.",
                 Inches(1.5), Inches(6.55), Inches(10.0), Inches(0.45),
                 font_size=10, italic=True, color=GREY, align=PP_ALIGN.CENTER)

    out = os.path.join(os.path.dirname(__file__), "SageAI_Tech_Pitch.pptx")
    prs.save(out)
    print(f"  ✅  Tech Pitch saved → {out}")


# ═══════════════════════════════════════════════════════════════════════════════
if __name__ == "__main__":
    print("Building presentations…")
    build_business_case()
    build_tech_pitch()
    print("Done.")
